#!/usr/bin/env python3
"""
slide_kit.py — 입찰 제안서 PPTX 공통 렌더링 툴킷 v3.7

컨설팅 스타일 + Modern 컬러 시스템
- 미세 그림자 + 그라디언트 깊이감
- 라운드 코너 카드 + 직각 도형 혼용
- 도형 위 텍스트는 항상 중앙 정렬
- 매 페이지 도식화/구조 중심
- Action Title (인사이트 기반 제목)
- Source 출처 표기
- 1 슬라이드 = 1 인사이트

v3.6: 컬러 유틸(darken/lighten) + 그라디언트 커버/클로징 + 그림자 확대 적용
      + SemiBold 타이포 + KPIS/GRID/TABLE 폴리시 + LINE_CHART smooth 버그 수정

v3.7: 안전성 강화 — 겹침/크래시 방지
      + _ensure_emu(): raw float → Inches 자동 변환 (VStack.next_raw() 호환)
      + 모든 컴포넌트 함수에 y/h 자동 변환 적용
      + QUOTE: text_color 파라미터 추가 (다크 배경 호환)
      + COMPARE: h 파라미터 추가 (본문 높이 조절 가능)
      + FLOW/TIMELINE/HIGHLIGHT/TABLE: 실제 높이 반환 (VStack 예약 정확도↑)
      + VStack: max_y 제한, remaining_safe, would_overflow 추가

사용법:
    import importlib.util
    spec = importlib.util.spec_from_file_location('sk', '경로/src/generators/slide_kit.py')
    sk = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(sk)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
import os


# ── 안전 유틸리티 (v3.7) ──────────────────────────────────────

_EMU_THRESHOLD = 914400 * 0.05  # 0.05" = 45720 EMU — 이보다 크면 EMU로 간주

def _ensure_emu(val):
    """raw float(인치) 또는 EMU를 안전하게 EMU int로 변환.

    v3.8 개선: float 값의 크기를 보고 인치 vs EMU를 자동 판별.
    - |val| < _EMU_THRESHOLD(45720) → 인치로 간주 → Inches(val)
    - |val| >= _EMU_THRESHOLD       → 이미 EMU → int(val)

    이전 버전 버그: CW * 0.5 같은 EMU 스케일 float 연산 결과가
    인치로 오인되어 Inches(6095847.5) = 5.57조 EMU 발생 → PPT 파손.
    """
    if val is None:
        return None
    if isinstance(val, float):
        if abs(val) >= _EMU_THRESHOLD:
            return int(val)       # EMU 스케일 float → int로 안전 변환
        return Inches(val)        # 인치 스케일 float → EMU 변환
    return val  # int (이미 EMU) 또는 기타


# ═══════════════════════════════════════════════════════════════
#  1. 디자인 상수
# ═══════════════════════════════════════════════════════════════

# ── 컬러 유틸리티 (v3.6) ──────────────────────────────────────

def _darken(r, g, b, amount=0.25):
    """RGB를 amount 비율만큼 어둡게"""
    return (max(0, int(r * (1 - amount))),
            max(0, int(g * (1 - amount))),
            max(0, int(b * (1 - amount))))


def _lighten(r, g, b, amount=0.35):
    """RGB를 amount 비율만큼 밝게"""
    return (min(255, int(r + (255 - r) * amount)),
            min(255, int(g + (255 - g) * amount)),
            min(255, int(b + (255 - b) * amount)))


def darken(color, amount=0.25):
    """RGBColor를 amount만큼 어둡게 → 새 RGBColor 반환"""
    r, g, b = color[0], color[1], color[2]
    dr, dg, db = _darken(r, g, b, amount)
    return RGBColor(dr, dg, db)


def lighten(color, amount=0.35):
    """RGBColor를 amount만큼 밝게 → 새 RGBColor 반환"""
    r, g, b = color[0], color[1], color[2]
    lr, lg, lb = _lighten(r, g, b, amount)
    return RGBColor(lr, lg, lb)


# 컬러 팔레트 (Modern 절제 스타일)
C = {
    "primary":   RGBColor(0, 44, 95),       # #002C5F  주색 (다크블루)
    "secondary": RGBColor(0, 170, 210),      # #00AAD2  보조색 (스카이블루)
    "teal":      RGBColor(0, 161, 156),      # #00A19C  틸
    "accent":    RGBColor(230, 51, 18),      # #E63312  강조 (레드)
    "dark":      RGBColor(33, 33, 33),       # #212121  본문 기본색
    "light":     RGBColor(245, 245, 245),    # #F5F5F5  밝은 배경
    "white":     RGBColor(255, 255, 255),    # #FFFFFF
    "gray":      RGBColor(117, 117, 117),    # #757575  보조 텍스트
    "lgray":     RGBColor(200, 200, 200),    # #C8C8C8  구분선
    "green":     RGBColor(46, 125, 50),      # #2E7D32  성과/긍정
    "orange":    RGBColor(245, 166, 35),     # #F5A623  주의
    "gold":      RGBColor(197, 151, 62),     # #C5973E  프리미엄
}

# 파생 컬러 (v3.6) — 깊이감/계층 표현용
C["primary_dark"]   = darken(C["primary"], 0.3)      # 더 진한 네이비
C["primary_light"]  = lighten(C["primary"], 0.85)     # 연한 블루 배경
C["secondary_dark"] = darken(C["secondary"], 0.25)    # 진한 스카이블루
C["secondary_light"]= lighten(C["secondary"], 0.80)   # 연한 스카이 배경
C["teal_light"]     = lighten(C["teal"], 0.80)        # 연한 틸 배경
C["accent_light"]   = lighten(C["accent"], 0.80)      # 연한 레드 배경
C["green_light"]    = lighten(C["green"], 0.80)       # 연한 그린 배경
C["card_bg"]        = RGBColor(250, 250, 252)         # 카드 배경 (약간 블루)
C["card_border"]    = RGBColor(230, 232, 236)         # 카드 테두리

# 슬라이드 규격 (16:9)
SW = Inches(13.333)
SH = Inches(7.5)
ML = Inches(0.8)       # 좌측 여백
MR = Inches(0.8)       # 우측 여백
MT_Y = Inches(0.4)     # 상단 여백
CW = SW - ML - MR      # 콘텐츠 너비

# 타이포그래피
FONT = "Pretendard"

# 폰트 웨이트 (v3.6) — python-pptx는 font.name으로 웨이트 구분
FONT_W = {
    "light":    "Pretendard Light",
    "regular":  "Pretendard",
    "medium":   "Pretendard Medium",
    "semibold": "Pretendard SemiBold",
    "bold":     "Pretendard Bold",
    "black":    "Pretendard Black",
}

# 폰트 사이즈 체계
SZ = {
    "hero":      60,   # 표지
    "divider":   40,   # 섹션 구분자
    "action":    20,   # Action Title (슬라이드 제목)
    "subtitle":  16,   # 부제
    "body":      13,   # 본문
    "body_sm":   11,   # 본문 소
    "caption":   10,   # 캡션
    "source":     8,   # 출처
}


def set_font(font_name, weights=None):
    """폰트 패밀리 변경 (v3.8)

    Args:
        font_name: 기본 폰트 이름 (예: "페이퍼로지", "Pretendard")
        weights: 웨이트별 폰트 이름 dict (없으면 기본 폰트에 bold 속성으로 대체)
    """
    global FONT, FONT_W
    FONT = font_name
    if weights:
        FONT_W.update(weights)
    else:
        FONT_W.update({
            "light":    font_name,
            "regular":  font_name,
            "medium":   font_name,
            "semibold": font_name,
            "bold":     font_name,
            "black":    font_name,
        })
    return font_name


def set_slide_size(width_in, height_in, margin_in=None, scale_fonts=True):
    """슬라이드 크기 변경 (v3.8/v4.1) — 연관 상수 자동 재계산

    Args:
        width_in: 슬라이드 너비 (인치)
        height_in: 슬라이드 높이 (인치)
        margin_in: 좌우 마진 (인치, 없으면 너비 비례 자동 계산)
        scale_fonts: True면 SZ 딕셔너리를 캔버스 크기에 비례해 자동 스케일
    """
    global SW, SH, ML, MR, CW, MT_Y, CW_IN, ML_IN, CGAP, GAP, Z, SZ
    # 기준 13.333 × 7.5 대비 스케일 (원본 폰트 사이즈 복원용)
    if scale_fonts:
        _scale = width_in / 13.333
        # SZ 전체를 기준값 대비 스케일링 (원본 보존용 _SZ_BASE)
        if not hasattr(set_slide_size, "_SZ_BASE"):
            set_slide_size._SZ_BASE = {k: v for k, v in SZ.items()}
        for k, base_v in set_slide_size._SZ_BASE.items():
            # 최소 6pt, 최대 원본값
            scaled = max(6, int(round(base_v * _scale)))
            SZ[k] = scaled
    SW = Inches(width_in)
    SH = Inches(height_in)
    if margin_in is not None:
        ML = MR = Inches(margin_in)
    else:
        ML = MR = Inches(round(width_in * 0.06, 2))  # ~6% 마진
    MT_Y = Inches(round(height_in * 0.053, 2))       # ~5.3% 상단
    CW = SW - ML - MR
    # Zone 상수 재계산
    ML_IN = ML / 914400
    CW_IN = CW / 914400
    CGAP = round(CW_IN * 0.024, 3)   # 컬럼 간격
    GAP = round(CW_IN * 0.016, 3)    # 작은 간격
    # 콘텐츠 Zone 재계산 (비율 기반)
    sh_in = height_in
    tb_h = round(sh_in * 0.117, 2)   # 타이틀바 높이 (~0.88" at 7.5)
    ct_y = round(sh_in * 0.147, 2)   # 콘텐츠 시작 (~1.1" at 7.5)
    ct_b = round(sh_in * 0.867, 2)   # 콘텐츠 하단 (~6.5" at 7.5)
    ct_h = round(ct_b - ct_y, 2)     # 콘텐츠 높이
    ft_y = round(sh_in * 0.893, 2)   # 푸터 시작
    ft_h = round(sh_in - ft_y, 2)    # 푸터 높이
    Z = {
        "tb_y":  0,
        "tb_h":  tb_h,
        "ct_y":  ct_y,
        "ct_h":  ct_h,
        "ct_b":  ct_b,
        "ft_y":  ft_y,
        "ft_h":  ft_h,
        "full":  (ML_IN, CW_IN),
        "left":  (ML_IN, CW_IN * 0.48),
        "right": (ML_IN + CW_IN * 0.52, CW_IN * 0.48),
        "center": (ML_IN + CW_IN * 0.15, CW_IN * 0.70),
    }
    return (SW, SH)


# ═══════════════════════════════════════════════════════════════
#  2. 프레젠테이션 / 슬라이드 생성
# ═══════════════════════════════════════════════════════════════

def new_presentation():
    """16:9 빈 프레젠테이션 생성"""
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


def new_slide(prs, inherit_bg=None):
    """빈 레이아웃 슬라이드 추가.

    Args:
        inherit_bg:
            True  → 강제 마스터 배경 상속
            False → 강제 흰 배경
            None  → 자동 (마스터 배경이 설정되어 있으면 상속, 아니면 흰 배경)
    """
    # Blank 레이아웃 자동 선택 (기본 인덱스 6, prune 후엔 0)
    blank_layout = None
    for lo in prs.slide_layouts:
        if lo.name.lower() in ("blank", "빈 화면", "빈화면"):
            blank_layout = lo
            break
    if blank_layout is None:
        # fallback: 마지막 레이아웃 (보통 Blank)
        try:
            blank_layout = prs.slide_layouts[6]
        except IndexError:
            blank_layout = prs.slide_layouts[-1] if len(prs.slide_layouts) else prs.slide_layouts[0]
    s = prs.slides.add_slide(blank_layout)

    # 자동 감지: 마스터에 커스텀 bg 요소가 있으면 상속
    if inherit_bg is None:
        try:
            from lxml import etree
            ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
            master = prs.slide_master
            cSld = master._element.find(f'{{{ns_p}}}cSld')
            bg_el = cSld.find(f'{{{ns_p}}}bg') if cSld is not None else None
            inherit_bg = bg_el is not None and len(bg_el) > 0
        except Exception:
            inherit_bg = False

    if inherit_bg:
        clear_slide_bg(s)
    else:
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = C["white"]
    return s


def _detect_dark_bg(s):
    """슬라이드/마스터 XML에서 배경색을 읽어 어두운지 판정.

    s.background 접근 없이 XML만 검사 → 상속 유지.
    """
    import re
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'

    def _check_element(el):
        """요소에서 solidFill srgbClr 찾기 → luma 체크."""
        for sf in el.iter(f'{{{ns_a}}}solidFill'):
            for c in sf.iter(f'{{{ns_a}}}srgbClr'):
                val = c.get('val')
                if val and len(val) == 6:
                    r = int(val[0:2], 16)
                    g = int(val[2:4], 16)
                    b = int(val[4:6], 16)
                    luma = r * 0.299 + g * 0.587 + b * 0.114
                    return luma < 100
        return None

    # 1. 슬라이드 자체 bg 확인
    try:
        cSld = s._element.find(f'{{{ns_p}}}cSld')
        if cSld is not None:
            bg = cSld.find(f'{{{ns_p}}}bg')
            if bg is not None:
                r = _check_element(bg)
                if r is not None:
                    return r
    except Exception:
        pass
    # 2. 마스터 bg 확인
    try:
        master = s.part.slide_layout.slide_master
        m_cSld = master._element.find(f'{{{ns_p}}}cSld')
        if m_cSld is not None:
            m_bg = m_cSld.find(f'{{{ns_p}}}bg')
            if m_bg is not None:
                r = _check_element(m_bg)
                if r is not None:
                    return r
    except Exception:
        pass
    return False


def clear_slide_bg(s):
    """슬라이드의 <p:bg> 요소 제거 — 마스터 배경 상속 활성화."""
    try:
        from lxml import etree
        ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        cSld = s._element.find(f'{{{ns}}}cSld')
        if cSld is not None:
            bg = cSld.find(f'{{{ns}}}bg')
            if bg is not None:
                cSld.remove(bg)
    except Exception:
        pass
    return s


def setup_master_background(prs, color=None, gradient=None):
    """프레젠테이션 마스터 배경 설정 — 모든 슬라이드에 자동 적용.

    마스터에 한 번 설정하면 inherit_bg=True로 생성한 슬라이드가 자동 상속.
    개별 슬라이드는 여전히 bg()/gradient_bg()로 덮어쓰기 가능.

    Args:
        prs: Presentation 객체
        color: 솔리드 배경색 (RGBColor 또는 (r,g,b))
        gradient: (c1, c2) 튜플이면 그라디언트 배경 (상→하)

    Example:
        prs = new_presentation()
        setup_master_background(prs, color=tok("surface/darker"))
        # 이후 new_slide(prs, inherit_bg=True)로 만들면 자동 다크 배경
    """
    m = prs.slide_master
    if gradient:
        c1, c2 = gradient
        gradient_shape(m.background._element.getparent().getparent(),
                        c1, c2) if False else None
        # 마스터 그라디언트는 XML 직접 편집 필요
        _master_gradient(m, c1, c2)
    elif color is not None:
        if isinstance(color, tuple):
            color = RGBColor(*color)
        m.background.fill.solid()
        m.background.fill.fore_color.rgb = color
    return m


def apply_palette(*, bg, text, key, sub1, sub2,
                    card=None, text_muted=None):
    """사용자 5색 팔레트 적용 — editorial_dark 테마 덮어쓰기.

    Args:
        bg:    배경 (hex "#RRGGBB" 또는 (r,g,b))
        text:  본문 흰색/밝은색
        key:   메인 키 컬러 (헤더 악센트, 뱃지)
        sub1:  서브 컬러 1 (두번째 라벨)
        sub2:  서브 컬러 2 (eyebrow/태그)
        card:  카드 표면 (None이면 bg 밝기 +10)
        text_muted: 뮤티드 텍스트 (None이면 text 40% 톤)

    Example:
        apply_palette(
            bg   = "#1C1F28",
            text = "#FFFFFF",
            key  = "#5F70FC",
            sub1 = "#6296FF",
            sub2 = "#66FFFF",
        )
    """
    def _to_rgb(v):
        if isinstance(v, str):
            v = v.lstrip("#").upper()
            return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))
        if isinstance(v, tuple):
            return RGBColor(*v)
        return v   # 이미 RGBColor

    bg_c    = _to_rgb(bg)
    text_c  = _to_rgb(text)
    key_c   = _to_rgb(key)
    sub1_c  = _to_rgb(sub1)
    sub2_c  = _to_rgb(sub2)

    # 카드 표면 기본값: bg보다 약간 밝게
    if card is None:
        card_c = RGBColor(min(255, bg_c[0]+13),
                          min(255, bg_c[1]+15),
                          min(255, bg_c[2]+18))
    else:
        card_c = _to_rgb(card)

    # 뮤티드 텍스트 기본값: text와 bg의 중간
    if text_muted is None:
        muted_c = RGBColor(
            (int(text_c[0]) + int(bg_c[0])*2) // 3,
            (int(text_c[1]) + int(bg_c[1])*2) // 3,
            (int(text_c[2]) + int(bg_c[2])*2) // 3,
        )
    else:
        muted_c = _to_rgb(text_muted)

    # 살짝 더 어두운 bg (darker surface)
    darker_c = RGBColor(max(0, bg_c[0]-8),
                        max(0, bg_c[1]-8),
                        max(0, bg_c[2]-8))

    # TOKENS 전체 업데이트
    TOKENS.update({
        "surface/base":    bg_c,
        "surface/raised":  card_c,
        "surface/dark":    card_c,
        "surface/darker":  bg_c,        # 마스터 배경
        "surface/overlay": darker_c,
        "border/subtle":   card_c,
        "border/dark":     card_c,

        "brand/primary":   key_c,
        "brand/secondary": sub1_c,
        "brand/deep":      key_c,

        "neon/cyan":       sub2_c,
        "neon/aqua":       sub2_c,
        "neon/electric":   sub2_c,
        "neon/mint":       sub2_c,

        "text/on_dark":    text_c,
        "text/on_light":   text_c,
        "text/muted":      muted_c,
        "text/subtle":     muted_c,
        "text/accent":     key_c,

        # aliases
        "primary":    key_c,
        "secondary":  sub2_c,
        "accent":     sub1_c,
        "text":       text_c,
        "muted":      muted_c,
        "bg":         bg_c,
        "card":       card_c,
        "border":     card_c,
    })

    # THEMES["editorial_dark"]도 업데이트 (apply_theme 시 C 동기)
    THEMES["editorial_dark"] = {
        "primary":   tuple(key_c),
        "secondary": tuple(sub1_c),
        "teal":      tuple(sub2_c),
        "accent":    tuple(key_c),
        "dark":      tuple(bg_c),
        "light":     tuple(text_c),
    }
    apply_theme("editorial_dark")
    return TOKENS


def propose_palette():
    """현재 팔레트 프로포절 출력 (콘솔). 새 덱 시작 시 호출."""
    print("\n=== slide_kit 팔레트 제안 (editorial_dark 기본) ===")
    print(f"  ① 배경   : #1C1F28  (딥 네이비)")
    print(f"  ② 글자   : #FFFFFF  (흰색) + 뮤티드 자동")
    print(f"  ③ 키     : #5F70FC  (브랜드 퍼플) ← 헤더 악센트/뱃지")
    print(f"  ④ 서브1  : #6296FF  (브랜드 블루) ← 두번째 라벨")
    print(f"  ⑤ 서브2  : #66FFFF  (사이언)       ← eyebrow/태그")
    print(f"\n  적용: apply_palette(bg=, text=, key=, sub1=, sub2=)")
    print(f"  기본 유지: apply_theme('editorial_dark')\n")


# ═══════════════════════════════════════════════════════════════════════
#  큐레이션 팔레트 라이브러리 + 추천 시스템
# ═══════════════════════════════════════════════════════════════════════

PALETTE_LIBRARY = {
    # ── DARK MODE ──────────────────────────────────────────────────
    "editorial_dark": {
        "name": "Editorial Dark",
        "desc": "에디토리얼 다크 — MIT Tech Review / Stripe 톤 (기본)",
        "colors": {"bg": "#1C1F28", "text": "#FFFFFF",
                    "key": "#5F70FC", "sub1": "#6296FF", "sub2": "#66FFFF"},
        "tags": ["다크", "에디토리얼", "테크", "모던", "게이밍", "IT"],
        "fit": ["it_system", "marketing_pr", "event", "consulting"],
        "mood": "professional, modern, focused",
    },
    "cyberpunk_neon": {
        "name": "Cyberpunk Neon",
        "desc": "사이버펑크 네온 — 게이밍/E스포츠",
        "colors": {"bg": "#0A0A1F", "text": "#FFFFFF",
                    "key": "#FF0080", "sub1": "#00FFE5", "sub2": "#F0FF00"},
        "tags": ["네온", "게이밍", "이스포츠", "강렬", "미래"],
        "fit": ["event", "marketing_pr"],
        "mood": "bold, electric, futuristic",
    },
    "midnight_forest": {
        "name": "Midnight Forest",
        "desc": "미드나잇 포레스트 — 환경/지속가능성",
        "colors": {"bg": "#0F1F1A", "text": "#F4F7F2",
                    "key": "#47C472", "sub1": "#8BD8A8", "sub2": "#F4C542"},
        "tags": ["자연", "환경", "친환경", "다크", "프리미엄"],
        "fit": ["public", "consulting", "marketing_pr"],
        "mood": "natural, sustainable, elegant",
    },
    "deep_luxury": {
        "name": "Deep Luxury",
        "desc": "딥 럭셔리 — 명품/프리미엄 브랜드",
        "colors": {"bg": "#1A1022", "text": "#F5E6D3",
                    "key": "#C9A961", "sub1": "#8B7ED8", "sub2": "#E8B4C8"},
        "tags": ["럭셔리", "프리미엄", "명품", "다크", "엘레강스"],
        "fit": ["marketing_pr", "consulting"],
        "mood": "luxurious, premium, refined",
    },
    "finance_navy": {
        "name": "Finance Navy",
        "desc": "파이낸스 네이비 — 금융/컨설팅",
        "colors": {"bg": "#0A1628", "text": "#FFFFFF",
                    "key": "#3B82F6", "sub1": "#60A5FA", "sub2": "#FCD34D"},
        "tags": ["금융", "컨설팅", "신뢰", "공식", "보수"],
        "fit": ["consulting", "it_system", "public"],
        "mood": "trustworthy, analytical, authoritative",
    },

    # ── LIGHT MODE ─────────────────────────────────────────────────
    "minimal_light": {
        "name": "Minimal Light",
        "desc": "미니멀 라이트 — Apple/Notion 톤",
        "colors": {"bg": "#FAFAFA", "text": "#1D1D1F",
                    "key": "#0071E3", "sub1": "#34C759", "sub2": "#FF9500"},
        "tags": ["미니멀", "라이트", "애플", "깔끔", "모던"],
        "fit": ["it_system", "consulting", "marketing_pr"],
        "mood": "clean, minimal, approachable",
    },
    "paper_warm": {
        "name": "Paper Warm",
        "desc": "페이퍼 웜 — 출판/에디토리얼 잡지",
        "colors": {"bg": "#FAF5EE", "text": "#2B2118",
                    "key": "#C94F3B", "sub1": "#D89E5B", "sub2": "#5F7A61"},
        "tags": ["웜", "페이퍼", "에디토리얼", "출판", "라이프스타일"],
        "fit": ["marketing_pr", "event"],
        "mood": "warm, human, storytelling",
    },
    "nordic_cool": {
        "name": "Nordic Cool",
        "desc": "노르딕 쿨 — 북유럽 미니멀",
        "colors": {"bg": "#F4F1EC", "text": "#2C3E50",
                    "key": "#34495E", "sub1": "#7FB3D3", "sub2": "#E8B4A5"},
        "tags": ["노르딕", "미니멀", "차분", "라이프스타일", "인테리어"],
        "fit": ["marketing_pr", "consulting"],
        "mood": "calm, thoughtful, understated",
    },
    "healthcare_mint": {
        "name": "Healthcare Mint",
        "desc": "헬스케어 민트 — 의료/웰니스",
        "colors": {"bg": "#F0F9F7", "text": "#1A3A3A",
                    "key": "#14B8A6", "sub1": "#7DD3C0", "sub2": "#F97316"},
        "tags": ["의료", "헬스", "웰니스", "신뢰", "케어"],
        "fit": ["public", "consulting"],
        "mood": "healthy, trustworthy, caring",
    },
    "corporate_blue": {
        "name": "Corporate Blue",
        "desc": "코퍼레이트 블루 — 전통 기업",
        "colors": {"bg": "#FFFFFF", "text": "#1E293B",
                    "key": "#002C5F", "sub1": "#00AAD2", "sub2": "#E63312"},
        "tags": ["기업", "전통", "공식", "신뢰", "보수"],
        "fit": ["consulting", "it_system", "public"],
        "mood": "reliable, professional, traditional",
    },

    # ── VIBRANT ────────────────────────────────────────────────────
    "tech_gradient": {
        "name": "Tech Gradient",
        "desc": "테크 그라디언트 — 스타트업/SaaS",
        "colors": {"bg": "#F8FAFF", "text": "#1E1B4B",
                    "key": "#6366F1", "sub1": "#A855F7", "sub2": "#EC4899"},
        "tags": ["테크", "스타트업", "그라디언트", "모던", "생동감"],
        "fit": ["it_system", "marketing_pr"],
        "mood": "innovative, dynamic, youthful",
    },
    "sunset_coral": {
        "name": "Sunset Coral",
        "desc": "선셋 코럴 — 크리에이티브/캠페인",
        "colors": {"bg": "#FFF8F5", "text": "#2D1B1B",
                    "key": "#FF6B6B", "sub1": "#FFB88C", "sub2": "#4ECDC4"},
        "tags": ["캠페인", "크리에이티브", "웜", "생동감", "이벤트"],
        "fit": ["event", "marketing_pr"],
        "mood": "energetic, warm, creative",
    },
    "youth_pop": {
        "name": "Youth Pop",
        "desc": "유스 팝 — 10-20대 타겟",
        "colors": {"bg": "#FFFCE8", "text": "#1F1147",
                    "key": "#7C3AED", "sub1": "#F59E0B", "sub2": "#EC4899"},
        "tags": ["젊음", "MZ", "팝", "캠페인", "SNS"],
        "fit": ["marketing_pr", "event"],
        "mood": "playful, bold, social",
    },
    "food_warm": {
        "name": "Food Warm",
        "desc": "푸드 웜 — F&B/요식업",
        "colors": {"bg": "#FFF5E8", "text": "#3D2817",
                    "key": "#D2691E", "sub1": "#8B4513", "sub2": "#228B22"},
        "tags": ["F&B", "요식", "웜", "전통", "품질"],
        "fit": ["marketing_pr", "event"],
        "mood": "appetizing, warm, inviting",
    },

    # ── SPECIALIZED ────────────────────────────────────────────────
    "heritage_gold": {
        "name": "Heritage Gold",
        "desc": "헤리티지 골드 — 공공기관/전통",
        "colors": {"bg": "#FAF5E8", "text": "#2C1810",
                    "key": "#8B1A1A", "sub1": "#C5973E", "sub2": "#4A5D23"},
        "tags": ["전통", "공공", "한국", "문화", "유산"],
        "fit": ["public", "consulting"],
        "mood": "authoritative, traditional, cultural",
    },
    "industrial_steel": {
        "name": "Industrial Steel",
        "desc": "인더스트리얼 스틸 — 제조/산업",
        "colors": {"bg": "#1F2937", "text": "#F9FAFB",
                    "key": "#F59E0B", "sub1": "#EF4444", "sub2": "#10B981"},
        "tags": ["제조", "산업", "다크", "강인", "중후"],
        "fit": ["it_system", "consulting"],
        "mood": "industrial, robust, technical",
    },
    "mono_elegant": {
        "name": "Mono Elegant",
        "desc": "모노 엘레강트 — 고급 컨설팅",
        "colors": {"bg": "#FFFFFF", "text": "#0A0A0A",
                    "key": "#404040", "sub1": "#737373", "sub2": "#B91C1C"},
        "tags": ["모노", "흑백", "엘레강트", "프리미엄", "미니멀"],
        "fit": ["consulting"],
        "mood": "sophisticated, minimal, authoritative",
    },
    "fintech_purple": {
        "name": "Fintech Purple",
        "desc": "핀테크 퍼플 — 디지털 금융",
        "colors": {"bg": "#0F0A1A", "text": "#FFFFFF",
                    "key": "#8B5CF6", "sub1": "#3B82F6", "sub2": "#10B981"},
        "tags": ["핀테크", "디지털", "금융", "테크", "모던"],
        "fit": ["it_system", "consulting"],
        "mood": "innovative, trustworthy, digital",
    },
    "event_gala": {
        "name": "Event Gala",
        "desc": "이벤트 갈라 — 시상식/행사",
        "colors": {"bg": "#16080E", "text": "#FFF9E6",
                    "key": "#D4AF37", "sub1": "#C0392B", "sub2": "#F4D06F"},
        "tags": ["이벤트", "갈라", "시상식", "프리미엄", "골드"],
        "fit": ["event"],
        "mood": "prestigious, celebratory, golden",
    },
    "eco_fresh": {
        "name": "Eco Fresh",
        "desc": "에코 프레시 — 친환경 캠페인",
        "colors": {"bg": "#F3F9F1", "text": "#1A3A1F",
                    "key": "#22C55E", "sub1": "#84CC16", "sub2": "#F59E0B"},
        "tags": ["친환경", "에코", "ESG", "자연", "캠페인"],
        "fit": ["public", "marketing_pr"],
        "mood": "fresh, sustainable, optimistic",
    },

    # ── 게임 IP 판타지 ─────────────────────────────────────────────
    "fantasy_mystic": {
        "name": "Fantasy Mystic",
        "desc": "판타지 미스틱 — 마비노기/RPG류 게임 쇼케이스",
        "colors": {"bg": "#14102A", "text": "#F5F0FF",
                    "key": "#9D7DFF", "sub1": "#D4AF66", "sub2": "#F0C4E0"},
        "tags": ["판타지", "게임", "RPG", "마법", "신비", "감성",
                  "다크", "보라", "골드", "쇼케이스", "스토리"],
        "fit": ["event", "marketing_pr"],
        "mood": "mystical, enchanting, narrative",
        "gradient_bg": ("#1A1333", "#14102A", "#0F0B24"),
        "gradient_key": ("#9D7DFF", "#F0C4E0"),
    },

    # ── LAON 실전 수주 레퍼런스 ────────────────────────────────────
    "vaetki_pastel": {
        "name": "VAETKI Pastel",
        "desc": "VAETKI Commerce 쇼케이스 실전 수주 팔레트 — 파스텔 그라디언트 에디토리얼",
        "colors": {"bg": "#FFFFFF", "text": "#0D0D15",
                    "key": "#6868F1", "sub1": "#DD6495", "sub2": "#B667DD"},
        "tags": ["파스텔", "에디토리얼", "쇼케이스", "AI", "커머스",
                  "인플루언서", "그라디언트", "라이트", "프리미엄"],
        "fit": ["marketing_pr", "event"],
        "mood": "elegant, modern, refined",
        "gradient_bg": ("#DEE2FB", "#FFFFFF", "#C2CAF8"),   # 3-stop 파스텔
        "gradient_key": ("#6868F1", "#DD6495"),               # 보라블루→핑크 헤드라인용
    },
}


def list_palettes():
    """팔레트 라이브러리 전체 출력."""
    print(f"\n=== 팔레트 라이브러리 ({len(PALETTE_LIBRARY)}종) ===")
    for key, p in PALETTE_LIBRARY.items():
        colors = p["colors"]
        print(f"\n  [{key}]  {p['name']}")
        print(f"    {p['desc']}")
        print(f"    태그: {', '.join(p['tags'][:5])}")
        print(f"    색상: bg={colors['bg']}  text={colors['text']}  "
               f"key={colors['key']}  sub1={colors['sub1']}  sub2={colors['sub2']}")
    print()


def recommend_palettes(project_type=None, industry=None, keywords=None,
                        mood=None, top_n=3):
    """프로젝트 특성 기반 팔레트 추천.

    Args:
        project_type: "marketing_pr" | "event" | "it_system" | "public" | "consulting"
        industry: 업종 키워드 (예: "game", "finance", "food", "healthcare")
        keywords: 추가 키워드 리스트 (예: ["다크", "모던"])
        mood: 분위기 설명 (예: "luxurious", "bold")
        top_n: 상위 N개 반환

    Returns:
        [{"key": str, "score": int, "palette": dict, "reason": str}, ...]
    """
    keywords = keywords or []
    # 정규화
    search_terms = set()
    if industry:
        search_terms.add(industry.lower())
    for k in keywords:
        search_terms.add(k.lower())
    if mood:
        search_terms.update(mood.lower().split())

    scored = []
    for key, pal in PALETTE_LIBRARY.items():
        score = 0
        reasons = []

        # project_type 매칭
        if project_type and project_type in pal.get("fit", []):
            score += 3
            reasons.append(f"'{project_type}' 적합")

        # 태그 매칭
        tag_set = {t.lower() for t in pal.get("tags", [])}
        matched_tags = search_terms & tag_set
        if matched_tags:
            score += len(matched_tags) * 2
            reasons.append(f"태그 매칭: {', '.join(matched_tags)}")

        # mood 매칭
        if mood and any(m in pal.get("mood", "").lower() for m in mood.lower().split()):
            score += 2
            reasons.append("분위기 일치")

        # desc에 키워드 포함?
        desc_lower = pal.get("desc", "").lower()
        for term in search_terms:
            if term in desc_lower:
                score += 1

        if score > 0:
            scored.append({
                "key": key,
                "score": score,
                "palette": pal,
                "reason": "; ".join(reasons) if reasons else "부분 매칭",
            })

    # 점수 내림차순
    scored.sort(key=lambda x: -x["score"])
    return scored[:top_n]


def print_recommendations(recs):
    """추천 결과 콘솔 출력."""
    if not recs:
        print("\n  [추천 없음] 키워드를 추가하거나 list_palettes()로 전체 확인\n")
        return
    print(f"\n=== 팔레트 추천 (상위 {len(recs)}) ===")
    for i, r in enumerate(recs, 1):
        p = r["palette"]
        c = p["colors"]
        print(f"\n  {i}위. [{r['key']}]  {p['name']}  (score={r['score']})")
        print(f"       {p['desc']}")
        print(f"       근거: {r['reason']}")
        print(f"       색상: bg={c['bg']}  text={c['text']}  key={c['key']}  "
               f"sub1={c['sub1']}  sub2={c['sub2']}")
    print(f"\n  적용: apply_from_library('<key>')  "
           f"예: apply_from_library('{recs[0]['key']}')\n")


def apply_from_library(palette_key):
    """라이브러리에서 팔레트 키로 적용.

    Example:
        apply_from_library("finance_navy")
    """
    if palette_key not in PALETTE_LIBRARY:
        raise ValueError(f"Unknown palette '{palette_key}'. "
                           f"Available: {list(PALETTE_LIBRARY.keys())}")
    pal = PALETTE_LIBRARY[palette_key]
    apply_palette(**pal["colors"])
    return pal


def start_deck_interactive(project_type=None, industry=None,
                             keywords=None, mood=None, auto_apply=True):
    """새 덱 시작 시 원샷 — 추천 + 자동 적용.

    Args:
        project_type, industry, keywords, mood: recommend_palettes 참조
        auto_apply: True면 1위 자동 적용

    Returns:
        적용된 팔레트 dict
    """
    propose_palette()
    recs = recommend_palettes(project_type=project_type, industry=industry,
                                keywords=keywords, mood=mood, top_n=3)
    print_recommendations(recs)
    if auto_apply and recs:
        top = recs[0]
        print(f"  → 1위 팔레트 자동 적용: {top['key']}\n")
        apply_from_library(top["key"])
        return PALETTE_LIBRARY[top["key"]]
    return None


def prune_slide_layouts(prs, keep_indexes=(6,)):
    """기본 슬라이드 레이아웃 정리 — keep_indexes 외 레이아웃 제거.

    python-pptx Presentation()은 기본 9개 레이아웃을 생성하지만 대부분 안 씀.
    layout[6] (Blank)만 남기고 제거해서 "Insert new slide" UI 깔끔하게.

    Args:
        keep_indexes: 유지할 레이아웃 인덱스 튜플 (기본: (6,) = Blank만)
    """
    try:
        from lxml import etree
    except ImportError:
        return prs
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

    master = prs.slide_master
    master_el = master._element
    sldLayoutIdLst = master_el.find(f'{{{ns_p}}}sldLayoutIdLst')
    if sldLayoutIdLst is None:
        return prs

    layout_ids = list(sldLayoutIdLst.findall(f'{{{ns_p}}}sldLayoutId'))
    keep_set = set(keep_indexes)

    # 레이아웃 사용 여부 확인 — 슬라이드가 참조하는 레이아웃은 반드시 유지
    used_layouts = set()
    for slide in prs.slides:
        used_layouts.add(slide.slide_layout.partname)

    removed = 0
    for i, layout_el in enumerate(layout_ids):
        if i in keep_set:
            continue
        # 이 레이아웃이 실제 슬라이드에 사용 중인지 확인
        rId = layout_el.get(f'{{{ns_r}}}id')
        # master.part.rels에서 target 확인
        try:
            layout_part = master.part.related_parts.get(rId)
            if layout_part and layout_part.partname in used_layouts:
                continue   # 사용 중 → 유지
        except Exception:
            pass
        # 제거: sldLayoutIdLst에서 엔트리, rels, part 순으로
        try:
            sldLayoutIdLst.remove(layout_el)
            master.part.drop_rel(rId)
            removed += 1
        except Exception:
            pass
    return prs


def setup_editorial_deck(prs, *, bg_color=None, gradient=None, prune=True):
    """에디토리얼 다크 덱 원샷 설정.

    - 마스터 배경 설정 (모든 슬라이드 자동 상속)
    - 사용 안 하는 기본 레이아웃 제거 (옵션)

    Args:
        bg_color: 마스터 배경색. None이면 tok("surface/darker")
        gradient: (c1, c2) 그라디언트 배경
        prune: True면 Blank 레이아웃 외 제거

    Example:
        prs = new_presentation()
        setup_editorial_deck(prs)
        # 이후 new_slide()는 자동으로 마스터 배경 상속
    """
    if gradient is None and bg_color is None:
        bg_color = tok("surface/darker")
    setup_master_background(prs, color=bg_color, gradient=gradient)
    if prune:
        prune_slide_layouts(prs, keep_indexes=(6,))
    return prs


def _master_gradient(master, c1, c2, angle=5400000):
    """마스터 배경에 그라디언트 설정 (XML 직접 조작)."""
    try:
        from lxml import etree
    except ImportError:
        return
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    cSld = master._element.find(f'{{{ns_p}}}cSld')
    # bg 요소 확보
    bg = cSld.find(f'{{{ns_p}}}bg')
    if bg is None:
        bg = etree.SubElement(cSld, f'{{{ns_p}}}bg')
        cSld.insert(0, bg)
    # 기존 자식 제거
    for child in list(bg):
        bg.remove(child)
    bgPr = etree.SubElement(bg, f'{{{ns_p}}}bgPr')
    gradFill = etree.SubElement(bgPr, f'{{{ns}}}gradFill',
                                 flip='none', rotWithShape='1')
    gsLst = etree.SubElement(gradFill, f'{{{ns}}}gsLst')
    for pos, rgb in [(0, c1), (100000, c2)]:
        gs = etree.SubElement(gsLst, f'{{{ns}}}gs', pos=str(pos))
        hex_val = f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}" \
                    if isinstance(rgb, tuple) else str(rgb).lstrip('#').upper()
        etree.SubElement(gs, f'{{{ns}}}srgbClr', val=hex_val)
    etree.SubElement(gradFill, f'{{{ns}}}lin', ang=str(angle), scaled='1')


# ═══════════════════════════════════════════════════════════════
#  3. 기본 도형
# ═══════════════════════════════════════════════════════════════

def _safe_int(val):
    """EMU 값을 안전하게 int로 변환 (float EMU → int, 비정상 값 방지).

    v3.8: CW*0.5 등 EMU 스케일 float가 add_shape()에 직접 전달될 때
    XML에 소수점이 포함되어 PPT 파손 방지.
    """
    if isinstance(val, float):
        return int(val)
    return val


def R(s, l, t, w, h, f=None, lc=None, lw=1):
    """직각 사각형

    Args:
        f: 채우기 색상, None이면 투명
        lc: 테두리 색상, None이면 없음
        lw: 테두리 두께(pt)
    """
    l, t, w, h = _safe_int(l), _safe_int(t), _safe_int(w), _safe_int(h)
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.line.fill.background()
    if f:
        sh.fill.solid()
        sh.fill.fore_color.rgb = f
    else:
        sh.fill.background()
    if lc:
        sh.line.color.rgb = lc
        sh.line.width = Pt(lw)
    return sh


def BOX(s, l, t, w, h, f, text="", sz=13, tc=None, b=False):
    """텍스트가 중앙 정렬된 직각 박스 (도형 + 텍스트 일체형)

    도형 위에 텍스트가 있을 때 항상 이 함수 사용.
    텍스트는 수평/수직 모두 중앙 정렬됨.
    """
    if tc is None:
        tc = C["white"]
    l, t, w, h = _safe_int(l), _safe_int(t), _safe_int(w), _safe_int(h)
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = f
    sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Pt(6))
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = tc
    p.font.bold = b
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh


def OBOX(s, l, t, w, h, text="", sz=13, tc=None, b=False, lc=None):
    """테두리만 있는 박스 (아웃라인 박스) — 텍스트 중앙 정렬

    배경 투명, 테두리만 표시.
    """
    if tc is None:
        tc = C["dark"]
    if lc is None:
        lc = C["primary"]
    l, t, w, h = _safe_int(l), _safe_int(t), _safe_int(w), _safe_int(h)
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.background()
    sh.line.color.rgb = lc
    sh.line.width = Pt(1.5)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Pt(6))
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = tc
    p.font.bold = b
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh


# ═══════════════════════════════════════════════════════════════
#  4. 텍스트 헬퍼
# ═══════════════════════════════════════════════════════════════

def T(s, l, t, w, h, text, sz=13, c=None, b=False, al=PP_ALIGN.LEFT, ls=1.4,
      fn=None):
    """단일 스타일 텍스트

    Args:
        fn: 폰트 이름 (None이면 FONT 사용, FONT_W["semibold"] 등 사용 가능)
    """
    if c is None:
        c = C["dark"]
    l, t, w, h = _safe_int(l), _safe_int(t), _safe_int(w), _safe_int(h)
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Pt(0))
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = fn or FONT
    p.alignment = al
    p.line_spacing = Pt(int(sz * ls))
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    return tb


def RT(s, l, t, w, h, parts, al=PP_ALIGN.LEFT, ls=1.4):
    """리치 텍스트 — [(text, size, color, bold), ...]"""
    l, t, w, h = _safe_int(l), _safe_int(t), _safe_int(w), _safe_int(h)
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Pt(0))
    p = tf.paragraphs[0]
    p.alignment = al
    max_sz = 13
    for text, sz, c, b in parts:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(sz)
        r.font.color.rgb = c
        r.font.bold = b
        r.font.name = FONT
        if sz > max_sz:
            max_sz = sz
    p.line_spacing = Pt(int(max_sz * ls))
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    return tb


def MT(s, l, t, w, h, lines, sz=13, c=None, b=False, al=PP_ALIGN.LEFT, ls=1.6, bul=False):
    """멀티라인 텍스트"""
    if c is None:
        c = C["dark"]
    l, t, w, h = _safe_int(l), _safe_int(t), _safe_int(w), _safe_int(h)
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " + ln) if bul else ln
        p.font.size = Pt(sz)
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = FONT
        p.alignment = al
        p.line_spacing = Pt(int(sz * ls))
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    return tb


# ═══════════════════════════════════════════════════════════════
#  5. 슬라이드 공통 요소
# ═══════════════════════════════════════════════════════════════

def bg(s, c):
    """단색 배경"""
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c


def gradient_bg(s, c1, c2, angle=270.0):
    """그래디언트 배경"""
    fill = s.background.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_angle = angle


def gradient_shape(shape, c1, c2, angle=270.0):
    """도형에 그라디언트 채우기 적용 (v3.6)"""
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_angle = angle
    return shape


# 그라디언트 프리셋 (v3.6)
GRAD = {
    "cover":     lambda: (C["primary_dark"], C["primary"]),     # 커버용 다크→블루
    "closing":   lambda: (C["primary"], darken(C["primary"], 0.15)),  # 클로징
    "highlight": lambda: (C["primary"], C["secondary_dark"]),   # 하이라이트 박스
    "section":   lambda: (darken(C["dark"], 0.3), C["dark"]),   # 섹션 구분자
    "teal":      lambda: (darken(C["teal"], 0.2), C["teal"]),   # 틸 계열
    "accent":    lambda: (C["accent"], lighten(C["accent"], 0.2)),  # 강조 계열
}


# 배경 프리셋 (v3.8)
def bg_preset(s, preset_name):
    """명명된 배경 프리셋 적용 — 배경 다양성 확보"""
    presets = {
        "white":         lambda: bg(s, C["white"]),
        "light":         lambda: bg(s, C["light"]),
        "warm_light":    lambda: bg(s, C["primary_light"]),
        "cool_light":    lambda: bg(s, C["secondary_light"]),
        "subtle_blue":   lambda: bg(s, lighten(C["primary"], 0.92)),
        "dark":          lambda: bg(s, C["dark"]),
        "gradient_dark": lambda: gradient_bg(s, darken(C["dark"], 0.3), C["dark"]),
        "gradient_blue": lambda: gradient_bg(s, C["primary_dark"], C["primary"]),
        "gradient_teal": lambda: gradient_bg(s, darken(C["teal"], 0.2), C["teal"]),
    }
    fn = presets.get(preset_name, presets["white"])
    fn()


def set_char_spacing(tb, spacing=200):
    """자간 설정 (100 = 1pt)"""
    try:
        for p in tb.text_frame.paragraphs:
            for r in p.runs:
                rPr = r._r.get_or_add_rPr()
                rPr.set('spc', str(spacing))
    except Exception:
        pass


def PN(s, n):
    """페이지 번호 (우하단) — 중복 방지: 이미 동일 위치에 페이지번호 존재 시 스킵"""
    pn_left = SW - Inches(1.0)
    pn_top = SH - Inches(0.4)
    # 중복 검사: 동일 위치에 이미 텍스트 존재하면 스킵
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == str(n):
            if abs(sh.left - pn_left) < Inches(0.05) and abs(sh.top - pn_top) < Inches(0.05):
                return  # 이미 존재 — 스킵
    T(s, pn_left, pn_top, Inches(0.7), Inches(0.25),
      str(n), sz=SZ["source"], c=C["gray"], al=PP_ALIGN.RIGHT)


def SRC(s, text):
    """출처 표기 (좌하단) — McKinsey 스타일"""
    T(s, ML, SH - Inches(0.4), Inches(8), Inches(0.25),
      f"Source: {text}", sz=SZ["source"], c=C["gray"])


def TB(s, text, pg=None, src=None):
    """Action Title 상단바 (McKinsey 스타일)

    - 좌측 프라이머리 라인
    - Action Title = 인사이트 기반 제목 (문장형)
    - 하단 구분선
    """
    R(s, Inches(0), Inches(0), Inches(0.08), SH, f=C["primary"])
    T(s, ML, Inches(0.35), CW, Inches(0.55),
      text, sz=SZ["action"], c=C["dark"], b=True,
      fn=FONT_W["semibold"])
    R(s, ML, Inches(0.88), CW, Pt(1), f=C["lgray"])
    if pg:
        PN(s, pg)
    if src:
        SRC(s, src)


def WB(s, theme_key, win_themes, x=None, y=None, w=None):
    """Win Theme 뱃지"""
    if x is None:
        x = ML
    if y is None:
        y = SH - Inches(0.9)
    if w is None:
        w = Inches(4.5)
    BOX(s, x, y, w, Inches(0.35), C["teal"],
        f"Win Theme  |  {win_themes.get(theme_key, theme_key)}",
        sz=SZ["caption"], tc=C["white"], b=True)


def IMG(s, l, t, w, h, desc="이미지 영역"):
    """이미지 플레이스홀더"""
    l, t, w, h = _safe_int(l), _safe_int(t), _safe_int(w), _safe_int(h)
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = C["light"]
    sh.line.color.rgb = C["lgray"]
    sh.line.width = Pt(1)
    sh.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = desc
    r.font.size = Pt(SZ["caption"])
    r.font.color.rgb = C["gray"]
    r.font.name = FONT
    return sh


# ═══════════════════════════════════════════════════════════════
#  6. 도식화 헬퍼 — 구조를 잡는 함수들
# ═══════════════════════════════════════════════════════════════

def _flow_arrow(s, x_emu, y_emu, gap_w, box_h, cy_offset=None):
    """FLOW/STEP_ARROW 사이 ▶ 삼각형 화살표 (그래픽, 수직 가운데 정렬)

    Args:
        x_emu: 화살표 영역 시작 X (EMU)
        y_emu: 박스 시작 Y (EMU)
        gap_w: 화살표 영역 너비 (인치 float)
        box_h: 박스 높이 (인치 float)
        cy_offset: Y 오프셋 (인치 float, None이면 box_h 기준 가운데)
    """
    tri_w = 0.18   # 삼각형 너비
    tri_h = 0.22   # 삼각형 높이
    # 수평 가운데
    ax = float(x_emu / 914400) + (gap_w - tri_w) / 2
    # 수직 가운데 (박스 높이 기준)
    if cy_offset is not None:
        ay = float(y_emu / 914400) + cy_offset + (box_h - tri_h) / 2
    else:
        ay = float(y_emu / 914400) + (box_h - tri_h) / 2
    sh = s.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(ax), Inches(ay), Inches(tri_w), Inches(tri_h))
    sh.rotation = 90.0  # 오른쪽 방향 ▶
    fill = sh.fill
    fill.solid()
    fill.fore_color.rgb = C["lgray"]
    sh.line.fill.background()  # 테두리 없음


def FLOW(s, items, y=None, h=None, colors=None):
    """프로세스 플로우 (가로 화살표 연결)

    Args:
        items: [("제목", "설명"), ...] 리스트 (3~5개 권장)
        y: 시작 Y 위치 (기본 1.2") — raw float 또는 Inches 객체
        h: 박스 높이 (기본 1.2") — raw float 또는 Inches 객체
        colors: 색상 리스트, None이면 primary 계열 자동

    Returns:
        float: 실제 전체 높이(인치). desc 있으면 h + 0.9, 없으면 h.
        ⚠ VStack 예약 시 이 반환값을 사용할 것.
    """
    y = _ensure_emu(y)
    h = _ensure_emu(h)
    if y is None:
        y = Inches(1.2)
    if h is None:
        h = Inches(1.2)
    n = len(items)
    if colors is None:
        palette = [C["primary"], C["secondary"], C["teal"], C["accent"], C["green"]]
        colors = [palette[i % len(palette)] for i in range(n)]
    arrow_w = 0.35
    total = float(CW / 914400)
    box_w = (total - arrow_w * (n - 1)) / n
    has_desc = False
    for i, (title, desc) in enumerate(items):
        x = ML + Inches((box_w + arrow_w) * i)
        BOX(s, x, y, Inches(box_w), h, colors[i],
            title, sz=SZ["body"], tc=C["white"], b=True)
        if desc:
            has_desc = True
            T(s, x, y + h + Inches(0.1), Inches(box_w), Inches(0.4),
              desc, sz=SZ["body_sm"], c=C["gray"], al=PP_ALIGN.CENTER, ls=1.3)
        if i < n - 1:
            _flow_arrow(s, x + Inches(box_w), y, arrow_w, float(h / 914400))
    h_in = float(h / 914400)
    return h_in + 0.6 if has_desc else h_in


def COLS(s, items, y=None, h=None, colors=None, show_header=True, shadow=True,
         x=None, w=None):
    """N-컬럼 카드 레이아웃 (v3.8 — 그림자 + x/w 커스텀 영역)

    Args:
        items: [{"title": "제목", "body": ["항목1", "항목2"]}, ...] (2~4개)
        y: 시작 Y (기본 1.2") — raw float 또는 Inches 객체
        h: 카드 높이 (기본 3.5") — raw float 또는 Inches 객체
        colors: 헤더 색상 리스트
        show_header: False면 헤더 없이 아웃라인 박스
        shadow: 그림자 적용 여부
        x: 시작 X (기본 ML) — 분할 레이아웃 시 우측 영역 지정 가능
        w: 전체 너비 (기본 CW) — x와 함께 사용하여 제한된 영역에 배치
    """
    y = _ensure_emu(y)
    h = _ensure_emu(h)
    x_start = _ensure_emu(x) if x is not None else None
    w_total = _ensure_emu(w) if w is not None else None
    if y is None:
        y = Inches(1.2)
    if h is None:
        h = Inches(3.5)
    if x_start is None:
        x_start = ML
    if w_total is None:
        w_total = CW
    # ★ 경계 클램핑: 슬라이드 하단(0.5" 마진) 넘지 않도록 높이 자동 조정
    y_in = float(y / 914400)
    h_in_raw = float(h / 914400)
    max_h = 7.0 - y_in    # 하단 0.5" 여유
    if h_in_raw > max_h and max_h > 1.0:
        h = Inches(max_h)
    n = len(items)
    gap = 0.2
    total = float(w_total / 914400)
    col_w = (total - gap * (n - 1)) / n
    if colors is None:
        palette = [C["primary"], C["secondary"], C["teal"], C["accent"]]
        colors = [palette[i % len(palette)] for i in range(n)]
    header_h = Inches(0.5)
    for i, item in enumerate(items):
        cx = x_start + Inches((col_w + gap) * i)   # ★ v3.8: x_start 기반 (커스텀 영역 지원)
        title = item.get("title", "")
        body = item.get("body", [])
        if show_header:
            # 카드 배경 (그림자)
            card_sh = R(s, cx, y, Inches(col_w), h,
                        f=C["card_bg"], lc=C["card_border"], lw=0.5)
            if shadow:
                add_shadow(card_sh, preset="card")
            BOX(s, cx, y, Inches(col_w), header_h, colors[i],
                title, sz=SZ["body"], tc=C["white"], b=True)
            MT(s, cx + Inches(0.15), y + header_h + Inches(0.1),
               Inches(col_w - 0.3), h - header_h - Inches(0.2),
               body, sz=SZ["body_sm"], bul=True)
        else:
            header_h_ns = Inches(0.5)
            OBOX(s, cx, y, Inches(col_w), header_h_ns, title, sz=SZ["body"],
                 tc=C["primary"], b=True, lc=colors[i])
            R(s, cx, y + header_h_ns, Inches(col_w), h - header_h_ns,
              lc=colors[i], lw=1.5)
            MT(s, cx + Inches(0.15), y + header_h_ns + Inches(0.1),
               Inches(col_w - 0.3), h - header_h_ns - Inches(0.2),
               body, sz=SZ["body_sm"], bul=True)


def PYRAMID(s, levels, y=None, w_max=None, h_total=None):
    """피라미드 구조 (위가 좁고 아래가 넓음 — McKinsey Pyramid Principle)

    Args:
        levels: [("최상위 메시지", color), ("중간", color), ("하단", color)]
                위에서 아래 순서
        y: 시작 Y (기본 1.2") — raw float 또는 Inches 객체
        w_max: 최대 너비 (기본 CW)
        h_total: 전체 높이 (기본 4.5")
    """
    y = _ensure_emu(y)
    if y is None:
        y = Inches(1.2)
    if w_max is None:
        w_max = float(CW / 914400)
    if h_total is None:
        h_total = 4.5
    n = len(levels)
    level_h = h_total / n
    center = float(ML / 914400) + w_max / 2
    for i, (text, clr) in enumerate(levels):
        ratio = 0.4 + 0.6 * (i / max(n - 1, 1))
        lw = w_max * ratio
        lx = center - lw / 2
        ly = float(y / 914400) + level_h * i
        BOX(s, Inches(lx), Inches(ly), Inches(lw), Inches(level_h - 0.08),
            clr, text, sz=SZ["body"], tc=C["white"], b=True)


def MATRIX(s, quadrants, x_label="", y_label="", y_start=None):
    """2x2 매트릭스 (McKinsey 전략 매트릭스)

    Args:
        quadrants: [("좌상", color), ("우상", color), ("좌하", color), ("우하", color)]
        x_label: X축 라벨
        y_label: Y축 라벨
        y_start: 시작 Y
    """
    if y_start is None:
        y_start = Inches(1.2)
    total = float(CW / 914400)
    label_w = 0.5
    gap = 0.1
    cell_w = (total - label_w - gap * 2) / 2
    cell_h = 2.2
    ox = float(ML / 914400) + label_w + gap
    oy = float(y_start / 914400)
    # Y축 라벨
    if y_label:
        T(s, ML, y_start, Inches(label_w), Inches(cell_h * 2 + gap),
          y_label, sz=SZ["body_sm"], c=C["gray"], b=True,
          al=PP_ALIGN.CENTER)
    # X축 라벨
    if x_label:
        T(s, Inches(ox), Inches(oy + cell_h * 2 + gap + 0.1),
          Inches(cell_w * 2 + gap), Inches(0.3),
          x_label, sz=SZ["body_sm"], c=C["gray"], b=True,
          al=PP_ALIGN.CENTER)
    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for idx, (text, clr) in enumerate(quadrants[:4]):
        col, row = positions[idx]
        bx = Inches(ox + col * (cell_w + gap))
        by = Inches(oy + row * (cell_h + gap))
        BOX(s, bx, by, Inches(cell_w), Inches(cell_h), clr,
            text, sz=SZ["body"], tc=C["white"], b=True)


def TABLE(s, headers, rows, y=None, col_widths=None):
    """데이터 테이블 — 실제 PPTX Table 객체 (테두리 + 배경 컬러)

    Args:
        headers: ["항목", "AS-IS", "TO-BE"]
        rows: [["항목1", "현재", "목표"], ...]
        y: 시작 Y — raw float 또는 Inches 객체
        col_widths: 열 너비 비율 리스트, None이면 균등

    Returns:
        float: 실제 전체 높이(인치).
    """
    from pptx.oxml.ns import qn
    from lxml import etree
    y = _ensure_emu(y)
    if y is None:
        y = Inches(1.2)
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header
    row_h_in = 0.45
    total_h = Inches(row_h_in * n_rows)
    # 열 너비 계산
    total_w_in = float(CW / 914400)
    if col_widths is None:
        w_list = [total_w_in / n_cols] * n_cols
    else:
        ratio_sum = sum(col_widths)
        w_list = [total_w_in * (r / ratio_sum) for r in col_widths]
    # 테이블 shape 생성
    table_shape = s.shapes.add_table(n_rows, n_cols, ML, y, CW, total_h)
    tbl = table_shape.table
    # 기본 테이블 스타일 제거 (수동 서식 적용)
    tbl_pr = tbl._tbl.tblPr
    # bandRow/firstRow 등 비활성화
    tbl_pr.set("bandRow", "0")
    tbl_pr.set("bandCol", "0")
    tbl_pr.set("firstRow", "0")
    tbl_pr.set("lastRow", "0")
    # 열 너비 설정
    for j, w in enumerate(w_list):
        tbl.columns[j].width = Inches(w)
    # 셀 서식 헬퍼
    def _fmt_cell(cell, text, bg_color, txt_color, bold=False):
        cell.text = str(text)
        # 폰트
        for para in cell.text_frame.paragraphs:
            para.font.name = FONT
            para.font.size = Pt(SZ["body_sm"])
            para.font.color.rgb = txt_color
            para.font.bold = bold
            para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        # 배경색
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
        # 테두리 (4면)
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        border_color = C.get("card_border", RGBColor(220, 220, 220))
        for edge in ["lnL", "lnR", "lnT", "lnB"]:
            ln = etree.SubElement(tcPr, qn(f"a:{edge}"))
            ln.set("w", "6350")  # 0.5pt in EMU
            ln.set("cmpd", "sng")
            sf = etree.SubElement(ln, qn("a:solidFill"))
            clr = etree.SubElement(sf, qn("a:srgbClr"))
            clr.set("val", f"{border_color[0]:02X}{border_color[1]:02X}{border_color[2]:02X}"
                     if isinstance(border_color, (tuple, list))
                     else f"{border_color.red:02X}{border_color.green:02X}{border_color.blue:02X}")
        # 마진
        tc.marL = Inches(0.08)
        tc.marR = Inches(0.08)
        tc.marT = Inches(0.04)
        tc.marB = Inches(0.04)
    # 헤더 행
    for j, hdr in enumerate(headers):
        _fmt_cell(tbl.cell(0, j), hdr, C["primary"], C["white"], bold=True)
    # 데이터 행
    for i, row in enumerate(rows):
        bg = C.get("card_bg", RGBColor(248, 249, 250)) if i % 2 == 0 else C["white"]
        for j, cell_text in enumerate(row):
            _fmt_cell(tbl.cell(i + 1, j), cell_text, bg, C["dark"])
    return row_h_in * n_rows


def HIGHLIGHT(s, text, sub="", y=None, color=None, grad=False):
    """핵심 메시지 강조 박스 (v3.6 — 그라디언트 옵션 + 라운드)

    Args:
        text: 메인 메시지
        sub: 보조 텍스트
        y: Y 위치 — raw float 또는 Inches 객체
        color: 배경색
        grad: True면 그라디언트 적용 (color→secondary_dark)

    Returns:
        float: 실제 높이(인치). sub 있으면 1.2, 없으면 0.8.
    """
    y = _ensure_emu(y)
    if y is None:
        y = Inches(1.2)
    if color is None:
        color = C["primary"]
    if sub:
        h = Inches(1.2)
        sh = RBOX(s, ML, y, CW, h, color, radius=0.04)
        if grad:
            gradient_shape(sh, color, darken(color, 0.15), angle=0.0)
        T(s, ML + Inches(0.3), y + Inches(0.1), CW - Inches(0.6), Inches(0.4),
          text, sz=SZ["subtitle"], c=C["white"], b=True, al=PP_ALIGN.CENTER,
          fn=FONT_W["semibold"])
        T(s, ML + Inches(0.3), y + Inches(0.6), CW - Inches(0.6), Inches(0.45),
          sub, sz=SZ["body_sm"], c=C["white"], al=PP_ALIGN.CENTER,
          fn=FONT_W["light"])
        return 1.2
    else:
        h = Inches(0.8)
        sh = RBOX(s, ML, y, CW, h, color, text, sz=SZ["subtitle"],
                  tc=C["white"], b=True, radius=0.04)
        if grad:
            gradient_shape(sh, color, darken(color, 0.15), angle=0.0)
        return 0.8


def KPIS(s, items, y=None, h=None, shadow=True):
    """KPI 카드 그리드 (v3.6 — 라운드 + 그림자 + 타이포 강화)

    Args:
        items: [{"value": "+30%", "label": "팔로워 성장", "basis": "산출근거"}, ...]
        y: 시작 Y — raw float 또는 Inches 객체
        h: 카드 높이 — raw float 또는 Inches 객체
        shadow: 그림자 적용 여부
    """
    y = _ensure_emu(y)
    h = _ensure_emu(h)
    if y is None:
        y = Inches(1.2)
    if h is None:
        h = Inches(1.8)
    n = len(items)
    gap = 0.15
    total = float(CW / 914400)
    card_w = (total - gap * (n - 1)) / n
    palette = [C["primary"], C["secondary"], C["teal"], C["accent"]]
    h_in = float(h / 914400)
    for i, item in enumerate(items):
        x = ML + Inches((card_w + gap) * i)
        clr = palette[i % len(palette)]
        # 카드 배경 (라운드 + 그림자)
        card_sh = RBOX(s, x, y, Inches(card_w), h + Pt(4),
                       C["card_bg"], radius=0.06)
        card_sh.line.color.rgb = C["card_border"]
        card_sh.line.width = Pt(0.5)
        if shadow:
            add_shadow(card_sh, preset="card")
        # 상단 컬러 바 (카드 상단에 겹침)
        R(s, x, y, Inches(card_w), Pt(4), f=clr)
        # 값 (비율 기반 배치)
        val_y = 0.18
        val_h = h_in * 0.35
        T(s, x + Inches(0.15), y + Inches(val_y), Inches(card_w - 0.3), Inches(val_h),
          item.get("value", ""), sz=28, c=clr, b=True, al=PP_ALIGN.CENTER,
          fn=FONT_W["bold"])
        # 라벨
        lbl_y = val_y + val_h
        lbl_h = h_in * 0.2
        T(s, x + Inches(0.15), y + Inches(lbl_y), Inches(card_w - 0.3), Inches(lbl_h),
          item.get("label", ""), sz=SZ["body_sm"], c=C["dark"], b=True,
          al=PP_ALIGN.CENTER, fn=FONT_W["semibold"])
        # 산출근거
        basis = item.get("basis", "")
        if basis:
            basis_y = lbl_y + lbl_h + 0.05
            basis_h = max(h_in - basis_y - 0.05, 0.3)
            T(s, x + Inches(0.1), y + Inches(basis_y), Inches(card_w - 0.2), Inches(basis_h),
              basis, sz=SZ["source"], c=C["gray"], al=PP_ALIGN.CENTER, ls=1.3)


def COMPARE(s, left_title, left_items, right_title, right_items,
            y=None, left_color=None, right_color=None, h=None):
    """좌우 비교 레이아웃 (AS-IS / TO-BE)

    Args:
        left_title, right_title: 좌/우 제목
        left_items, right_items: 좌/우 항목 리스트
        y: 시작 Y — raw float 또는 Inches 객체
        h: 본문 영역 높이 (기본 3.0") — 전체 높이는 h + 0.5(헤더)

    Returns:
        float: 실제 전체 높이(인치) = h + 0.5.
    """
    y = _ensure_emu(y)
    if y is None:
        y = Inches(1.2)
    if left_color is None:
        left_color = C["gray"]
    if right_color is None:
        right_color = C["primary"]
    body_h = h if h is not None else 3.0
    if isinstance(body_h, (int, float)):
        body_h_in = body_h
    else:
        body_h_in = float(body_h / 914400)
    total = float(CW / 914400)
    gap = 0.3                              # 좌우 간격 (화살표 제거)
    half = (total - gap) / 2
    # 좌측
    BOX(s, ML, y, Inches(half), Inches(0.5), left_color,
        left_title, sz=SZ["body"], tc=C["white"], b=True)
    R(s, ML, y + Inches(0.5), Inches(half), Inches(body_h_in),
      f=C["light"], lc=C["lgray"])
    MT(s, ML + Inches(0.15), y + Inches(0.6),
       Inches(half - 0.3), Inches(body_h_in - 0.2),
       left_items, sz=SZ["body_sm"], bul=True)
    # 중앙 구분선 (화살표 대체)
    div_x = ML + Inches(half + gap / 2 - 0.01)
    R(s, div_x, y + Inches(0.1), Pt(2), Inches(body_h_in + 0.3), f=C["lgray"])
    # 우측
    rx = ML + Inches(half + gap)
    BOX(s, rx, y, Inches(half), Inches(0.5), right_color,
        right_title, sz=SZ["body"], tc=C["white"], b=True)
    R(s, rx, y + Inches(0.5), Inches(half), Inches(body_h_in),
      f=C["light"], lc=C["lgray"])
    MT(s, rx + Inches(0.15), y + Inches(0.6),
       Inches(half - 0.3), Inches(body_h_in - 0.2),
       right_items, sz=SZ["body_sm"], bul=True)
    return body_h_in + 0.5


def TIMELINE(s, items, y=None, h=None):
    """타임라인 (가로 배치)

    Args:
        items: [("기간", "내용"), ...] 3~6개
        y: 시작 Y — raw float 또는 Inches 객체
        h: 콘텐츠 텍스트 높이 — raw float 또는 Inches 객체

    Returns:
        float: 실제 전체 높이(인치) = 0.7 + h.
        ⚠ 전체 높이는 h가 아닌 h + 0.7 (기간/마커 영역).
        VStack 예약 시 반환값을 사용할 것.
    """
    y = _ensure_emu(y)
    h = _ensure_emu(h)
    if y is None:
        y = Inches(1.2)
    if h is None:
        h = Inches(1.0)
    n = len(items)
    total = float(CW / 914400)
    cell_w = total / n
    palette = [C["primary"], C["secondary"], C["teal"], C["accent"], C["green"], C["gold"]]
    # 가로 바
    R(s, ML, y + Inches(0.5), CW, Pt(3), f=C["lgray"])
    for i, item in enumerate(items):
        # dict/tuple 모두 지원
        if isinstance(item, dict):
            period = item.get("label", item.get("period", ""))
            content = item.get("desc", item.get("content", ""))
        else:
            period, content = item[0], item[1]
        x = ML + Inches(cell_w * i)
        clr = palette[i % len(palette)]
        # 마커
        BOX(s, x + Inches(cell_w / 2 - 0.12), y + Inches(0.38),
            Inches(0.24), Inches(0.24), clr, "", sz=8)
        # 기간
        T(s, x, y, Inches(cell_w), Inches(0.35),
          period, sz=SZ["body_sm"], c=clr, b=True, al=PP_ALIGN.CENTER)
        # 내용
        T(s, x + Inches(0.05), y + Inches(0.7), Inches(cell_w - 0.1), h,
          content, sz=SZ["caption"], c=C["dark"], al=PP_ALIGN.CENTER, ls=1.5)
    return 0.7 + float(h / 914400)


# ═══════════════════════════════════════════════════════════════
#  7. 슬라이드 템플릿
# ═══════════════════════════════════════════════════════════════

def slide_cover(prs, project_name, client_name, year="2026",
                tagline="", company_name="[수행사명]"):
    """표지 슬라이드 (v3.6 — 그라디언트 배경 + 시각 폴리시)"""
    s = new_slide(prs)
    # 그라디언트 배경 (좌하→우상 대각선 다크→블루)
    c1, c2 = GRAD["cover"]()
    gradient_bg(s, c1, c2, angle=225.0)
    # 상단 악센트 라인
    R(s, Inches(0), Inches(0), SW, Pt(4), f=C["secondary"])
    # 좌측 세로 악센트 (미세한 시각 앵커)
    R(s, Inches(0.35), Inches(1.5), Pt(3), Inches(4.0), f=C["secondary"])
    # 프로젝트명
    title_tb = T(s, ML, Inches(1.8), CW, Inches(2.2),
                 project_name, sz=SZ["hero"], c=C["white"], b=True,
                 al=PP_ALIGN.LEFT, fn=FONT_W["bold"])
    set_char_spacing(title_tb, 100)
    # 구분선
    R(s, ML, Inches(4.2), Inches(3), Pt(2), f=C["secondary"])
    # 부제
    if tagline:
        T(s, ML, Inches(4.5), CW, Inches(0.5),
          tagline, sz=SZ["subtitle"], c=C["lgray"],
          fn=FONT_W["light"])
    # 연도 + 발주처
    T(s, ML, Inches(5.2), CW, Inches(0.4),
      f"{year}  |  {client_name}", sz=SZ["body"], c=C["gray"],
      fn=FONT_W["medium"])
    # 수행사명
    T(s, ML, SH - Inches(0.8), CW, Inches(0.4),
      company_name, sz=SZ["body_sm"], c=C["lgray"],
      fn=FONT_W["light"])
    # 하단 악센트 라인
    R(s, Inches(0), SH - Pt(3), SW, Pt(3), f=C["secondary"])
    return s


def slide_section_divider(prs, num, title, subtitle="", story="",
                          win_theme_key=None, win_themes=None):
    """섹션 구분자 슬라이드 (v3.6 — 그라디언트 + 대형 숫자 아웃라인)"""
    s = new_slide(prs)
    c1, c2 = GRAD["section"]()
    gradient_bg(s, c1, c2, angle=270.0)
    # 상단 악센트 라인
    R(s, Inches(0), Inches(0), SW, Pt(3), f=C["secondary"])
    # 대형 번호 (아웃라인 스타일 — 연한 색)
    T(s, ML, Inches(1.2), Inches(3), Inches(2.2),
      num, sz=110, c=darken(C["secondary"], 0.3), b=True,
      fn=FONT_W["black"])
    # 제목
    title_tb = T(s, ML, Inches(3.5), CW, Inches(0.7),
                 title, sz=SZ["divider"], c=C["white"], b=True,
                 fn=FONT_W["bold"])
    set_char_spacing(title_tb, 80)
    # 부제
    if subtitle:
        T(s, ML, Inches(4.2), CW, Inches(0.4),
          subtitle, sz=SZ["subtitle"], c=C["lgray"],
          fn=FONT_W["light"])
    # 스토리
    if story:
        T(s, ML, Inches(4.8), CW, Inches(0.4),
          story, sz=SZ["body"], c=C["secondary"])
    # 하단 라인
    R(s, ML, Inches(5.4), Inches(2), Pt(2), f=C["secondary"])
    # Win Theme
    if win_theme_key and win_themes:
        WB(s, win_theme_key, win_themes, ML, Inches(6.0))
    return s


def slide_toc(prs, title, items, pg=None):
    """목차 슬라이드

    Args:
        items: [("01", "HOOK", "설명"), ...]
    """
    s = new_slide(prs)
    TB(s, title, pg)
    y_start = 1.1
    row_h = min(0.55, 5.5 / max(len(items), 1))
    text_h = min(0.3, row_h - 0.08)
    text_pad = (row_h - 0.04 - text_h) / 2
    for i, (num, name, desc) in enumerate(items):
        y = Inches(y_start + row_h * i)
        bgc = C["light"] if i % 2 == 0 else C["white"]
        R(s, ML, y, CW, Inches(row_h - 0.04), f=bgc)
        T(s, ML + Inches(0.2), y + Inches(text_pad), Inches(0.5), Inches(text_h),
          num, sz=SZ["body"], c=C["secondary"], b=True)
        T(s, ML + Inches(0.9), y + Inches(text_pad), Inches(3), Inches(text_h),
          name, sz=SZ["body"], c=C["primary"], b=True)
        T(s, ML + Inches(4.2), y + Inches(text_pad), Inches(7), Inches(text_h),
          desc, sz=SZ["body_sm"], c=C["gray"])
    return s


def slide_exec_summary(prs, title, one_liner, win_themes_dict, kpis, why_us_points):
    """Executive Summary (v3.6 — 그라디언트 하이라이트 + 카드 그림자)"""
    s = new_slide(prs)
    TB(s, title)
    # One Sentence Pitch (그라디언트)
    HIGHLIGHT(s, one_liner, y=Inches(1.1), grad=True)
    # Win Theme (라운드 박스)
    themes = list(win_themes_dict.items())
    colors = [C["primary"], C["secondary"], C["teal"]]
    for i, (key, desc) in enumerate(themes[:3]):
        x = ML + Inches(3.95 * i)
        sh = RBOX(s, x, Inches(2.1), Inches(3.75), Inches(0.55), colors[i % 3],
                  f"Win Theme {i+1}: {desc}", sz=SZ["body_sm"], tc=C["white"],
                  b=True, radius=0.06)
        add_shadow(sh, preset="subtle")
    # KPI
    T(s, ML, Inches(2.85), Inches(3), Inches(0.25),
      "핵심 KPI", sz=SZ["body_sm"], c=C["primary"], b=True,
      fn=FONT_W["semibold"])
    KPIS(s, kpis, y=Inches(3.1), h=Inches(1.5))
    # Why Us
    T(s, ML, Inches(4.8), Inches(3), Inches(0.25),
      "Why Us", sz=SZ["body_sm"], c=C["primary"], b=True,
      fn=FONT_W["semibold"])
    R(s, ML, Inches(5.05), CW, Inches(1.2), f=C["card_bg"], lc=C["card_border"])
    for i, pt in enumerate(why_us_points[:3]):
        x = ML + Inches(3.95 * i) + Inches(0.15)
        T(s, x, Inches(5.15), Inches(3.6), Inches(0.9),
          f"— {pt}", sz=SZ["body_sm"], c=C["dark"], ls=1.5)
    return s


def slide_next_step(prs, headline, steps, contact=""):
    """Next Step / CTA 슬라이드 (v3.6 — 그라디언트 배경)"""
    s = new_slide(prs)
    c1, c2 = GRAD["section"]()
    gradient_bg(s, c1, c2, angle=270.0)
    R(s, Inches(0), Inches(0), SW, Pt(4), f=C["secondary"])
    # NEXT STEP
    ns_tb = T(s, ML, Inches(0.8), CW, Inches(0.5),
              "NEXT STEP", sz=32, c=C["white"], b=True)
    set_char_spacing(ns_tb, 200)
    R(s, ML, Inches(1.3), Inches(1.5), Pt(2), f=C["secondary"])
    # 헤드라인
    T(s, ML, Inches(1.6), CW, Inches(0.4),
      headline, sz=SZ["subtitle"], c=C["lgray"])
    # 스텝 카드
    n = len(steps)
    gap = 0.3
    total = float(CW / 914400)
    card_w = (total - gap * (n - 1)) / n
    card_y = 2.3
    card_h = 3.2
    for i, (step_label, title, desc, clr) in enumerate(steps):
        x = ML + Inches((card_w + gap) * i)
        # 라운드 카드 + 그림자 (v3.6)
        card_sh = RBOX(s, x, Inches(card_y), Inches(card_w), Inches(card_h),
                       clr, radius=0.06)
        add_shadow(card_sh, preset="elevated")
        # 스텝 라벨만 별도 (상단)
        T(s, x, Inches(card_y + 0.2), Inches(card_w), Inches(0.3),
          step_label, sz=SZ["caption"], c=C["white"], al=PP_ALIGN.CENTER,
          fn=FONT_W["medium"])
        T(s, x, Inches(card_y + 0.65), Inches(card_w), Inches(0.5),
          title, sz=20, c=C["white"], b=True, al=PP_ALIGN.CENTER,
          fn=FONT_W["bold"])
        T(s, x + Inches(0.2), Inches(card_y + 1.3), Inches(card_w - 0.4), Inches(1.2),
          desc, sz=SZ["body_sm"], c=C["white"], al=PP_ALIGN.CENTER, ls=1.5)
        if i < n - 1:
            T(s, x + Inches(card_w), Inches(card_y + 1.0), Inches(gap), Inches(0.5),
              "→", sz=22, c=C["lgray"], b=True, al=PP_ALIGN.CENTER)
    # 연락처
    if contact:
        T(s, ML, SH - Inches(0.8), CW, Inches(0.5),
          f"Contact: {contact}", sz=SZ["body_sm"], c=C["lgray"], al=PP_ALIGN.CENTER)
    return s


def slide_closing(prs, message="감사합니다", tagline="",
                  project_title="", contact=""):
    """마지막 감사 슬라이드 (v3.8 — 동적 Y 포지셔닝 + 겹침 방지)

    v3.8 변경: SZ["hero"]=60pt 기준 메시지 줄 수에 따라 아래 요소 자동 배치.
    - 1줄 → msg_h=1.0"  / 2줄 → msg_h=1.6"  / 3줄+ → msg_h=2.2"
    - 구분선/태그라인/프로젝트 제목이 메시지 하단에 동적으로 배치되어 겹침 방지
    """
    s = new_slide(prs)
    c1, c2 = GRAD["closing"]()
    gradient_bg(s, c1, c2, angle=225.0)
    R(s, Inches(0), Inches(0), SW, Pt(4), f=C["secondary"])

    # ★ 동적 높이 계산: hero(60pt) ≈ 0.83" per line + 자간(100) 보정
    # CW=11.8" 기준, 60pt 한글 ~0.83"/자 → 한 줄에 ~14자
    cw_in = float(CW / 914400)
    chars_per_line = max(1, int(cw_in / 0.83))
    line_count = max(1, -(-len(message) // chars_per_line))   # ceil division
    # 수동 줄바꿈 카운트 반영
    newline_count = message.count("\n")
    if newline_count > 0:
        line_count = max(line_count, newline_count + 1)
    msg_h = {1: 1.0, 2: 1.6}.get(line_count, 2.2)

    msg_y = Inches(2.5)
    # 좌측 세로 악센트 (메시지 영역에 맞춰 높이 조정)
    R(s, Inches(0.35), Inches(2.2), Pt(3), Inches(msg_h + 0.6), f=C["secondary"])

    # 메인 메시지 (동적 높이)
    msg_tb = T(s, ML, msg_y, CW, Inches(msg_h),
               message, sz=SZ["hero"], c=C["white"], b=True, al=PP_ALIGN.LEFT,
               fn=FONT_W["bold"])
    set_char_spacing(msg_tb, 100)

    # ★ 동적 Y: 메시지 하단 + 0.3" 여백 기준
    cursor_y = 2.5 + msg_h + 0.3

    # 구분선
    R(s, ML, Inches(cursor_y), Inches(2), Pt(2), f=C["secondary"])
    cursor_y += 0.3

    # 태그라인
    if tagline:
        T(s, ML, Inches(cursor_y), CW, Inches(0.4),
          tagline, sz=SZ["subtitle"], c=C["lgray"],
          fn=FONT_W["light"])
        cursor_y += 0.5

    # 프로젝트
    if project_title:
        T(s, ML, Inches(cursor_y), CW, Inches(0.4),
          project_title, sz=SZ["body"], c=C["gray"],
          fn=FONT_W["medium"])

    # 연락처 (항상 하단 고정)
    if contact:
        T(s, ML, SH - Inches(0.8), CW, Inches(0.4),
          contact, sz=SZ["body_sm"], c=C["lgray"])
    # 하단 라인
    R(s, Inches(0), SH - Pt(3), SW, Pt(3), f=C["secondary"])
    return s


# ═══════════════════════════════════════════════════════════════
#  8. 유틸리티
# ═══════════════════════════════════════════════════════════════

def save_pptx(prs, output_path):
    """프레젠테이션 저장"""
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    n = len(prs.slides)
    print(f"생성 완료: {output_path} ({n}장)")
    return output_path


# ═══════════════════════════════════════════════════════════════
#  9. 레이아웃 Zone 시스템 — 겹침 방지 표준 영역
# ═══════════════════════════════════════════════════════════════

# 표준 영역 (TB 타이틀바 포함 기준)
Z = {
    "tb_y":  0,        # 타이틀바 시작
    "tb_h":  0.88,     # 타이틀바 높이 (TB 함수 기준)
    "ct_y":  1.1,      # 콘텐츠 시작
    "ct_h":  5.4,      # 콘텐츠 높이
    "ct_b":  6.5,      # 콘텐츠 하단
    "ft_y":  6.7,      # 푸터 시작
    "ft_h":  0.8,      # 푸터 높이
}

# 안전 간격 (인치 float — Zone 시스템 호환)
# ★ 주의: EMU int와 연산 시 float 결과 발생 → _ensure_emu가 임계값으로 자동 판별
GAP = 0.2    # 요소 간 수직 간격 (인치)
CGAP = 0.15  # 컬럼 간 수평 간격 (인치)
CW_IN = float(CW / 914400)   # CW in inches
ML_IN = float(ML / 914400)   # ML in inches


def _cols(n, gap=CGAP):
    """N등분 컬럼 너비 계산 (inches)"""
    return (CW_IN - gap * (n - 1)) / n


# ═══════════════════════════════════════════════════════════════
#  10. LAYOUTS 데이터 — 20가지 레이아웃 프리셋
# ═══════════════════════════════════════════════════════════════

"""
레이아웃 데이터 규격:
- 각 레이아웃은 name, desc, zones 를 가짐
- zones: [{"id": str, "x": float, "y": float, "w": float, "h": float, "role": str}]
  - x, y, w, h 는 모두 인치(inches) 단위
  - role: "header"|"body"|"image"|"table"|"card"|"kpi"|"footer"
- 모든 위치는 TB() 이후 콘텐츠 영역 기준으로 사전 계산됨
"""

LAYOUTS = {

    # ── 1. 풀바디 ────────────────────────────────────────────────
    "FULL_BODY": {
        "desc": "타이틀 + 전체 너비 본문 텍스트",
        "zones": [
            {"id": "body", "x": ML_IN, "y": Z["ct_y"], "w": CW_IN, "h": Z["ct_h"], "role": "body"},
        ],
    },

    # ── 2. 하이라이트 + 본문 ──────────────────────────────────────
    "HIGHLIGHT_BODY": {
        "desc": "강조 메시지 + 본문",
        "zones": [
            {"id": "highlight", "x": ML_IN, "y": Z["ct_y"], "w": CW_IN, "h": 0.8, "role": "header"},
            {"id": "body", "x": ML_IN, "y": Z["ct_y"] + 1.0, "w": CW_IN, "h": Z["ct_h"] - 1.0, "role": "body"},
        ],
    },

    # ── 3. 2단 컬럼 ────────────────────────────────────────────────
    "TWO_COL": {
        "desc": "타이틀 + 좌우 2단 레이아웃",
        "zones": [
            {"id": "left",  "x": ML_IN, "y": Z["ct_y"],
             "w": _cols(2), "h": Z["ct_h"], "role": "body"},
            {"id": "right", "x": ML_IN + _cols(2) + CGAP, "y": Z["ct_y"],
             "w": _cols(2), "h": Z["ct_h"], "role": "body"},
        ],
    },

    # ── 4. 3단 컬럼 ────────────────────────────────────────────────
    "THREE_COL": {
        "desc": "타이틀 + 3단 비교 레이아웃",
        "zones": [
            {"id": "col1", "x": ML_IN, "y": Z["ct_y"],
             "w": _cols(3), "h": Z["ct_h"], "role": "body"},
            {"id": "col2", "x": ML_IN + (_cols(3) + CGAP), "y": Z["ct_y"],
             "w": _cols(3), "h": Z["ct_h"], "role": "body"},
            {"id": "col3", "x": ML_IN + (_cols(3) + CGAP) * 2, "y": Z["ct_y"],
             "w": _cols(3), "h": Z["ct_h"], "role": "body"},
        ],
    },

    # ── 5. 4단 컬럼 ────────────────────────────────────────────────
    "FOUR_COL": {
        "desc": "타이틀 + 4단 카드 레이아웃",
        "zones": [
            {"id": f"col{i+1}",
             "x": ML_IN + (_cols(4) + CGAP) * i, "y": Z["ct_y"],
             "w": _cols(4), "h": Z["ct_h"], "role": "card"}
            for i in range(4)
        ],
    },

    # ── 6. 좌우 비교 (AS-IS / TO-BE) ─────────────────────────────
    "COMPARE_LR": {
        "desc": "좌우 비교 (Before/After)",
        "zones": [
            {"id": "left_header",  "x": ML_IN, "y": Z["ct_y"],
             "w": (CW_IN - 0.6) / 2, "h": 0.5, "role": "header"},
            {"id": "left_body",    "x": ML_IN, "y": Z["ct_y"] + 0.5,
             "w": (CW_IN - 0.6) / 2, "h": Z["ct_h"] - 0.5, "role": "body"},
            {"id": "arrow",        "x": ML_IN + (CW_IN - 0.6) / 2, "y": Z["ct_y"],
             "w": 0.6, "h": Z["ct_h"], "role": "body"},
            {"id": "right_header", "x": ML_IN + (CW_IN - 0.6) / 2 + 0.6, "y": Z["ct_y"],
             "w": (CW_IN - 0.6) / 2, "h": 0.5, "role": "header"},
            {"id": "right_body",   "x": ML_IN + (CW_IN - 0.6) / 2 + 0.6, "y": Z["ct_y"] + 0.5,
             "w": (CW_IN - 0.6) / 2, "h": Z["ct_h"] - 0.5, "role": "body"},
        ],
    },

    # ── 7. 하이라이트 + 3단 카드 ─────────────────────────────────
    "HIGHLIGHT_THREE_CARD": {
        "desc": "강조 메시지 + 3단 카드",
        "zones": [
            {"id": "highlight", "x": ML_IN, "y": Z["ct_y"], "w": CW_IN, "h": 0.8, "role": "header"},
            {"id": "card1", "x": ML_IN, "y": Z["ct_y"] + 1.1,
             "w": _cols(3), "h": Z["ct_h"] - 1.1, "role": "card"},
            {"id": "card2", "x": ML_IN + (_cols(3) + CGAP), "y": Z["ct_y"] + 1.1,
             "w": _cols(3), "h": Z["ct_h"] - 1.1, "role": "card"},
            {"id": "card3", "x": ML_IN + (_cols(3) + CGAP) * 2, "y": Z["ct_y"] + 1.1,
             "w": _cols(3), "h": Z["ct_h"] - 1.1, "role": "card"},
        ],
    },

    # ── 8. KPI 카드 그리드 ──────────────────────────────────────
    "KPI_GRID": {
        "desc": "타이틀 + KPI 카드 + 산출근거",
        "zones": [
            {"id": "kpi_row", "x": ML_IN, "y": Z["ct_y"],
             "w": CW_IN, "h": 2.0, "role": "kpi"},
            {"id": "detail", "x": ML_IN, "y": Z["ct_y"] + 2.3,
             "w": CW_IN, "h": Z["ct_h"] - 2.3, "role": "body"},
        ],
    },

    # ── 9. 프로세스 플로우 + 설명 ────────────────────────────────
    "PROCESS_DESC": {
        "desc": "프로세스 플로우 + 하단 상세 설명",
        "zones": [
            {"id": "flow", "x": ML_IN, "y": Z["ct_y"],
             "w": CW_IN, "h": 1.2, "role": "header"},
            {"id": "flow_desc", "x": ML_IN, "y": Z["ct_y"] + 1.4,
             "w": CW_IN, "h": 0.8, "role": "body"},
            {"id": "detail", "x": ML_IN, "y": Z["ct_y"] + 2.4,
             "w": CW_IN, "h": Z["ct_h"] - 2.4, "role": "body"},
        ],
    },

    # ── 10. 타임라인 + 하단 설명 ─────────────────────────────────
    "TIMELINE_DESC": {
        "desc": "타임라인 + 하단 본문",
        "zones": [
            {"id": "timeline", "x": ML_IN, "y": Z["ct_y"],
             "w": CW_IN, "h": 2.2, "role": "header"},
            {"id": "body", "x": ML_IN, "y": Z["ct_y"] + 2.5,
             "w": CW_IN, "h": Z["ct_h"] - 2.5, "role": "body"},
        ],
    },

    # ── 11. 피라미드 + 우측 설명 ─────────────────────────────────
    "PYRAMID_DESC": {
        "desc": "좌측 피라미드 + 우측 설명",
        "zones": [
            {"id": "pyramid", "x": ML_IN, "y": Z["ct_y"],
             "w": CW_IN * 0.45, "h": Z["ct_h"], "role": "body"},
            {"id": "desc", "x": ML_IN + CW_IN * 0.5, "y": Z["ct_y"],
             "w": CW_IN * 0.5, "h": Z["ct_h"], "role": "body"},
        ],
    },

    # ── 12. 2×2 매트릭스 + 하단 설명 ─────────────────────────────
    "MATRIX_DESC": {
        "desc": "매트릭스 + 하단 시사점",
        "zones": [
            {"id": "matrix", "x": ML_IN, "y": Z["ct_y"],
             "w": CW_IN, "h": 4.0, "role": "body"},
            {"id": "insight", "x": ML_IN, "y": Z["ct_y"] + 4.2,
             "w": CW_IN, "h": Z["ct_h"] - 4.2, "role": "body"},
        ],
    },

    # ── 13. 이미지 갤러리 (3×2 그리드) ─────────────────────────────
    "GALLERY_3x2": {
        "desc": "3열 2행 이미지 갤러리 + 캡션",
        "zones": [
            {"id": f"img_{r}_{c}",
             "x": ML_IN + (_cols(3) + CGAP) * c,
             "y": Z["ct_y"] + (2.5 + GAP) * r,
             "w": _cols(3), "h": 2.3, "role": "image"}
            for r in range(2) for c in range(3)
        ],
    },

    # ── 14. 키비주얼 (좌측 이미지 + 우측 텍스트) ──────────────────
    "KEY_VISUAL": {
        "desc": "좌측 대형 이미지 + 우측 텍스트",
        "zones": [
            {"id": "image", "x": ML_IN, "y": Z["ct_y"],
             "w": CW_IN * 0.45, "h": Z["ct_h"], "role": "image"},
            {"id": "title", "x": ML_IN + CW_IN * 0.5, "y": Z["ct_y"],
             "w": CW_IN * 0.5, "h": 0.6, "role": "header"},
            {"id": "body",  "x": ML_IN + CW_IN * 0.5, "y": Z["ct_y"] + 0.8,
             "w": CW_IN * 0.5, "h": Z["ct_h"] - 0.8, "role": "body"},
        ],
    },

    # ── 15. 테이블 + 인사이트 ──────────────────────────────────────
    "TABLE_INSIGHT": {
        "desc": "데이터 테이블 + 하단 인사이트 박스",
        "zones": [
            {"id": "table", "x": ML_IN, "y": Z["ct_y"],
             "w": CW_IN, "h": 3.5, "role": "table"},
            {"id": "insight", "x": ML_IN, "y": Z["ct_y"] + 3.8,
             "w": CW_IN, "h": Z["ct_h"] - 3.8, "role": "body"},
        ],
    },

    # ── 16. 프로그램 카드 (이미지+내용+포인트) ─────────────────────
    "PROGRAM_CARD_3": {
        "desc": "3단 프로그램 카드 (이미지 + 본문 + 포인트)",
        "zones": [
            {"id": f"card{i+1}_img",
             "x": ML_IN + (_cols(3) + CGAP) * i, "y": Z["ct_y"],
             "w": _cols(3), "h": 2.0, "role": "image"}
            for i in range(3)
        ] + [
            {"id": f"card{i+1}_body",
             "x": ML_IN + (_cols(3) + CGAP) * i, "y": Z["ct_y"] + 2.1,
             "w": _cols(3), "h": 2.3, "role": "body"}
            for i in range(3)
        ] + [
            {"id": f"card{i+1}_point",
             "x": ML_IN + (_cols(3) + CGAP) * i, "y": Z["ct_y"] + 4.5,
             "w": _cols(3), "h": 0.9, "role": "footer"}
            for i in range(3)
        ],
    },

    # ── 17. 4분할 공간 ──────────────────────────────────────────
    "QUAD_GRID": {
        "desc": "2×2 이미지 그리드 + 각 캡션",
        "zones": [
            {"id": f"quad_{r}_{c}",
             "x": ML_IN + (_cols(2) + CGAP) * c,
             "y": Z["ct_y"] + (2.6 + GAP) * r,
             "w": _cols(2), "h": 2.4, "role": "image"}
            for r in range(2) for c in range(2)
        ],
    },

    # ── 18. 조직도 (3단 계층) ──────────────────────────────────────
    "ORG_CHART": {
        "desc": "PM + 감독 + 팀 3단 계층",
        "zones": [
            {"id": "pm", "x": ML_IN + CW_IN * 0.35, "y": Z["ct_y"],
             "w": CW_IN * 0.3, "h": 1.2, "role": "card"},
            {"id": "dir1", "x": ML_IN, "y": Z["ct_y"] + 1.6,
             "w": _cols(4), "h": 1.2, "role": "card"},
            {"id": "dir2", "x": ML_IN + (_cols(4) + CGAP), "y": Z["ct_y"] + 1.6,
             "w": _cols(4), "h": 1.2, "role": "card"},
            {"id": "dir3", "x": ML_IN + (_cols(4) + CGAP) * 2, "y": Z["ct_y"] + 1.6,
             "w": _cols(4), "h": 1.2, "role": "card"},
            {"id": "dir4", "x": ML_IN + (_cols(4) + CGAP) * 3, "y": Z["ct_y"] + 1.6,
             "w": _cols(4), "h": 1.2, "role": "card"},
            {"id": "team_row", "x": ML_IN, "y": Z["ct_y"] + 3.2,
             "w": CW_IN, "h": Z["ct_h"] - 3.2, "role": "body"},
        ],
    },

    # ── 19. 리스크 카드 (2열 × 3단 대응) ──────────────────────────
    "RISK_CARD": {
        "desc": "좌우 리스크 + 3단 대응 방안",
        "zones": [
            {"id": "risk1_title", "x": ML_IN, "y": Z["ct_y"],
             "w": _cols(2), "h": 0.5, "role": "header"},
            {"id": "risk1_body",  "x": ML_IN, "y": Z["ct_y"] + 0.6,
             "w": _cols(2), "h": 2.0, "role": "body"},
            {"id": "risk1_resp",  "x": ML_IN, "y": Z["ct_y"] + 2.8,
             "w": _cols(2), "h": Z["ct_h"] - 2.8, "role": "body"},
            {"id": "risk2_title", "x": ML_IN + _cols(2) + CGAP, "y": Z["ct_y"],
             "w": _cols(2), "h": 0.5, "role": "header"},
            {"id": "risk2_body",  "x": ML_IN + _cols(2) + CGAP, "y": Z["ct_y"] + 0.6,
             "w": _cols(2), "h": 2.0, "role": "body"},
            {"id": "risk2_resp",  "x": ML_IN + _cols(2) + CGAP, "y": Z["ct_y"] + 2.8,
             "w": _cols(2), "h": Z["ct_h"] - 2.8, "role": "body"},
        ],
    },

    # ── 20. 연간 간트 차트 (12개월) ───────────────────────────────
    "GANTT": {
        "desc": "월별 간트 차트 (좌측 카테고리 + 12개월 그리드)",
        "zones": [
            {"id": "categories", "x": ML_IN, "y": Z["ct_y"],
             "w": 2.0, "h": Z["ct_h"], "role": "body"},
            {"id": "grid", "x": ML_IN + 2.1, "y": Z["ct_y"],
             "w": CW_IN - 2.1, "h": Z["ct_h"], "role": "table"},
        ],
    },
}


def get_zones(layout_name):
    """레이아웃 프리셋의 zone 목록 반환 (인치 단위 dict 리스트)

    Usage:
        zones = get_zones("TWO_COL")
        left = zones["left"]   # {"x": ..., "y": ..., "w": ..., "h": ..., "role": ...}
        right = zones["right"]
    Returns:
        dict[str, dict] — id를 key로 하는 zone 딕셔너리
    """
    layout = LAYOUTS.get(layout_name)
    if not layout:
        raise ValueError(f"Unknown layout: {layout_name}. "
                         f"Available: {list(LAYOUTS.keys())}")
    return {z["id"]: z for z in layout["zones"]}


def zone_to_inches(z):
    """zone dict → (Inches(x), Inches(y), Inches(w), Inches(h)) 튜플"""
    return Inches(z["x"]), Inches(z["y"]), Inches(z["w"]), Inches(z["h"])


# ═══════════════════════════════════════════════════════════════
#  11. 추가 도식화 헬퍼 — 레퍼런스 분석 기반
# ═══════════════════════════════════════════════════════════════

def GRID(s, items, cols=3, y=None, h=None, gap=CGAP, shadow=True):
    """N×M 카드 그리드 (v3.6 — 라운드 헤더 + 그림자)

    Args:
        items: [{"title": "제목", "body": "본문" or ["줄1","줄2"], "color": RGBColor}, ...]
        cols: 열 수 (2~4)
        y: 시작 Y (기본 1.1") — raw float 또는 Inches 객체
        h: 카드 1개 높이 (기본 자동) — raw float 또는 Inches 객체
            ⚠ h는 카드 1개의 높이. 전체 높이 = h × rows + gap × (rows-1).
        gap: 간격
        shadow: 그림자 적용 여부
    """
    y = _ensure_emu(y)
    h = _ensure_emu(h)
    if y is None:
        y = Inches(Z["ct_y"])
    n = len(items)
    rows = (n + cols - 1) // cols
    col_w = _cols(cols, gap)
    if h is None:
        card_h = min(2.5, (Z["ct_h"] - gap * (rows - 1)) / rows)
    else:
        card_h = float(h / 914400)
    palette = [C["primary"], C["secondary"], C["teal"], C["accent"],
               C["green"], C["gold"]]
    for idx, item in enumerate(items):
        c_idx = idx % cols
        r_idx = idx // cols
        x = ML + Inches((col_w + gap) * c_idx)
        iy = y + Inches((card_h + gap) * r_idx)
        clr = item.get("color", palette[idx % len(palette)])
        title = item.get("title", "")
        body = item.get("body", "")
        # 카드 배경 (그림자)
        card_bg_sh = R(s, x, iy, Inches(col_w), Inches(card_h),
                       f=C["card_bg"], lc=C["card_border"], lw=0.5)
        if shadow:
            add_shadow(card_bg_sh, preset="card")
        # 헤더
        RBOX(s, x, iy, Inches(col_w), Inches(0.45), clr,
             title, sz=SZ["body"], tc=C["white"], b=True, radius=0.0)
        # 바디
        body_y = iy + Inches(0.45)
        body_h = Inches(card_h - 0.45)
        if isinstance(body, list):
            MT(s, x + Inches(0.1), body_y + Inches(0.05),
               Inches(col_w - 0.2), body_h - Inches(0.1),
               body, sz=SZ["body_sm"], bul=True)
        else:
            T(s, x + Inches(0.1), body_y + Inches(0.05),
              Inches(col_w - 0.2), body_h - Inches(0.1),
              body, sz=SZ["body_sm"])


def STAT_ROW(s, items, y=None, h=None, shadow=True):
    """통계/수치 강조 행 (v3.6 — 그림자 + 타이포 강화)

    Args:
        items: [{"value": "87%", "label": "달성률", "color": RGBColor}, ...]
        y: 시작 Y — raw float 또는 Inches 객체
        h: 높이 — raw float 또는 Inches 객체
        shadow: 그림자 적용 여부
    """
    y = _ensure_emu(y)
    h = _ensure_emu(h)
    if y is None:
        y = Inches(Z["ct_y"])
    if h is None:
        h = Inches(1.2)
    n = len(items)
    col_w = _cols(n)
    palette = [C["primary"], C["secondary"], C["teal"], C["accent"]]
    for i, item in enumerate(items):
        x = ML + Inches((col_w + CGAP) * i)
        clr = item.get("color", palette[i % len(palette)])
        # 카드 배경 (그림자)
        card_sh = R(s, x, y, Inches(col_w), h + Pt(4),
                    f=C["card_bg"], lc=C["card_border"], lw=0.5)
        if shadow:
            add_shadow(card_sh, preset="subtle")
        # 상단 라인
        R(s, x, y, Inches(col_w), Pt(4), f=clr)
        # 수치
        T(s, x, y + Inches(0.12), Inches(col_w), Inches(0.6),
          item.get("value", ""), sz=32, c=clr, b=True, al=PP_ALIGN.CENTER,
          fn=FONT_W["bold"])
        # 라벨
        T(s, x, y + Inches(0.72), Inches(col_w), Inches(0.35),
          item.get("label", ""), sz=SZ["body_sm"], c=C["dark"], al=PP_ALIGN.CENTER,
          fn=FONT_W["semibold"])


def GANTT_CHART(s, categories, months, data, y=None, colors=None):
    """간트 차트 (좌측 카테고리 + 월별 컬러 바)

    Args:
        categories: ["기획", "실행", "보고"]
        months: ["3월", "4월", "5월", ...]
        data: [[1,1,0,0,...], [0,1,1,1,...], ...] — 각 카테고리별 활성 월 (1/0)
        y: 시작 Y — raw float 또는 Inches 객체
        colors: 카테고리별 색상
    """
    y = _ensure_emu(y)
    if y is None:
        y = Inches(Z["ct_y"])
    if colors is None:
        palette = [C["primary"], C["secondary"], C["teal"], C["accent"],
                   C["green"], C["gold"]]
        colors = [palette[i % len(palette)] for i in range(len(categories))]
    cat_w = 2.0
    n_months = len(months)
    month_w = (CW_IN - cat_w - 0.1) / n_months
    row_h = min(0.5, (Z["ct_h"] - 0.5) / max(len(categories), 1))
    # 헤더
    BOX(s, ML, y, Inches(cat_w), Inches(0.4), C["primary"],
        "구분", sz=SZ["body_sm"], tc=C["white"], b=True)
    for j, m in enumerate(months):
        mx = ML + Inches(cat_w + 0.1 + month_w * j)
        BOX(s, mx, y, Inches(month_w - 0.02), Inches(0.4), C["primary"],
            m, sz=SZ["caption"], tc=C["white"], b=True)
    # 데이터 행
    for i, cat in enumerate(categories):
        ry = y + Inches(0.45 + row_h * i)
        bgc = C["light"] if i % 2 == 0 else C["white"]
        # 카테고리명
        R(s, ML, ry, Inches(cat_w), Inches(row_h - 0.02), f=bgc)
        T(s, ML + Inches(0.1), ry + Inches(0.02), Inches(cat_w - 0.2),
          Inches(row_h - 0.06), cat, sz=SZ["body_sm"], c=C["dark"], b=True)
        # 월별 바
        for j in range(n_months):
            mx = ML + Inches(cat_w + 0.1 + month_w * j)
            if i < len(data) and j < len(data[i]) and data[i][j]:
                R(s, mx, ry + Inches(0.06), Inches(month_w - 0.02),
                  Inches(row_h - 0.14), f=colors[i])
            else:
                R(s, mx, ry, Inches(month_w - 0.02),
                  Inches(row_h - 0.02), f=bgc)


def ORG(s, pm, directors, teams=None, y=None):
    """조직도 (PM + 감독 + 팀원)

    Args:
        pm: {"name": "PM명", "role": "프로젝트 매니저", "detail": "상세"}
        directors: [{"name": "감독1", "role": "역할"}, ...]
        teams: [{"name": "팀원1", "role": "역할"}, ...] (선택)
        y: 시작 Y — raw float 또는 Inches 객체
    """
    y = _ensure_emu(y)
    if y is None:
        y = Inches(Z["ct_y"])
    # PM 박스
    pm_w = 3.0
    pm_x = ML_IN + (CW_IN - pm_w) / 2
    pm_h = 1.2 if pm.get("detail") else 1.0
    BOX(s, Inches(pm_x), y, Inches(pm_w), Inches(pm_h), C["primary"],
        f"{pm.get('name', 'PM')}\n{pm.get('role', '')}", sz=SZ["body"], tc=C["white"], b=True)
    if pm.get("detail"):
        T(s, Inches(pm_x), y + Inches(0.88), Inches(pm_w), Inches(0.25),
          pm["detail"], sz=SZ["source"], c=C["white"], al=PP_ALIGN.CENTER)
    # 연결선
    line_y = y + Inches(pm_h)
    R(s, Inches(pm_x + pm_w / 2 - 0.01), line_y, Pt(2), Inches(0.3), f=C["lgray"])
    # Directors
    dir_y = line_y + Inches(0.3)
    n = len(directors)
    dir_w = _cols(n, CGAP)
    R(s, ML, dir_y, CW, Pt(2), f=C["lgray"])
    for i, d in enumerate(directors):
        dx = ML + Inches((dir_w + CGAP) * i)
        BOX(s, dx, dir_y + Inches(0.1), Inches(dir_w), Inches(0.9), C["secondary"],
            f"{d.get('name', '')}\n{d.get('role', '')}", sz=SZ["body_sm"], tc=C["white"], b=True)
    # Teams
    if teams:
        team_y = dir_y + Inches(1.3)
        n_t = len(teams)
        t_w = _cols(n_t, 0.1)
        for i, t in enumerate(teams):
            tx = ML + Inches((t_w + 0.1) * i)
            OBOX(s, tx, team_y, Inches(t_w), Inches(0.7),
                 f"{t.get('name', '')} — {t.get('role', '')}",
                 sz=SZ["caption"], tc=C["dark"], lc=C["lgray"])


def ICON_CARDS(s, items, y=None, h=None):
    """아이콘 + 텍스트 카드 행 (포인트 태그)

    Args:
        items: [{"icon": "★", "title": "제목", "desc": "설명"}, ...]
        y: 시작 Y — raw float 또는 Inches 객체
        h: 카드 높이 — raw float 또는 Inches 객체
    """
    y = _ensure_emu(y)
    h = _ensure_emu(h)
    if y is None:
        y = Inches(Z["ct_y"])
    if h is None:
        h = Inches(1.5)
    n = len(items)
    col_w = _cols(n)
    palette = [C["primary"], C["secondary"], C["teal"], C["accent"]]
    for i, item in enumerate(items):
        x = ML + Inches((col_w + CGAP) * i)
        clr = palette[i % len(palette)]
        # 아이콘 서클
        BOX(s, x + Inches(col_w / 2 - 0.3), y, Inches(0.6), Inches(0.6), clr,
            item.get("icon", "●"), sz=24, tc=C["white"], b=True)
        # 타이틀
        T(s, x, y + Inches(0.7), Inches(col_w), Inches(0.3),
          item.get("title", ""), sz=SZ["body"], c=C["dark"], b=True, al=PP_ALIGN.CENTER)
        # 설명
        T(s, x + Inches(0.1), y + Inches(1.0), Inches(col_w - 0.2),
          h - Inches(1.0),
          item.get("desc", ""), sz=SZ["body_sm"], c=C["gray"], al=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
#  12. 시각화 헬퍼
# ═══════════════════════════════════════════════════════════════

def IMG_PH(s, x, y, w, h, label="이미지 영역", on_dark=None):
    """이미지 플레이스홀더 — 회색 박스 + 아이콘 + 라벨.

    Args:
        on_dark: True면 다크 배경용(어두운 회색 박스), False면 라이트 배경용.
                 None이면 자동 감지(슬라이드 배경 체크).
    """
    # 자동 감지: 슬라이드/마스터 XML에서 직접 배경색 추출
    # (s.background 접근은 <p:bg> 요소를 강제 생성해서 마스터 상속을 깨뜨림 — 사용 금지)
    if on_dark is None:
        on_dark = _detect_dark_bg(s)

    if on_dark:
        # 다크 배경용 — 카드 톤 박스 + 얇은 보더
        fill_c = RGBColor(41, 46, 58)     # surface/dark
        border_c = RGBColor(58, 64, 78)   # border/dark
        icon_c = RGBColor(88, 96, 112)
        label_c = RGBColor(128, 136, 152)
    else:
        fill_c = C["light"]
        border_c = C["lgray"]
        icon_c = C["lgray"]
        label_c = C["gray"]

    R(s, x, y, w, h, f=fill_c, lc=border_c)
    h_in = float(h / 914400)
    T(s, x, y + Inches(0.05), w, Inches(h_in * 0.5),
      "[IMG]", sz=28, c=icon_c, al=PP_ALIGN.CENTER)
    T(s, x, y + Inches(h_in * 0.6), w, Inches(h_in * 0.3),
      label, sz=SZ["body_sm"], c=label_c, al=PP_ALIGN.CENTER)


def PROGRESS_BAR(s, x, y, w, label, value, max_val=100, color=None, show_pct=True):
    """프로그레스 바 — 라벨 + 바 + 수치"""
    if color is None:
        color = C["secondary"]
    bar_h = Inches(0.22)
    # 라벨
    T(s, x, y, Inches(2.5), Inches(0.25),
      label, sz=SZ["body_sm"], c=C["dark"], b=True)
    # 배경 바
    bar_x = x + Inches(2.6)
    bar_w = w - Inches(3.5)
    R(s, bar_x, y + Inches(0.02), bar_w, bar_h, f=C["light"])
    # 채움 바
    fill_w = Inches(float(bar_w / 914400) * min(value / max_val, 1.0))
    if float(fill_w / 914400) > 0.05:
        R(s, bar_x, y + Inches(0.02), fill_w, bar_h, f=color)
    # 수치
    txt = f"{value}%" if show_pct else str(value)
    T(s, x + w - Inches(0.8), y, Inches(0.8), Inches(0.25),
      txt, sz=SZ["body_sm"], c=color, b=True, al=PP_ALIGN.RIGHT)


def METRIC_CARD(s, x, y, w, h, value, label, sub="", color=None, shadow=True):
    """메트릭 카드 (v3.6 — 그림자 + 라운드)"""
    x = _ensure_emu(x)
    y = _ensure_emu(y)
    w = _ensure_emu(w)
    h = _ensure_emu(h)
    if color is None:
        color = C["primary"]
    h_in = float(h / 914400)
    # 카드 배경 (그림자)
    card_sh = R(s, x, y, w, h, f=C["card_bg"], lc=C["card_border"], lw=0.5)
    if shadow:
        add_shadow(card_sh, preset="subtle")
    # 상단 컬러 바
    R(s, x, y, w, Pt(4), f=color)
    # 비율 기반 배치 — 카드 높이에 비례
    val_sz = max(24, min(40, int(h_in * 24)))
    T(s, x, y + Inches(h_in * 0.08), w, Inches(h_in * 0.38),
      str(value), sz=val_sz, c=color, b=True, al=PP_ALIGN.CENTER)
    T(s, x, y + Inches(h_in * 0.46), w, Inches(h_in * 0.22),
      label, sz=SZ["body_sm"], c=C["dark"], b=True, al=PP_ALIGN.CENTER)
    if sub:
        T(s, x + Inches(0.08), y + Inches(h_in * 0.70), w - Inches(0.16),
          Inches(h_in * 0.28),
          sub, sz=SZ["caption"], c=C["gray"], al=PP_ALIGN.CENTER)


def STEP_ARROW(s, items, y=None, h=None):
    """화살표 스텝 다이어그램 — 숫자 원 + 제목 + 설명 (가로)"""
    y = _ensure_emu(y)
    h = _ensure_emu(h)
    if y is None:
        y = Inches(Z["ct_y"])
    if h is None:
        h = Inches(1.8)
    n = len(items)
    total = float(CW / 914400)
    arrow_w = 0.35
    item_w = (total - arrow_w * (n - 1)) / n
    palette = [C["primary"], C["secondary"], C["teal"], C["accent"], C["green"]]
    for i, (num, title, desc) in enumerate(items):
        clr = palette[i % len(palette)]
        x = ML + Inches((item_w + arrow_w) * i)
        # 원형 숫자
        circle_sz = 0.5
        cx = float(x / 914400) + item_w / 2 - circle_sz / 2
        BOX(s, Inches(cx), y, Inches(circle_sz), Inches(circle_sz), clr,
            str(num), sz=18, tc=C["white"], b=True)
        # 제목 (2줄 텍스트 허용)
        T(s, x, y + Inches(0.6), Inches(item_w), Inches(0.55),
          title, sz=SZ["body"], c=C["dark"], b=True, al=PP_ALIGN.CENTER)
        # 설명
        T(s, x + Inches(0.1), y + Inches(1.2), Inches(item_w - 0.2), h - Inches(1.25),
          desc, sz=SZ["body_sm"], c=C["gray"], al=PP_ALIGN.CENTER, ls=1.4)
        # 화살표 (▶ 삼각형 — 원형 숫자 기준 가운데 정렬)
        if i < n - 1:
            _flow_arrow(s, x + Inches(item_w), y, arrow_w, circle_sz,
                        cy_offset=0.0)


def DONUT_LABEL(s, x, y, w, value, label, color=None):
    """도넛 차트 스타일 라벨 — 원형 + 큰 숫자 + 라벨 (세로 배치)"""
    if color is None:
        color = C["primary"]
    # 원형 배경
    circle_d = min(float(w / 914400) * 0.7, 1.2)
    cx = float(x / 914400) + (float(w / 914400) - circle_d) / 2
    BOX(s, Inches(cx), y, Inches(circle_d), Inches(circle_d), color,
        str(value), sz=24, tc=C["white"], b=True)
    # 라벨
    T(s, x, y + Inches(circle_d + 0.1), w, Inches(0.3),
      label, sz=SZ["body_sm"], c=C["dark"], b=True, al=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
#  13. 유틸리티
# ═══════════════════════════════════════════════════════════════

def list_layouts():
    """사용 가능한 모든 레이아웃 목록 출력"""
    for name, layout in LAYOUTS.items():
        n_zones = len(layout["zones"])
        print(f"  {name:25s} — {layout['desc']}  ({n_zones} zones)")


# ═══════════════════════════════════════════════════════════════
#  14. v3.5 — VStack 자동 수직 스택
# ═══════════════════════════════════════════════════════════════

class VStack:
    """수직 자동 스택 — Y좌표를 자동 계산하여 겹침/공백 방지 (v3.7)

    Usage:
        v = VStack()
        HIGHLIGHT(s, "메시지", y=v.next(0.8))
        COLS(s, items, y=v.next(3.5), h=Inches(3.5))
        MT(s, ML, v.next(1.4), CW, Inches(1.4), lines, bul=True)

    v3.7 업데이트:
        - max_y: 최대 허용 Y 위치 (WB 뱃지 등 하단 고정 요소 보호)
        - next_raw()에서 y/h 값이 raw float이어도 안전 사용 가능
        - 컴포넌트 함수들이 _ensure_emu()로 자동 변환

    ⚠ VStack 예약 규칙:
        - HIGHLIGHT (sub 있음): 1.2"
        - HIGHLIGHT (sub 없음): 0.8"
        - FLOW: h + 0.9" (desc 있을 때) 또는 h (desc 없을 때)
        - TIMELINE: 0.7" + h (기간/마커 영역이 h 위에 추가됨)
        - COMPARE: body_h + 0.5" (헤더)
        - COLS: 기본 3.5" (반드시 명시적 h 지정 권장)
        - TABLE: 0.45 × (1 + 행 수)
        - GRID: 카드h × 행수 + gap × (행수-1)
    """

    # WB 뱃지 기본 위치: SH - 0.9" = 6.6"
    WB_SAFE_Y = 6.5  # WB 뱃지가 있을 때 콘텐츠 최대 Y

    def __init__(self, y_start=None, gap=GAP, max_y=None):
        self.y = y_start if y_start is not None else Z["ct_y"]
        self.gap = gap
        self._max_y = max_y  # None이면 무제한

    def next(self, height):
        """다음 요소의 Y위치(Inches) 반환 후 커서를 height+gap만큼 이동"""
        y = self.y
        self.y += height + self.gap
        return Inches(y)

    def next_raw(self, height, gap_override=None):
        """다음 요소의 Y위치(float, inches) 반환 — Inches 래핑 없이

        Args:
            height: 요소 높이 (inches)
            gap_override: 이 요소 뒤 간격 오버라이드 (None이면 기본 gap)

        ⚠ 반환값은 raw float. 컴포넌트 함수의 y= 파라미터에 직접 전달 가능
        (v3.7: 모든 컴포넌트가 _ensure_emu()로 자동 변환)
        """
        y = self.y
        g = gap_override if gap_override is not None else self.gap
        self.y += height + g
        return y

    def skip(self, amount=0.2):
        """추가 여백 삽입"""
        self.y += amount
        return self

    def breathe(self, amount=0.4):
        """시각적 호흡 — skip보다 큰 여백 (섹션 간 전환 등)"""
        self.y += amount
        return self

    def peek(self):
        """현재 Y위치 반환 (커서 이동 없음)"""
        return Inches(self.y)

    def peek_raw(self):
        """현재 Y위치 반환 (float, 커서 이동 없음)"""
        return self.y

    @property
    def remaining(self):
        """남은 콘텐츠 영역 높이 (inches)"""
        limit = self._max_y if self._max_y else Z["ct_b"]
        return limit - self.y

    @property
    def remaining_safe(self):
        """WB 뱃지 영역을 고려한 남은 높이 (inches)"""
        return self.WB_SAFE_Y - self.y

    @property
    def is_full(self):
        """남은 공간이 0.5" 미만이면 True"""
        return self.remaining < 0.5

    def would_overflow(self, height):
        """height 만큼 배치 시 max_y를 초과하는지 체크"""
        if self._max_y is None:
            return self.y + height > Z["ct_b"]
        return self.y + height > self._max_y


# ═══════════════════════════════════════════════════════════════
#  15. v3.5 — 라운드 코너 박스
# ═══════════════════════════════════════════════════════════════

def RBOX(s, l, t, w, h, f, text="", sz=13, tc=None, b=False, radius=0.12):
    """라운드 코너 텍스트 박스 — 카드/배지에 부드러운 인상

    Args:
        radius: 코너 반경 비율 (0.0~0.5). 0.12 = 적당히 둥글게
    """
    if tc is None:
        tc = C["white"]
    l, t, w, h = _safe_int(l), _safe_int(t), _safe_int(w), _safe_int(h)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    # 라운드 비율 설정 (0.0 = 직각, 0.5 = 완전 원형)
    sh.adjustments[0] = min(radius, 0.5)
    sh.fill.solid()
    sh.fill.fore_color.rgb = f
    sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Pt(6))
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = tc
    p.font.bold = b
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh


def ORBOX(s, l, t, w, h, text="", sz=13, tc=None, b=False, lc=None, radius=0.12):
    """라운드 아웃라인 박스 — 배경 투명, 테두리 + 라운드"""
    if tc is None:
        tc = C["dark"]
    if lc is None:
        lc = C["primary"]
    l, t, w, h = _safe_int(l), _safe_int(t), _safe_int(w), _safe_int(h)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.adjustments[0] = min(radius, 0.5)
    sh.fill.background()
    sh.line.color.rgb = lc
    sh.line.width = Pt(1.5)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Pt(6))
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = tc
    p.font.bold = b
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh


# ═══════════════════════════════════════════════════════════════
#  16. v3.5 — 미세 그림자 + 반투명 오버레이
# ═══════════════════════════════════════════════════════════════

# 그림자 프리셋 (v3.6)
SHADOW = {
    "subtle":   {"blur_pt": 2, "offset_pt": 1, "alpha": 78000},   # 은은한 떠있는 느낌
    "normal":   {"blur_pt": 3, "offset_pt": 2, "alpha": 65000},   # 기본 깊이감
    "elevated": {"blur_pt": 5, "offset_pt": 3, "alpha": 55000},   # 강한 부유감
    "card":     {"blur_pt": 4, "offset_pt": 2, "alpha": 72000},   # 카드 전용 (부드러운)
}


def add_shadow(shape, blur_pt=3, offset_pt=2, direction=2700000, alpha=60000,
               preset=None):
    """미세 그림자 — 카드/박스에 깊이감 부여

    Args:
        blur_pt: 블러 반경 (pt)
        offset_pt: 그림자 거리 (pt)
        direction: 각도 (2700000 = 우하단, 단위: 60000분의 1도)
        alpha: 불투명도 (0=불투명, 100000=완전투명)
        preset: SHADOW 프리셋 키 ("subtle"/"normal"/"elevated"/"card")
    """
    if preset and preset in SHADOW:
        p = SHADOW[preset]
        blur_pt = p["blur_pt"]
        offset_pt = p["offset_pt"]
        alpha = p["alpha"]
    try:
        from lxml import etree
    except ImportError:
        return shape
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    spPr = shape._element.spPr
    # 기존 effectLst 제거
    for old in spPr.findall(f'{{{ns}}}effectLst'):
        spPr.remove(old)
    effectLst = etree.SubElement(spPr, f'{{{ns}}}effectLst')
    outerShdw = etree.SubElement(effectLst, f'{{{ns}}}outerShdw',
                                  blurRad=str(blur_pt * 12700),
                                  dist=str(offset_pt * 12700),
                                  dir=str(direction), algn='tl')
    srgb = etree.SubElement(outerShdw, f'{{{ns}}}srgbClr', val='000000')
    etree.SubElement(srgb, f'{{{ns}}}alpha', val=str(alpha))
    return shape


def OVERLAY(s, l, t, w, h, color, alpha=50000):
    """반투명 오버레이 — 이미지 위 텍스트 가독성 확보

    Args:
        color: 오버레이 색상
        alpha: 0=완전불투명, 100000=완전투명 (50000=반투명)
    """
    sh = R(s, l, t, w, h, f=color)
    try:
        from lxml import etree
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        # shape의 spPr 내 solidFill/srgbClr에 alpha 추가
        spPr = sh._element.spPr
        srgbClr = spPr.find(f'.//{{{ns}}}srgbClr')
        if srgbClr is not None:
            for old in srgbClr.findall(f'{{{ns}}}alpha'):
                srgbClr.remove(old)
            etree.SubElement(srgbClr, f'{{{ns}}}alpha', val=str(alpha))
    except (ImportError, AttributeError):
        pass
    return sh


# ═══════════════════════════════════════════════════════════════
#  17. v3.5 — 구분선 / 악센트 요소
# ═══════════════════════════════════════════════════════════════

def DIVIDER(s, y, style="line", color=None, w=None):
    """수평 구분선

    Args:
        style: "line" (실선), "thick" (두꺼운 선), "double" (이중선)
        color: 색상 (기본 lgray)
        w: 너비 (기본 CW)
    """
    if color is None:
        color = C["lgray"]
    if w is None:
        w = CW
    if style == "line":
        R(s, ML, Inches(y), w, Pt(1), f=color)
    elif style == "thick":
        R(s, ML, Inches(y), w, Pt(3), f=color)
    elif style == "double":
        R(s, ML, Inches(y), w, Pt(1), f=color)
        R(s, ML, Inches(y + 0.06), w, Pt(1), f=color)


def ACCENT_LINE(s, x, y, h, color=None, w_pt=3):
    """좌측 악센트 라인 — 인용문/강조 블록 좌측 수직선

    Args:
        x, y: 위치 (inches)
        h: 높이 (inches)
        w_pt: 선 두께 (pt)
    """
    if color is None:
        color = C["secondary"]
    R(s, Inches(x), Inches(y), Pt(w_pt), Inches(h), f=color)


# ═══════════════════════════════════════════════════════════════
#  18. v3.5 — 인용문 / 번호 리스트
# ═══════════════════════════════════════════════════════════════

def QUOTE(s, text, author="", y=None, color=None, style="modern", text_color=None):
    """인용문 블록

    Args:
        text: 인용문 텍스트
        author: 출처/저자
        y: Y 위치 — raw float 또는 Inches 객체
        color: 악센트 라인/테두리 색상
        style: "modern" (좌측 악센트 라인) / "box" (박스형)
        text_color: 인용문 텍스트 색상 (None이면 C["dark"])
            ⚠ 다크 배경에서 사용 시 C["white"] 등으로 지정할 것.
    """
    y = _ensure_emu(y)
    if y is None:
        y = Inches(Z["ct_y"])
    if color is None:
        color = C["secondary"]
    tc = text_color if text_color is not None else C["dark"]
    if style == "modern":
        total_h = 1.0 if not author else 1.3
        # 좌측 악센트 라인
        R(s, ML, y, Pt(4), Inches(total_h), f=color)
        # 인용문 텍스트
        T(s, ML + Inches(0.3), y + Inches(0.08), CW - Inches(0.3), Inches(0.7),
          f'\u201c{text}\u201d', sz=16, c=tc, al=PP_ALIGN.LEFT, ls=1.5)
        if author:
            T(s, ML + Inches(0.3), y + Inches(0.85), CW - Inches(0.3), Inches(0.3),
              f"\u2014 {author}", sz=SZ["body_sm"], c=C["gray"])
        return total_h
    elif style == "box":
        total_h = 1.2 if not author else 1.5
        R(s, ML, y, CW, Inches(total_h), f=C["light"], lc=color, lw=1.5)
        # 큰 따옴표
        T(s, ML + Inches(0.2), y + Inches(0.02), Inches(0.5), Inches(0.5),
          "\u201c", sz=36, c=color, b=True)
        T(s, ML + Inches(0.5), y + Inches(0.15), CW - Inches(0.8), Inches(0.7),
          text, sz=14, c=tc, al=PP_ALIGN.LEFT, ls=1.5)
        if author:
            T(s, ML + Inches(0.5), y + Inches(0.9), CW - Inches(0.8), Inches(0.3),
              f"\u2014 {author}", sz=SZ["body_sm"], c=C["gray"], al=PP_ALIGN.RIGHT)
        return total_h


def NUMBERED_LIST(s, x, y, w, items, sz=13, gap=0.55):
    """번호 리스트 — 색상 원형 번호 + 제목 + 설명

    Args:
        items: [("제목", "설명"), ...] 또는 ["항목1", "항목2", ...]
        gap: 항목 간 간격
    Returns:
        float: 전체 높이 (inches)
    """
    palette = [C["primary"], C["secondary"], C["teal"], C["accent"], C["green"]]
    w_in = float(w / 914400) if hasattr(w, '__class__') and w.__class__.__name__ != 'float' else w
    total_h = 0
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            title, desc = item
        else:
            title, desc = item, ""
        iy = y + Inches(gap * i)
        clr = palette[i % len(palette)]
        # 번호 원 (라운드 박스)
        RBOX(s, x, iy, Inches(0.38), Inches(0.38), clr,
             str(i + 1), sz=12, tc=C["white"], b=True, radius=0.5)
        # 제목
        T(s, x + Inches(0.52), iy + Inches(0.02), Inches(w_in - 0.52), Inches(0.25),
          title, sz=sz, c=C["dark"], b=True)
        # 설명
        if desc:
            T(s, x + Inches(0.52), iy + Inches(0.28), Inches(w_in - 0.52), Inches(0.22),
              desc, sz=SZ["body_sm"], c=C["gray"])
        total_h = gap * i + gap
    return total_h


# ═══════════════════════════════════════════════════════════════
#  19. v3.5 — 네이티브 차트 (BAR / PIE / LINE)
# ═══════════════════════════════════════════════════════════════

def BAR_CHART(s, x, y, w, h, categories, series_data, title="",
              chart_type="column", colors=None):
    """바 차트 — 비교 데이터 시각화

    Args:
        categories: ["항목A", "항목B", "항목C"]
        series_data: [("시리즈명", [10, 20, 30]), ...]
        chart_type: "column" (세로) / "bar" (가로) / "stacked"
        colors: 시리즈별 색상 리스트
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    type_map = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "stacked": XL_CHART_TYPE.COLUMN_STACKED,
    }
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data:
        chart_data.add_series(name, values)

    graphic = s.shapes.add_chart(type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED),
                                  x, y, w, h, chart_data)
    chart = graphic.chart
    chart.has_legend = len(series_data) > 1
    if chart.has_legend:
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(SZ["caption"])
        chart.legend.font.name = FONT

    # Modern 스타일 적용
    if colors is None:
        colors = [C["primary"], C["secondary"], C["teal"], C["accent"]]
    plot = chart.plots[0]
    plot.gap_width = 80
    for i, series in enumerate(plot.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = colors[i % len(colors)]

    # 축 폰트
    if chart.category_axis:
        chart.category_axis.tick_labels.font.size = Pt(SZ["caption"])
        chart.category_axis.tick_labels.font.name = FONT
    if chart.value_axis:
        chart.value_axis.tick_labels.font.size = Pt(SZ["caption"])
        chart.value_axis.tick_labels.font.name = FONT
        chart.value_axis.has_major_gridlines = True

    return graphic


def PIE_CHART(s, x, y, w, h, categories, values, title="", colors=None, donut=False):
    """파이/도넛 차트

    Args:
        categories: ["항목A", "항목B", "항목C"]
        values: [30, 50, 20]
        donut: True면 도넛 차트
        colors: 항목별 색상 리스트
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    chart_type = XL_CHART_TYPE.DOUGHNUT if donut else XL_CHART_TYPE.PIE
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("값", values)

    graphic = s.shapes.add_chart(chart_type, x, y, w, h, chart_data)
    chart = graphic.chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(SZ["caption"])
    chart.legend.font.name = FONT

    # Modern 스타일 색상
    if colors is None:
        colors = [C["primary"], C["secondary"], C["teal"], C["accent"],
                  C["green"], C["gold"], C["gray"], C["orange"]]
    plot = chart.plots[0]
    for i, point in enumerate(plot.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = colors[i % len(colors)]

    return graphic


def LINE_CHART(s, x, y, w, h, categories, series_data, title="",
               colors=None, smooth=False):
    """라인 차트 — 추세 데이터

    Args:
        categories: ["1월", "2월", "3월", ...]
        series_data: [("시리즈명", [10, 20, 30, ...]), ...]
        smooth: True면 곡선
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    chart_type = XL_CHART_TYPE.LINE_MARKERS if not smooth else XL_CHART_TYPE.LINE
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data:
        chart_data.add_series(name, values)

    graphic = s.shapes.add_chart(chart_type, x, y, w, h, chart_data)
    chart = graphic.chart
    chart.has_legend = len(series_data) > 1
    if chart.has_legend:
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(SZ["caption"])
        chart.legend.font.name = FONT

    # Modern 스타일 적용
    if colors is None:
        colors = [C["primary"], C["secondary"], C["teal"], C["accent"]]
    plot = chart.plots[0]
    if smooth:
        plot.smooth = True
    for i, series in enumerate(plot.series):
        series.format.line.color.rgb = colors[i % len(colors)]
        series.format.line.width = Pt(2.5)
        series.smooth = smooth

    # 축 폰트
    if chart.category_axis:
        chart.category_axis.tick_labels.font.size = Pt(SZ["caption"])
        chart.category_axis.tick_labels.font.name = FONT
    if chart.value_axis:
        chart.value_axis.tick_labels.font.size = Pt(SZ["caption"])
        chart.value_axis.tick_labels.font.name = FONT
        chart.value_axis.has_major_gridlines = True

    return graphic


# ═══════════════════════════════════════════════════════════════
#  20. v3.5 — 테마 시스템
# ═══════════════════════════════════════════════════════════════

THEMES = {
    "default_blue": {
        "primary": (0, 44, 95), "secondary": (0, 170, 210),
        "teal": (0, 161, 156), "accent": (230, 51, 18),
        "dark": (33, 33, 33), "light": (245, 245, 245),
    },
    "warm": {
        "primary": (139, 69, 19), "secondary": (210, 105, 30),
        "teal": (160, 82, 45), "accent": (220, 20, 60),
        "dark": (50, 30, 15), "light": (255, 248, 240),
    },
    "forest": {
        "primary": (27, 94, 32), "secondary": (76, 175, 80),
        "teal": (0, 150, 136), "accent": (255, 111, 0),
        "dark": (30, 40, 30), "light": (241, 248, 233),
    },
    "corporate": {
        "primary": (38, 50, 56), "secondary": (255, 111, 0),
        "teal": (0, 137, 123), "accent": (213, 0, 0),
        "dark": (33, 33, 33), "light": (245, 245, 245),
    },
    "purple": {
        "primary": (74, 20, 140), "secondary": (156, 39, 176),
        "teal": (0, 150, 136), "accent": (233, 30, 99),
        "dark": (40, 20, 60), "light": (243, 229, 245),
    },
}

_original_colors = {k: v for k, v in C.items()}


def _recalculate_derived_colors():
    """파생 컬러 재계산 — apply_theme/register_theme 후 호출"""
    C["primary_dark"]    = darken(C["primary"], 0.3)
    C["primary_light"]   = lighten(C["primary"], 0.85)
    C["secondary_dark"]  = darken(C["secondary"], 0.25)
    C["secondary_light"] = lighten(C["secondary"], 0.80)
    C["teal_light"]      = lighten(C["teal"], 0.80)
    C["accent_light"]    = lighten(C["accent"], 0.80)
    if "green" in C:
        C["green_light"] = lighten(C["green"], 0.80)


def register_theme(name, colors_dict):
    """외부 테마 동적 등록 (DesignAgent용)

    Args:
        name: 테마 이름 (예: "nikke_dark", "popup_warm")
        colors_dict: {"primary": (r,g,b), "secondary": (r,g,b), ...}
                     최소 primary, secondary 필수. 나머지는 기존 값 유지.

    Returns:
        등록된 테마 이름

    Example:
        register_theme("nikke_dark", {
            "primary": (15, 15, 40),
            "secondary": (255, 180, 0),
            "accent": (230, 51, 18),
        })
        apply_theme("nikke_dark")
    """
    THEMES[name] = colors_dict
    return name


def apply_theme(theme_name):
    """컬러 팔레트 일괄 변경 + 파생 컬러 자동 재계산

    Args:
        theme_name: THEMES 딕셔너리의 키 (예: "default_blue", "warm", "forest")
    """
    theme = THEMES.get(theme_name)
    if not theme:
        raise ValueError(f"Unknown theme: {theme_name}. Available: {list(THEMES.keys())}")
    for key, rgb in theme.items():
        if key in C:
            C[key] = RGBColor(*rgb)
    _recalculate_derived_colors()
    return theme_name


def reset_theme():
    """테마를 기본값(default_blue)으로 복원"""
    for k, v in _original_colors.items():
        C[k] = v


def list_themes():
    """사용 가능한 테마 목록 출력"""
    for name in THEMES:
        print(f"  {name}")


# ═══════════════════════════════════════════════════════════════
#  21. v3.5 — 레이아웃 시퀀스 검증
# ═══════════════════════════════════════════════════════════════

def validate_slide_shapes(prs, verbose=True):
    """v3.8 — 생성된 PPTX의 도형 배치 검증 (겹침, 경계 초과, 가시성 문제 탐지)

    생성 스크립트 실행 후 save_pptx() 전에 호출하여 품질 이슈를 사전 탐지.

    Returns:
        list[dict] — 각 이슈: {"slide": int, "severity": "error"|"warn", "msg": str}
    """
    issues = []
    SH_IN = 7.5
    SW_IN = 13.333
    BOTTOM_MARGIN = 7.0   # 하단 0.5" 마진 (페이지번호 + 출처)
    MAX_EMU = 914400 * 20  # 20" — 정상 EMU 최대치

    for idx, slide in enumerate(prs.slides, 1):
        shapes = []
        for sp in slide.shapes:
            try:
                # ★ 규칙 0: 비정상 EMU 탐지 (PPT 파손 원인)
                raw_l = sp.left if sp.left is not None else 0
                raw_t = sp.top if sp.top is not None else 0
                raw_w = sp.width if sp.width is not None else 0
                raw_h = sp.height if sp.height is not None else 0
                for attr_name, raw_val in [("left", raw_l), ("top", raw_t),
                                            ("width", raw_w), ("height", raw_h)]:
                    if abs(raw_val) > MAX_EMU:
                        issues.append({
                            "slide": idx, "severity": "error",
                            "msg": f"CORRUPT EMU: '{sp.name}' {attr_name}={raw_val} "
                                   f"({float(raw_val/914400):.0f}in) > 20in max"
                        })
                l = float(raw_l / 914400)
                t = float(raw_t / 914400)
                w = float(raw_w / 914400)
                h = float(raw_h / 914400)
            except (TypeError, ZeroDivisionError):
                continue
            # 비정상 크기 필터: 슬라이드보다 100배 이상 큰 값은 무시
            if l > 100 or t > 100 or w > 100 or h > 100:
                continue
            name = sp.name or ""
            shapes.append({"l": l, "t": t, "w": w, "h": h, "name": name, "sp": sp})

        for s in shapes:
            bottom = s["t"] + s["h"]
            right = s["l"] + s["w"]

            # 규칙 1: 하단 경계 초과 (0.5" 이상 침범)
            if bottom > BOTTOM_MARGIN + 0.5 and s["h"] > 0.1:
                issues.append({
                    "slide": idx, "severity": "error",
                    "msg": f"bottom overflow: '{s['name']}' bottom={bottom:.1f}\" > {BOTTOM_MARGIN}\" (+{bottom - BOTTOM_MARGIN:.1f}\")"
                })

            # 규칙 2: 우측 경계 초과
            if right > SW_IN + 0.1 and s["w"] > 0.1:
                issues.append({
                    "slide": idx, "severity": "warn",
                    "msg": f"right overflow: '{s['name']}' right={right:.1f}\" > {SW_IN}\""
                })

        # 규칙 3: 독립 컴포넌트 간 겹침 탐지
        # 제외: (a) 배경 도형 (슬라이드 전체 크기), (b) 부모-자식 관계
        big_shapes = [s for s in shapes
                      if s["w"] > 2.0 and s["h"] > 1.0
                      and not (s["w"] > SW_IN * 0.9 and s["h"] > SH_IN * 0.9)  # 전체 배경 제외
                      and not (s["w"] > SW_IN * 0.8)]  # 전체 너비 구조 요소 제외 (COMPARE 패널 등)
        for i, a in enumerate(big_shapes):
            for b in big_shapes[i+1:]:
                # 부모-자식 판별: 한쪽이 다른 쪽을 완전히 포함하면 → 컴포넌트 내부 구조 (정상)
                a_contains_b = (a["l"] <= b["l"] and a["t"] <= b["t"] and
                                a["l"]+a["w"] >= b["l"]+b["w"] and a["t"]+a["h"] >= b["t"]+b["h"])
                b_contains_a = (b["l"] <= a["l"] and b["t"] <= a["t"] and
                                b["l"]+b["w"] >= a["l"]+a["w"] and b["t"]+b["h"] >= a["t"]+a["h"])
                if a_contains_b or b_contains_a:
                    continue  # 부모-자식 관계 → 정상
                # 겹침 영역 계산
                ox = max(0, min(a["l"]+a["w"], b["l"]+b["w"]) - max(a["l"], b["l"]))
                oy = max(0, min(a["t"]+a["h"], b["t"]+b["h"]) - max(a["t"], b["t"]))
                overlap_area = ox * oy
                min_area = min(a["w"]*a["h"], b["w"]*b["h"])
                if min_area > 0 and overlap_area / min_area > 0.3:
                    issues.append({
                        "slide": idx, "severity": "error",
                        "msg": f"overlap: '{a['name']}' <> '{b['name']}' ({overlap_area/min_area:.0%})"
                    })

    if verbose and issues:
        errors = [i for i in issues if i["severity"] == "error"]
        warns = [i for i in issues if i["severity"] == "warn"]
        print(f"\n{'='*60}")
        print(f"  [!] Layout Validation: {len(errors)} errors, {len(warns)} warnings")
        print(f"{'='*60}")
        for i in issues:
            icon = "[ERR]" if i["severity"] == "error" else "[WRN]"
            print(f"  {icon} Slide {i['slide']}: {i['msg']}")
        print()

    return issues


def validate_sequence(slide_info):
    """레이아웃 시퀀스 검증 → 경고 메시지 리스트 반환

    Args:
        slide_info: [{"layout": "THREE_COL", "has_image": False, "has_highlight": True}, ...]

    Returns:
        list[str] — 경고 메시지
    """
    warnings = []
    layouts = [s.get("layout", "") for s in slide_info]

    # 규칙 1: 같은 레이아웃 3회 연속 금지
    for i in range(len(layouts) - 2):
        if layouts[i] and layouts[i] == layouts[i + 1] == layouts[i + 2]:
            warnings.append(
                f"[시각 단조] 슬라이드 {i+1}-{i+3}: '{layouts[i]}' 3회 연속 → 다른 레이아웃 권장")

    # 규칙 2: 5장 연속 이미지 없음 경고
    no_img_streak = 0
    for i, s in enumerate(slide_info):
        if s.get("has_image", False):
            no_img_streak = 0
        else:
            no_img_streak += 1
            if no_img_streak >= 5:
                warnings.append(
                    f"[이미지 부족] 슬라이드 {i-3}~{i+1}: 5장 연속 이미지 없음 → IMG_PH 추가 권장")
                no_img_streak = 0

    # 규칙 3: 5장 연속 HIGHLIGHT 없음 경고
    no_hl_streak = 0
    for i, s in enumerate(slide_info):
        if s.get("has_highlight", False):
            no_hl_streak = 0
        else:
            no_hl_streak += 1
            if no_hl_streak >= 5:
                warnings.append(
                    f"[강조 부족] 슬라이드 {i-3}~{i+1}: 5장 연속 HIGHLIGHT 없음 → 핵심 메시지 강조 권장")
                no_hl_streak = 0

    return warnings


# ═══════════════════════════════════════════════════════════════
#  22. v3.5 — new_presentation 템플릿 지원
# ═══════════════════════════════════════════════════════════════

def new_presentation_from_template(template_path):
    """기존 PPTX 템플릿 기반 프레젠테이션 생성

    마스터 슬라이드의 로고, 푸터, 테마 색상을 그대로 활용.
    """
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")
    prs = Presentation(template_path)
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


# ═══════════════════════════════════════════════════════════════
#  23. v3.5 — 카드 컴포넌트 (그림자 + 라운드)
# ═══════════════════════════════════════════════════════════════

def CARD(s, x, y, w, h, title, body="", color=None, shadow=True, rounded=True):
    """고급 카드 — 라운드 + 미세 그림자 + 컬러 상단바

    Args:
        title: 카드 제목
        body: 본문 텍스트 또는 ["줄1", "줄2"]
        color: 상단 바 색상
        shadow: 그림자 여부
        rounded: 라운드 코너 여부
    """
    if color is None:
        color = C["primary"]
    h_in = float(h / 914400) if hasattr(h, '__class__') and h.__class__.__name__ != 'float' else h
    h = Inches(h_in) if not hasattr(h, '__class__') or h.__class__.__name__ == 'float' else h

    # 배경 카드
    if rounded:
        card_sh = RBOX(s, x, y, w, h, C["white"], radius=0.08)
        card_sh.line.color.rgb = C["lgray"]
        card_sh.line.width = Pt(0.5)
    else:
        card_sh = R(s, x, y, w, h, f=C["white"], lc=C["lgray"], lw=0.5)

    if shadow:
        add_shadow(card_sh)

    # 상단 컬러 바
    w_in = float(w / 914400) if hasattr(w, '__class__') and w.__class__.__name__ != 'float' else w
    R(s, x, y, w, Pt(5), f=color)

    # 제목
    T(s, x + Inches(0.15), y + Inches(0.15), Inches(w_in - 0.3), Inches(0.35),
      title, sz=SZ["body"], c=C["dark"], b=True)

    # 본문
    if body:
        body_y = y + Inches(0.55)
        body_h = Inches(h_in - 0.7)
        if isinstance(body, list):
            MT(s, x + Inches(0.15), body_y, Inches(w_in - 0.3), body_h,
               body, sz=SZ["body_sm"], bul=True)
        else:
            T(s, x + Inches(0.15), body_y, Inches(w_in - 0.3), body_h,
              body, sz=SZ["body_sm"], c=C["gray"], ls=1.5)

    return card_sh


# ═══════════════════════════════════════════════════════════════
#  21. v3.8 — 비주얼 컴포넌트 (이벤트/에이전시 제안서용)
# ═══════════════════════════════════════════════════════════════

def HERO_IMAGE(s, desc="이미지 영역", title="", subtitle="",
               overlay_alpha=45000, title_pos="bottom_left"):
    """전폭 이미지 플레이스홀더 + 오버레이 텍스트 (v3.8)

    전체 콘텐츠 영역을 덮는 이미지 위에 반투명 오버레이와 텍스트를 배치.
    존 컨셉 비주얼, 무드보드 히어로, 키비주얼 등에 사용.

    Args:
        desc: 이미지 설명 (프로덕션 시 실제 이미지로 교체)
        title: 오버레이 타이틀 (흰색 대형)
        subtitle: 오버레이 서브타이틀
        overlay_alpha: 오버레이 투명도 (0=불투명, 100000=투명)
        title_pos: "bottom_left", "center", "bottom_center"
    """
    # 이미지 플레이스홀더 (전폭)
    img_y = Inches(Z["ct_y"])
    img_h = Inches(Z["ct_h"])
    IMG_PH(s, ML, img_y, CW, img_h, desc)
    # 하단 그라데이션 오버레이
    ov_h = Inches(2.5)
    ov_y = Inches(Z["ct_b"] - 2.5)
    OVERLAY(s, ML, ov_y, CW, ov_h, C["dark"], alpha=overlay_alpha)
    # 타이틀 배치
    if title:
        if title_pos == "center":
            ty = Inches(Z["ct_y"] + Z["ct_h"] / 2 - 0.6)
            al = PP_ALIGN.CENTER
        elif title_pos == "bottom_center":
            ty = Inches(Z["ct_b"] - 2.0)
            al = PP_ALIGN.CENTER
        else:  # bottom_left
            ty = Inches(Z["ct_b"] - 2.0)
            al = PP_ALIGN.LEFT
        T(s, ML + Inches(0.2), ty, CW - Inches(0.4), Inches(0.8),
          title, sz=SZ["divider"], c=C["white"], b=True, al=al,
          fn=FONT_W["bold"])
    if subtitle:
        sub_y = ty + Inches(0.85) if title else Inches(Z["ct_b"] - 1.0)
        T(s, ML + Inches(0.2), sub_y, CW - Inches(0.4), Inches(0.5),
          subtitle, sz=SZ["action"], c=C["secondary_light"], al=al,
          fn=FONT_W["medium"])


def MOOD_BOARD(s, images, y=None, h=None, cols=3, captions=None):
    """다중 이미지 그리드 — 레퍼런스/무드보드 (v3.8)

    이미지 플레이스홀더를 그리드로 배치. 텍스트 카드가 아닌 이미지 중심.

    Args:
        images: 이미지 설명 리스트 ["이미지1", "이미지2", ...]
        y: 시작 Y (raw float 또는 Inches)
        h: 전체 높이 (raw float 또는 Inches, 기본 3.5")
        cols: 열 수 (2~4)
        captions: 캡션 리스트 (None이면 생략)
    """
    y = _ensure_emu(y)
    h = _ensure_emu(h)
    if y is None:
        y = Inches(Z["ct_y"])
    if h is None:
        h = Inches(3.5)
    n = len(images)
    rows = (n + cols - 1) // cols
    gap = 0.12
    col_w = (float(CW / 914400) - gap * (cols - 1)) / cols
    h_in = float(h / 914400)
    cap_h = 0.25 if captions else 0.0
    row_h = (h_in - gap * (rows - 1) - cap_h * rows) / rows

    for idx, desc in enumerate(images):
        r, c_idx = divmod(idx, cols)
        ix = ML + Inches((col_w + gap) * c_idx)
        iy = y + Inches((row_h + cap_h + gap) * r)
        IMG_PH(s, ix, iy, Inches(col_w), Inches(row_h), desc)
        if captions and idx < len(captions):
            T(s, ix, iy + Inches(row_h + 0.02), Inches(col_w), Inches(cap_h),
              captions[idx], sz=SZ["caption"], c=C["gray"], al=PP_ALIGN.CENTER)


def ZONE_MAP(s, zones, y=None, h=None, bg_color=None, title=""):
    """부스 평면도 — 존별 비율 배치 다이어그램 (v3.8)

    각 존을 비율 좌표(0.0~1.0)로 배치하여 공간 레이아웃을 시각화.
    부스 배치도, 플로어맵, 공간 구성도 등에 사용.

    Args:
        zones: [{"name": "ZONE 1", "label": "Ark Gate",
                 "x": 0.0, "y": 0.0, "w": 0.3, "h": 0.5,
                 "color": C["primary"]}, ...]
                x, y, w, h는 전체 영역 대비 비율 (0.0~1.0)
        y: 시작 Y
        h: 전체 높이
        bg_color: 배경 색상 (None이면 C["light"])
        title: 상단 타이틀 (빈 문자열이면 생략)
    """
    y_emu = _ensure_emu(y)
    h_emu = _ensure_emu(h)
    if y_emu is None:
        y_emu = Inches(Z["ct_y"])
    if h_emu is None:
        h_emu = Inches(4.5)
    if bg_color is None:
        bg_color = C["light"]

    map_w = float(CW / 914400)
    map_h = float(h_emu / 914400)
    map_x = float(ML / 914400)
    map_y = float(y_emu / 914400)

    title_offset = 0.0
    if title:
        T(s, ML, y_emu, CW, Inches(0.35),
          title, sz=SZ["body"], c=C["dark"], b=True, al=PP_ALIGN.CENTER,
          fn=FONT_W["semibold"])
        title_offset = 0.4
        map_h -= title_offset

    # 배경 프레임
    R(s, ML, Inches(map_y + title_offset), CW, Inches(map_h),
      f=bg_color, lc=C["lgray"], lw=0.5)

    palette = [C["primary"], C["secondary"], C["teal"], C["accent"],
               C["green"], C["orange"]]

    # ratio 기반 자동 레이아웃 지원: ratio만 있으면 2행 그리드로 배치
    has_explicit_pos = any("x" in z and "w" in z for z in zones)
    if not has_explicit_pos:
        n = len(zones)
        cols = min(n, 3)
        rows = (n + cols - 1) // cols
        for idx, zone in enumerate(zones):
            r = idx // cols
            c = idx % cols
            zone.setdefault("x", c / cols)
            zone.setdefault("y", r / rows)
            zone.setdefault("w", 1.0 / cols)
            zone.setdefault("h", 1.0 / rows)

    for idx, zone in enumerate(zones):
        zx = zone.get("x", 0.0)
        zy = zone.get("y", 0.0)
        zw = zone.get("w", 0.2)
        zh = zone.get("h", 0.2)
        color = zone.get("color", palette[idx % len(palette)])
        name = zone.get("name", "")
        label = zone.get("label", "")

        pad = 0.06
        rx = map_x + pad + (map_w - pad * 2) * zx
        ry = map_y + title_offset + pad + (map_h - pad * 2) * zy
        rw = (map_w - pad * 2) * zw - pad
        rh = (map_h - pad * 2) * zh - pad

        RBOX(s, Inches(rx), Inches(ry), Inches(rw), Inches(rh),
             color, "", radius=0.08)
        # 존 이름 (중앙 상단)
        if name:
            T(s, Inches(rx + 0.05), Inches(ry + 0.08), Inches(rw - 0.1), Inches(0.25),
              name, sz=SZ["caption"], c=C["white"], b=True, al=PP_ALIGN.CENTER,
              fn=FONT_W["bold"])
        # 존 라벨 (중앙)
        if label:
            T(s, Inches(rx + 0.05), Inches(ry + rh * 0.35), Inches(rw - 0.1), Inches(0.35),
              label, sz=SZ["body_sm"], c=C["white"], al=PP_ALIGN.CENTER,
              fn=FONT_W["medium"])


def DETAIL_CARD(s, x, y, w, h, title, body="", image_desc="",
                image_ratio=0.45, color=None, shadow=True):
    """이미지+텍스트 카드 — 상단 이미지, 하단 텍스트 (v3.8)

    COLS/GRID의 텍스트 전용 카드와 달리 이미지 영역이 포함된 리치 카드.

    Args:
        x, y, w, h: 카드 위치/크기 (EMU 또는 Inches)
        title: 카드 제목
        body: 본문 텍스트 (문자열 또는 리스트)
        image_desc: 이미지 설명 (빈 문자열이면 이미지 영역 없이 큰 텍스트)
        image_ratio: 이미지 영역 비율 (0.3~0.6)
        color: 제목 바 색상
        shadow: 그림자 적용 여부
    """
    x = _ensure_emu(x) if isinstance(x, float) else x
    y = _ensure_emu(y) if isinstance(y, float) else y
    if color is None:
        color = C["primary"]
    w_in = float(w / 914400) if not isinstance(w, float) else w
    h_in = float(h / 914400) if not isinstance(h, float) else h
    if isinstance(w, float):
        w = Inches(w)
    if isinstance(h, float):
        h = Inches(h)

    # 카드 배경
    card_sh = RBOX(s, x, y, w, h, C["white"], radius=0.08)
    card_sh.line.color.rgb = C["lgray"]
    card_sh.line.width = Pt(0.5)
    if shadow:
        add_shadow(card_sh)

    # 이미지 영역
    if image_desc:
        img_h = h_in * image_ratio
        IMG_PH(s, x + Inches(0.08), y + Inches(0.08),
               Inches(w_in - 0.16), Inches(img_h - 0.08), image_desc)
        text_y = y + Inches(img_h + 0.05)
        text_h = h_in - img_h - 0.15
    else:
        text_y = y + Inches(0.08)
        text_h = h_in - 0.16

    # 제목 바
    R(s, x, text_y, w, Pt(4), f=color)
    T(s, x + Inches(0.12), text_y + Inches(0.08), Inches(w_in - 0.24), Inches(0.3),
      title, sz=SZ["body"], c=C["dark"], b=True, fn=FONT_W["semibold"])

    # 본문
    if body:
        body_y = text_y + Inches(0.42)
        body_h_in = text_h - 0.5
        if isinstance(body, list):
            MT(s, x + Inches(0.12), body_y, Inches(w_in - 0.24), Inches(body_h_in),
               body, sz=SZ["body_sm"], bul=True)
        else:
            T(s, x + Inches(0.12), body_y, Inches(w_in - 0.24), Inches(body_h_in),
              body, sz=SZ["body_sm"], c=C["gray"], ls=1.4)

    return card_sh


def SPLIT_VISUAL(s, image_desc, title, body_items, y=None, h=None,
                 image_side="left", image_ratio=0.45, color=None):
    """좌우 이미지/텍스트 분할 레이아웃 (v3.8)

    한쪽은 이미지 플레이스홀더, 다른 쪽은 제목+불릿 텍스트.
    존 디테일, 프로그램 상세 등 Show-Don't-Tell 슬라이드에 적합.

    Args:
        image_desc: 이미지 설명
        title: 텍스트 영역 제목
        body_items: 불릿 리스트 또는 본문 문자열
        y: 시작 Y
        h: 높이 (기본 4.5")
        image_side: "left" 또는 "right"
        image_ratio: 이미지 너비 비율 (0.35~0.55)
        color: 제목 악센트 색상
    """
    y = _ensure_emu(y)
    h = _ensure_emu(h)
    if y is None:
        y = Inches(Z["ct_y"])
    if h is None:
        h = Inches(4.5)
    # ★ 경계 클램핑
    y_in = float(y / 914400)
    h_in_raw = float(h / 914400)
    max_h = 7.0 - y_in
    if h_in_raw > max_h and max_h > 1.0:
        h = Inches(max_h)
    if color is None:
        color = C["primary"]

    total_w = float(CW / 914400)
    gap = 0.25
    img_w = total_w * image_ratio
    txt_w = total_w - img_w - gap

    if image_side == "left":
        img_x = ML
        txt_x = ML + Inches(img_w + gap)
    else:
        txt_x = ML
        img_x = ML + Inches(txt_w + gap)

    # 이미지 플레이스홀더
    IMG_PH(s, img_x, y, Inches(img_w), h, image_desc)

    # 악센트 라인
    ACCENT_LINE(s, float(txt_x / 914400), float(y / 914400) + 0.05,
                float(h / 914400) * 0.15, color=color)

    # 제목
    T(s, txt_x + Inches(0.15), y + Inches(0.1), Inches(txt_w - 0.15), Inches(0.5),
      title, sz=SZ["subtitle"], c=C["dark"], b=True, fn=FONT_W["semibold"])

    # 본문
    body_y = y + Inches(0.7)
    body_h = h - Inches(0.8)
    if isinstance(body_items, list):
        MT(s, txt_x + Inches(0.15), body_y, Inches(txt_w - 0.15), body_h,
           body_items, sz=SZ["body"], bul=True)
    else:
        T(s, txt_x + Inches(0.15), body_y, Inches(txt_w - 0.15), body_h,
          body_items, sz=SZ["body"], c=C["gray"], ls=1.5)


def SECTION_BRIDGE(s, from_text, to_text, connector="", bg_style="dark"):
    """내러티브 전환 슬라이드 — 페이즈 간 스토리텔링 연결 (v3.8)

    "문제를 알았다 → 이제 해법을 보여드립니다" 같은 전환 메시지.
    최소 시각 요소로 평가위원의 mental shift를 유도.

    Args:
        from_text: 이전 섹션 요약 (작은 글씨)
        to_text: 다음 섹션 예고 (큰 글씨, 강조)
        connector: 연결 텍스트 (빈 문자열이면 기본 악센트 라인)
        bg_style: "dark" 또는 "light"
    """
    if bg_style == "dark":
        gradient_bg(s, darken(C["dark"], 0.3), C["dark"])
        from_c = C["lgray"]
        to_c = C["white"]
        accent_c = C["secondary"]
    else:
        bg(s, C["light"])
        from_c = C["gray"]
        to_c = C["dark"]
        accent_c = C["primary"]

    # from (이전 섹션 요약)
    T(s, ML + Inches(1.0), Inches(2.5), CW - Inches(2.0), Inches(0.5),
      from_text, sz=SZ["body"], c=from_c, al=PP_ALIGN.CENTER,
      fn=FONT_W["regular"])

    # 악센트 구분선
    if connector:
        T(s, ML + Inches(1.0), Inches(3.2), CW - Inches(2.0), Inches(0.4),
          connector, sz=SZ["body"], c=accent_c, b=True, al=PP_ALIGN.CENTER)
    else:
        DIVIDER(s, 3.3, style="accent", color=accent_c)

    # to (다음 섹션 예고)
    T(s, ML + Inches(0.5), Inches(3.8), CW - Inches(1.0), Inches(1.0),
      to_text, sz=SZ["action"], c=to_c, b=True, al=PP_ALIGN.CENTER,
      fn=FONT_W["bold"])


# ═══════════════════════════════════════════════════════════════════════
#  v4.0 "NEON" UPGRADE — reference-driven expansion (2026-04-17)
# ═══════════════════════════════════════════════════════════════════════
"""
레퍼런스: [LAON]2026_메이플스토리_월드_메커톤_제안서_0401.pptx (65p)

분석 결과 → 엔진 확장:
  - 타이포그래피 밀도화 (9~11pt 본문 tier 추가)
  - Semantic 컬러 토큰 + 다크 네온 팔레트
  - 3-stop 브랜드 그라디언트
  - 유리/네온/글로우 이펙트
  - 10개 신규 컴포넌트 (NEON_FRAME, GLASS_CARD, CHEVRON_FLOW 등)
  - 5개 신규 슬라이드 템플릿 (stat_hero, manifesto, dashboard 등)
  - gaming_tech 테마

모든 추가는 기존 v3.8 API와 완전 호환. 신규 이름은 충돌 없음.
"""

__version__ = "4.1"

# v4.0 NEON 실험은 레퍼런스 오분석에 기반하여 실전 적용 비권장.
# v4.1 에디토리얼 다크가 공식 권장 방향. (아래 v4.1 섹션 참조)
# v4.0 함수들(NEON_FRAME, DOT_PATTERN, GLASS_CARD, NEON_KPI, DENSE_GRID,
# SPLIT_DIAGONAL, TIMELINE_RIBBON, slide_manifesto, slide_dashboard 등)은
# 하위 호환을 위해 남겨두지만 신규 제안서에는 사용 지양.

# ───────────────────────────────────────────────────────────────
#  v4.0 / 1. 확장 타이포그래피 tier (backward compat: 기존 SZ 유지)
# ───────────────────────────────────────────────────────────────

SZ.update({
    # v4.0: 프레젠테이션 환경에 맞춰 레퍼런스 대비 2~4pt 상향
    "micro":       10,   # 각주, 디스클레이머 (원본 7 → 10)
    "caption_sm":  12,   # 밀도형 본문 (원본 9 → 12)
    "label":       13,   # 범례/라벨 (원본 10 → 13)
    "fine":        14,   # 촘촘 본문 (원본 11 → 14)
    "stat_hero":   110,  # 대형 수치 (HERO) — 더 큰 임팩트 (원본 96 → 110)
    "manifesto":   64,   # 선언문 제목 (원본 60 → 64)
    "eyebrow":     13,   # 오버라인 kicker (원본 11 → 13)
    "card_title":  22,   # 카드 타이틀 (신규)
    "section_lg":  32,   # 섹션 대형 (신규)
})


# ───────────────────────────────────────────────────────────────
#  v4.0 / 2. Semantic Color Tokens (9-step neutrals + accents)
# ───────────────────────────────────────────────────────────────
#
# TOKENS는 C와 별개 네임스페이스. 기존 C는 그대로 유지.
# 테마에 따라 업데이트되며, 새 v4.0 컴포넌트들이 참조.

NEUTRAL_LIGHT = [
    RGBColor(250, 251, 253),   # 50  — 최상위 표면
    RGBColor(242, 245, 250),   # 100 — 페이퍼 화이트
    RGBColor(235, 237, 240),   # 200 — 서브틀 표면 ⭐ 레퍼런스 dominant
    RGBColor(215, 220, 228),   # 300 — 경계선
    RGBColor(185, 192, 204),   # 400 — 보조
    RGBColor(128, 136, 152),   # 500 — 뮤티드 텍스트
    RGBColor(88, 96, 112),     # 600 — 서브 텍스트
    RGBColor(58, 64, 78),      # 700 — 텍스트
    RGBColor(41, 46, 58),      # 800 — 다크 표면 ⭐
    RGBColor(28, 31, 40),      # 900 — 딥 다크
]

TOKENS = {
    # 표면 / 배경
    "surface/raised":   NEUTRAL_LIGHT[2],        # EBEDF0
    "surface/base":     NEUTRAL_LIGHT[1],
    "surface/dark":     NEUTRAL_LIGHT[8],        # 292E3A
    "surface/darker":   NEUTRAL_LIGHT[9],        # 1C1F28
    "surface/overlay":  RGBColor(20, 24, 34),

    # 경계
    "border/subtle":    NEUTRAL_LIGHT[3],
    "border/dark":      RGBColor(47, 51, 64),    # 2F3340

    # 브랜드
    "brand/primary":    RGBColor(95, 112, 252),  # 5F70FC ⭐ 핵심 액센트
    "brand/secondary":  RGBColor(98, 150, 255),  # 6296FF
    "brand/deep":       RGBColor(0, 42, 128),    # 002A80

    # 네온 (다크 BG 위에서 빛나는)
    "neon/cyan":        RGBColor(102, 255, 255), # 66FFFF
    "neon/aqua":        RGBColor(105, 226, 255), # 69E2FF
    "neon/electric":    RGBColor(42, 246, 255),  # 2AF6FF
    "neon/mint":        RGBColor(82, 255, 196),
    "neon/yellow":      RGBColor(255, 224, 51),

    # 그라디언트 정지점 (브랜드 gradient)
    "grad/start":       RGBColor(15, 106, 199),  # 0F6AC7
    "grad/mid":         RGBColor(95, 112, 252),  # 5F70FC
    "grad/end":         RGBColor(185, 103, 255), # 보조 보라

    # 텍스트
    "text/on_dark":     RGBColor(255, 255, 255),
    "text/on_light":    NEUTRAL_LIGHT[9],
    "text/muted":       NEUTRAL_LIGHT[5],
    "text/subtle":      NEUTRAL_LIGHT[6],
    "text/accent":      RGBColor(95, 112, 252),
}

# v4.1 편의 alias — 짧은 이름으로 자주 참조되는 컬러
TOKENS.update({
    "primary":      TOKENS["brand/primary"],      # 퍼플 악센트
    "secondary":    TOKENS["neon/cyan"],           # 사이언 eyebrow 전용
    "accent":       TOKENS["brand/deep"],
    "text":         TOKENS["text/on_dark"],
    "muted":        TOKENS["text/muted"],
    "bg":           TOKENS["surface/darker"],
    "card":         TOKENS["surface/dark"],
    "border":       TOKENS["border/subtle"],
})


def tok(key):
    """TOKENS 빠른 접근자. tok('brand/primary') → RGBColor."""
    return TOKENS[key]


# ───────────────────────────────────────────────────────────────
#  v4.0 / 3. 신규 테마: gaming_tech (다크 네온)
# ───────────────────────────────────────────────────────────────

THEMES["gaming_tech"] = {
    "primary":   (95, 112, 252),   # 5F70FC
    "secondary": (98, 150, 255),   # 6296FF
    "teal":      (102, 255, 255),  # 66FFFF
    "accent":    (185, 103, 255),  # purple
    "dark":      (28, 31, 40),     # 1C1F28
    "light":     (235, 237, 240),  # EBEDF0
}


# ───────────────────────────────────────────────────────────────
#  v4.0 / 4. 확장 그라디언트 — 3-stop & 레퍼런스 기반
# ───────────────────────────────────────────────────────────────

GRAD.update({
    # 2-stop 기본 브랜드 그라디언트 (블루 → 퍼플)
    "brand":         lambda: (tok("grad/start"), tok("grad/mid")),
    # 3-stop 전체 브랜드 (사용 시 gradient_shape_3stop 호출)
    "brand_3stop":   lambda: (tok("grad/start"), tok("grad/mid"), tok("grad/end")),
    # 네온 사이언 페이드
    "neon_cyan":     lambda: (tok("neon/aqua"), tok("neon/electric")),
    # 다크 아웃라인 글로우
    "dark_glow":     lambda: (tok("surface/darker"), tok("surface/dark")),
    # 섹션 디바이더용 (딥 네이비 → 브랜드)
    "section_dark":  lambda: (tok("brand/deep"), tok("brand/primary")),
})


def gradient_shape_3stop(shape, c1, c2, c3, angle=5400000):
    """3-stop 그라디언트 채우기. angle=5400000 → 90° (수직).

    angle: 21600000 = 360°. 0=좌→우, 5400000=상→하, 10800000=우→좌.
    """
    try:
        from lxml import etree
    except ImportError:
        return shape
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    spPr = shape._element.spPr
    # 기존 fill 제거
    for tag in ('solidFill', 'gradFill', 'noFill', 'blipFill', 'pattFill'):
        for old in spPr.findall(f'{{{ns}}}{tag}'):
            spPr.remove(old)
    # 앞에 insert (fill은 spPr 첫 자식 중 하나)
    gradFill = etree.SubElement(spPr, f'{{{ns}}}gradFill',
                                 flip='none', rotWithShape='1')
    gsLst = etree.SubElement(gradFill, f'{{{ns}}}gsLst')
    for pos, rgb in [(0, c1), (50000, c2), (100000, c3)]:
        gs = etree.SubElement(gsLst, f'{{{ns}}}gs', pos=str(pos))
        etree.SubElement(gs, f'{{{ns}}}srgbClr',
                          val=str(rgb).lstrip('#').upper() if isinstance(rgb, str)
                              else f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")
    etree.SubElement(gradFill, f'{{{ns}}}lin', ang=str(angle), scaled='1')
    # spPr 첫 자식 뒤로 이동 (xfrm 뒤)
    return shape


# ───────────────────────────────────────────────────────────────
#  v4.0 / 5. 신규 이펙트 — glow, alpha, gradient_text, duotone
# ───────────────────────────────────────────────────────────────

def add_glow(shape, color=None, blur_pt=20, alpha=60):
    """네온 글로우 이펙트 (outer glow).

    Args:
        color: RGBColor 또는 (r,g,b). None이면 tok("neon/cyan") 사용
        blur_pt: 글로우 반경 (pt). 클수록 부드러운 퍼짐
        alpha: 0~100 (%). 작을수록 투명
    """
    try:
        from lxml import etree
    except ImportError:
        return shape
    if color is None:
        color = tok("neon/cyan")
    rgb_hex = _rgb_to_hex(color)

    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    spPr = shape._element.spPr
    effectLst = spPr.find(f'{{{ns}}}effectLst')
    if effectLst is None:
        effectLst = etree.SubElement(spPr, f'{{{ns}}}effectLst')
    glow = etree.SubElement(effectLst, f'{{{ns}}}glow',
                             rad=str(blur_pt * 12700))
    srgb = etree.SubElement(glow, f'{{{ns}}}srgbClr', val=rgb_hex)
    # alpha는 100000 스케일 (100% = 100000)
    etree.SubElement(srgb, f'{{{ns}}}alpha', val=str(alpha * 1000))
    return shape


def add_alpha(shape, alpha_pct=50):
    """도형의 fill에 알파(투명도) 적용. 0=불투명, 100=완전투명."""
    try:
        from lxml import etree
    except ImportError:
        return shape
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    spPr = shape._element.spPr
    # solidFill 내 srgbClr 밑에 a:alpha 삽입
    solidFill = spPr.find(f'{{{ns}}}solidFill')
    if solidFill is not None:
        srgb = solidFill.find(f'{{{ns}}}srgbClr')
        if srgb is not None:
            # 기존 alpha 제거
            for old in srgb.findall(f'{{{ns}}}alpha'):
                srgb.remove(old)
            etree.SubElement(srgb, f'{{{ns}}}alpha',
                              val=str((100 - alpha_pct) * 1000))
    return shape


def _rgb_to_hex(c):
    """RGBColor 또는 (r,g,b) 튜플 → 'RRGGBB' hex (대문자)."""
    if hasattr(c, '__iter__') and not isinstance(c, str):
        parts = list(c)
        if len(parts) >= 3:
            return f"{int(parts[0]):02X}{int(parts[1]):02X}{int(parts[2]):02X}"
    # RGBColor는 __iter__로 (r,g,b) 반환
    try:
        s = str(c)
        if len(s) == 6:
            return s.upper()
    except Exception:
        pass
    return "000000"


def gradient_text(run, color_start, color_end):
    """텍스트 run에 그라디언트 fill 적용.

    주의: 작은 폰트(<20pt)에서는 효과 미미. 36pt 이상 권장.

    Args:
        run: TextFrame.runs[i] (python-pptx)
        color_start, color_end: RGBColor 또는 hex str
    """
    try:
        from lxml import etree
    except ImportError:
        return run
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    rPr = run._r.get_or_add_rPr()
    # 기존 fill 제거
    for tag in ('solidFill', 'gradFill', 'noFill'):
        for old in rPr.findall(f'{{{ns}}}{tag}'):
            rPr.remove(old)
    # gradFill 생성 (아직 insert 안 함)
    gradFill = etree.Element(f'{{{ns}}}gradFill', flip='none',
                              rotWithShape='1')
    gsLst = etree.SubElement(gradFill, f'{{{ns}}}gsLst')
    gs1 = etree.SubElement(gsLst, f'{{{ns}}}gs', pos='0')
    etree.SubElement(gs1, f'{{{ns}}}srgbClr', val=_rgb_to_hex(color_start))
    gs2 = etree.SubElement(gsLst, f'{{{ns}}}gs', pos='100000')
    etree.SubElement(gs2, f'{{{ns}}}srgbClr', val=_rgb_to_hex(color_end))
    # OOXML 순서: fill 요소는 latin/ea/cs보다 앞에 와야 함
    # 첫 latin 또는 ea 또는 cs 엘리먼트 앞에 삽입
    insert_idx = len(rPr)
    for i, child in enumerate(rPr):
        tag_name = child.tag.split('}')[-1]
        if tag_name in ('latin', 'ea', 'cs', 'sym', 'hlinkClick',
                         'hlinkMouseOver', 'rtl'):
            insert_idx = i
            break
    rPr.insert(insert_idx, gradFill)
    etree.SubElement(gradFill, f'{{{ns}}}lin', ang='0', scaled='1')
    return run


def duotone_overlay(s, l, t, w, h, dark_color=None, light_color=None, alpha=70):
    """이미지 위에 듀오톤 풍의 그라디언트 오버레이 반투명 레이어.

    실제 duotone 필터가 아닌 그라디언트+알파 조합으로 유사 효과.
    """
    dark_color = dark_color or tok("surface/darker")
    light_color = light_color or tok("brand/primary")
    rect = R(s, l, t, w, h, dark_color)
    gradient_shape(rect, dark_color, light_color, angle=2700000)
    add_alpha(rect, alpha)
    return rect


# ───────────────────────────────────────────────────────────────
#  v4.0 / 6. 신규 컴포넌트 10종
# ───────────────────────────────────────────────────────────────

def NEON_FRAME(s, l, t, w, h, color=None, glow_blur=18, line_pt=1.5,
               rounded=True, corner_radius=0.02):
    """네온 글로우 테두리 프레임 — 배경 없는 빛나는 외곽선.

    다크 배경 위에서 극적인 임팩트. 컨셉 reveal / 강조 박스용.
    """
    color = color or tok("neon/cyan")
    if rounded:
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    _ei(l), _ei(t), _ei(w), _ei(h))
        try:
            shape.adjustments[0] = corner_radius
        except Exception:
            pass
    else:
        shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    _ei(l), _ei(t), _ei(w), _ei(h))
    # 배경 투명
    shape.fill.background()
    # 아웃라인: 네온 컬러
    shape.line.color.rgb = color
    shape.line.width = Pt(line_pt)
    # 글로우
    add_glow(shape, color=color, blur_pt=glow_blur, alpha=70)
    return shape


def GLASS_CARD(s, l, t, w, h, *, alpha=40, border=True, border_color=None,
               tint=None, rounded=True):
    """반투명 유리 카드 (glassmorphism 근사).

    Args:
        alpha: 0~100. 40 = 60% 불투명. 레퍼런스 빈번 값.
        border: True면 얇은 테두리
        tint: fill 기본색. None이면 tok("surface/raised")
        rounded: 라운드 사각형 여부
    """
    tint = tint or tok("surface/raised")
    border_color = border_color or tok("border/subtle")

    if rounded:
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    _ei(l), _ei(t), _ei(w), _ei(h))
        try:
            shape.adjustments[0] = 0.03
        except Exception:
            pass
    else:
        shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    _ei(l), _ei(t), _ei(w), _ei(h))

    shape.fill.solid()
    shape.fill.fore_color.rgb = tint
    add_alpha(shape, alpha)

    if border:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()

    return shape


def _ei(v):
    """EMU-safe int helper (기존 _safe_int alias, 짧게)."""
    return _safe_int(v)


def CHEVRON_FLOW(s, l_in, t_in, w_in, h_in, items, *,
                  fill_color=None, text_color=None, gap_in=0.05):
    """쉐브론 기반 프로세스 플로우 — 모던 게임UI 느낌.

    Args:
        items: [str, ...] — 각 쉐브론의 텍스트
        gap_in: 쉐브론 간 겹침/간격
    """
    fill_color = fill_color or tok("brand/primary")
    text_color = text_color or tok("text/on_dark")
    n = len(items)
    if n == 0:
        return
    total_w = w_in
    each_w = (total_w - gap_in * (n - 1)) / n

    for i, text in enumerate(items):
        left = l_in + i * (each_w + gap_in)
        shape = s.shapes.add_shape(MSO_SHAPE.CHEVRON,
                                    Inches(left), Inches(t_in),
                                    Inches(each_w), Inches(h_in))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        # 텍스트
        tf = shape.text_frame
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = FONT_W["semibold"]
        run.font.size = Pt(SZ["fine"])
        run.font.color.rgb = text_color
        run.font.bold = True


def PARALLELOGRAM_BADGE(s, l_in, t_in, w_in, h_in, text, *,
                         color=None, text_color=None, sz_pt=None):
    """평행사변형 뱃지 — 섹션 라벨/카테고리 태그."""
    color = color or tok("brand/primary")
    text_color = text_color or tok("text/on_dark")
    sz_pt = sz_pt or SZ["label"]

    shape = s.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
                                Inches(l_in), Inches(t_in),
                                Inches(w_in), Inches(h_in))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.25
    except Exception:
        pass

    tf = shape.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_W["bold"]
    run.font.size = Pt(sz_pt)
    run.font.color.rgb = text_color
    run.font.bold = True
    return shape


def DOT_PATTERN(s, l_in, t_in, w_in, h_in, *,
                 dot_size_pt=3, gap_in=0.6, color=None, alpha_pct=50,
                 max_dots=150):
    """장식 도트 패턴 배경 — 섹션 구분/배경 텍스처.

    성능: max_dots 상한 강제 (PPTX에 수백 개 도형 생성 방지).

    Args:
        gap_in: 도트 간격 (기본 0.6" — 성능 고려)
        max_dots: 총 도트 수 상한 (초과 시 gap 자동 증가)
    """
    color = color or tok("text/muted")
    # 성능 보호: 도트 수 상한
    est = (w_in / gap_in) * (h_in / gap_in)
    if est > max_dots:
        scale = (est / max_dots) ** 0.5
        gap_in = gap_in * scale
    cols = int(w_in / gap_in) + 1
    rows = int(h_in / gap_in) + 1
    dot_size_in = dot_size_pt / 72
    count = 0
    for r in range(rows):
        for c in range(cols):
            if count >= max_dots:
                return
            x = l_in + c * gap_in
            y = t_in + r * gap_in
            if x + dot_size_in > l_in + w_in:
                continue
            if y + dot_size_in > t_in + h_in:
                continue
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                      Inches(x), Inches(y),
                                      Inches(dot_size_in), Inches(dot_size_in))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()
            add_alpha(dot, alpha_pct)
            count += 1


def STAT_HERO(s, l_in, t_in, w_in, h_in, value, label, *,
               unit="", sub="", color=None, label_color=None,
               value_sz=None, label_sz=None, align="center"):
    """대형 수치 히어로 — 핵심 통계 강조.

    레퍼런스 feel: 96pt 수치 + 얇은 라벨 + 작은 설명
    """
    color = color or tok("brand/primary")
    label_color = label_color or tok("text/muted")
    value_sz = value_sz or SZ["stat_hero"]
    label_sz = label_sz or SZ["label"]

    al = {"center": PP_ALIGN.CENTER, "left": PP_ALIGN.LEFT,
          "right": PP_ALIGN.RIGHT}[align]

    # 값 + 단위
    val_h = value_sz / 72 * 1.25
    label_h = label_sz / 72 * 2.0

    val_y = t_in + (h_in - val_h - label_h - 0.15) / 2

    # 큰 숫자
    val_box = s.shapes.add_textbox(Inches(l_in), Inches(val_y),
                                    Inches(w_in), Inches(val_h))
    tf = val_box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]
    p.alignment = al
    run = p.add_run()
    run.text = str(value)
    run.font.name = FONT_W["black"]
    run.font.size = Pt(value_sz)
    run.font.color.rgb = color
    run.font.bold = True
    if unit:
        run2 = p.add_run()
        run2.text = " " + unit
        run2.font.name = FONT_W["medium"]
        run2.font.size = Pt(int(value_sz * 0.35))
        run2.font.color.rgb = label_color

    # 라벨
    lab_box = s.shapes.add_textbox(Inches(l_in),
                                    Inches(val_y + val_h + 0.05),
                                    Inches(w_in), Inches(label_h))
    ltf = lab_box.text_frame
    ltf.margin_left = ltf.margin_right = Inches(0.05)
    lp = ltf.paragraphs[0]
    lp.alignment = al
    lrun = lp.add_run()
    lrun.text = label
    lrun.font.name = FONT_W["medium"]
    lrun.font.size = Pt(label_sz)
    lrun.font.color.rgb = label_color

    if sub:
        sub_box = s.shapes.add_textbox(Inches(l_in),
                                        Inches(val_y + val_h + 0.05 + label_sz/72*1.3),
                                        Inches(w_in), Inches(0.3))
        stf = sub_box.text_frame
        stf.margin_left = stf.margin_right = Inches(0.05)
        sp = stf.paragraphs[0]
        sp.alignment = al
        srun = sp.add_run()
        srun.text = sub
        srun.font.name = FONT_W["regular"]
        srun.font.size = Pt(SZ["caption_sm"])
        srun.font.color.rgb = tok("text/subtle")


def NEON_KPI(s, l_in, t_in, w_in, h_in, items, *,
              cols=None, gap_in=0.15, card_color=None, accent=None):
    """네온 아웃라인 KPI 카드 그리드.

    items: [{"value": "150%", "label": "성장률", "sub": "YoY"}, ...]
    """
    card_color = card_color or tok("surface/dark")
    accent = accent or tok("neon/cyan")
    n = len(items)
    cols = cols or n
    rows = (n + cols - 1) // cols
    each_w = (w_in - gap_in * (cols - 1)) / cols
    each_h = (h_in - gap_in * (rows - 1)) / rows

    for i, it in enumerate(items):
        r = i // cols
        c = i % cols
        x = l_in + c * (each_w + gap_in)
        y = t_in + r * (each_h + gap_in)

        # 다크 카드 배경
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y),
                                   Inches(each_w), Inches(each_h))
        try:
            card.adjustments[0] = 0.04
        except Exception:
            pass
        card.fill.solid()
        card.fill.fore_color.rgb = card_color
        card.line.color.rgb = accent
        card.line.width = Pt(1.2)
        add_glow(card, color=accent, blur_pt=12, alpha=50)

        # 값
        val_box = s.shapes.add_textbox(Inches(x), Inches(y + each_h * 0.18),
                                        Inches(each_w), Inches(each_h * 0.45))
        tf = val_box.text_frame
        tf.margin_left = tf.margin_right = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(it.get("value", ""))
        run.font.name = FONT_W["black"]
        run.font.size = Pt(36)
        run.font.color.rgb = accent
        run.font.bold = True

        # 라벨
        lab_box = s.shapes.add_textbox(Inches(x),
                                        Inches(y + each_h * 0.62),
                                        Inches(each_w),
                                        Inches(each_h * 0.2))
        ltf = lab_box.text_frame
        ltf.margin_left = ltf.margin_right = Inches(0.05)
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        lrun = lp.add_run()
        lrun.text = it.get("label", "")
        lrun.font.name = FONT_W["semibold"]
        lrun.font.size = Pt(SZ["fine"])
        lrun.font.color.rgb = tok("text/on_dark")

        # 서브
        if it.get("sub"):
            sub_box = s.shapes.add_textbox(Inches(x),
                                            Inches(y + each_h * 0.80),
                                            Inches(each_w),
                                            Inches(each_h * 0.14))
            stf = sub_box.text_frame
            stf.margin_left = stf.margin_right = Inches(0.05)
            sp = stf.paragraphs[0]
            sp.alignment = PP_ALIGN.CENTER
            srun = sp.add_run()
            srun.text = it["sub"]
            srun.font.name = FONT_W["regular"]
            srun.font.size = Pt(SZ["caption_sm"])
            srun.font.color.rgb = tok("text/muted")


def SPLIT_DIAGONAL(s, *, left_color=None, right_color=None,
                    angle_deg=12, split_ratio=0.5):
    """다이아고널 2-분할 배경 레이어.

    슬라이드 전체를 두 색으로 분할하되, 경계선이 기울어진 형태.

    Args:
        split_ratio: 0~1. 좌측 영역 비율 (기본 0.5)
        angle_deg: 분할선 기울기 (양수=우하향)
    """
    left_color = left_color or tok("surface/dark")
    right_color = right_color or tok("surface/darker")
    # 좌측 전체 채움
    R(s, 0, 0, SW, SH, left_color)
    # 우측을 평행사변형으로 덮기
    from math import tan, radians
    sh_in = SH / 914400
    sw_in = SW / 914400
    offset_in = tan(radians(angle_deg)) * sh_in
    # 평행사변형: 좌상단 (split_ratio*sw), 우상단 (sw), 우하단 (sw), 좌하단 (split_ratio*sw - offset)
    # python-pptx add_freeform은 복잡하므로 add_shape(MSO_SHAPE.PARALLELOGRAM) 사용
    para_l = split_ratio * sw_in - offset_in / 2
    para_w = sw_in - para_l + offset_in
    para = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                               Inches(para_l), Inches(0),
                               Inches(para_w * 1.2), Inches(sh_in))
    # 삼각형은 정확한 대각선 분할에 적합
    para.fill.solid()
    para.fill.fore_color.rgb = right_color
    para.line.fill.background()
    # 180도 플립해서 좌하향 → 우상단으로
    try:
        para.rotation = 0
    except Exception:
        pass
    return para


def DENSE_GRID(s, l_in, t_in, w_in, h_in, items, *,
                cols=4, gap_in=0.12, card_color=None, text_color=None):
    """9~12개 아이템 대시보드 그리드 — 정보 밀도 높은 레이아웃.

    items: [{"title": str, "value": str, "desc": str}, ...]
    """
    card_color = card_color or tok("surface/raised")
    text_color = text_color or tok("text/on_light")
    n = len(items)
    rows = (n + cols - 1) // cols
    each_w = (w_in - gap_in * (cols - 1)) / cols
    each_h = (h_in - gap_in * (rows - 1)) / rows

    for i, it in enumerate(items):
        r = i // cols
        c = i % cols
        x = l_in + c * (each_w + gap_in)
        y = t_in + r * (each_h + gap_in)

        GLASS_CARD(s, x, y, each_w, each_h, alpha=0, border=True,
                    tint=card_color)

        # 타이틀
        t_box = s.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.12),
                                      Inches(each_w - 0.3), Inches(0.3))
        tf = t_box.text_frame
        tf.margin_left = tf.margin_right = Inches(0.0)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = it.get("title", "")
        run.font.name = FONT_W["bold"]
        run.font.size = Pt(SZ["caption_sm"])
        run.font.color.rgb = tok("brand/primary")
        run.font.bold = True

        # 값
        v_box = s.shapes.add_textbox(Inches(x + 0.15),
                                      Inches(y + each_h * 0.35),
                                      Inches(each_w - 0.3),
                                      Inches(each_h * 0.35))
        vtf = v_box.text_frame
        vp = vtf.paragraphs[0]
        vr = vp.add_run()
        vr.text = it.get("value", "")
        vr.font.name = FONT_W["bold"]
        vr.font.size = Pt(18)
        vr.font.color.rgb = text_color
        vr.font.bold = True

        # 설명
        d_box = s.shapes.add_textbox(Inches(x + 0.15),
                                      Inches(y + each_h * 0.72),
                                      Inches(each_w - 0.3),
                                      Inches(each_h * 0.25))
        dtf = d_box.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dr = dp.add_run()
        dr.text = it.get("desc", "")
        dr.font.name = FONT_W["regular"]
        dr.font.size = Pt(SZ["caption_sm"])
        dr.font.color.rgb = tok("text/muted")


def TIMELINE_RIBBON(s, items, *, y_in=None, h_in=0.7,
                     color=None, text_color=None):
    """리본 스타일 타임라인 — 가로로 흐르는 얇은 리본.

    items: [("2024 Q1", "런칭"), ("2024 Q2", "성장"), ...]
    """
    color = color or tok("brand/primary")
    text_color = text_color or tok("text/on_dark")
    if y_in is None:
        y_in = 3.5
    l_in = 0.5
    w_in = 12.0
    n = len(items)
    seg_w = w_in / n

    # 리본 본체
    ribbon = R(s, Inches(l_in), Inches(y_in), Inches(w_in), Inches(h_in),
                color)
    try:
        # 리본을 라운드로
        pass
    except Exception:
        pass

    for i, item in enumerate(items):
        label, desc = item if isinstance(item, tuple) else (item.get("label", ""), item.get("desc", ""))
        x = l_in + i * seg_w
        # 구분선
        if i > 0:
            line = R(s, Inches(x), Inches(y_in + 0.05),
                      Inches(0.01), Inches(h_in - 0.1),
                      tok("text/on_dark"))
            add_alpha(line, 50)

        # 라벨 (상단)
        l_box = s.shapes.add_textbox(Inches(x + 0.05),
                                      Inches(y_in - 0.5),
                                      Inches(seg_w - 0.1),
                                      Inches(0.35))
        ltf = l_box.text_frame
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        lrun = lp.add_run()
        lrun.text = label
        lrun.font.name = FONT_W["bold"]
        lrun.font.size = Pt(SZ["fine"])
        lrun.font.color.rgb = tok("brand/primary")
        lrun.font.bold = True

        # 설명 (리본 내부)
        d_box = s.shapes.add_textbox(Inches(x + 0.05),
                                      Inches(y_in + 0.12),
                                      Inches(seg_w - 0.1),
                                      Inches(h_in - 0.24))
        dtf = d_box.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.alignment = PP_ALIGN.CENTER
        dr = dp.add_run()
        dr.text = desc
        dr.font.name = FONT_W["medium"]
        dr.font.size = Pt(SZ["label"])
        dr.font.color.rgb = text_color


# ───────────────────────────────────────────────────────────────
#  v4.0 / 7. 신규 슬라이드 템플릿 5종
# ───────────────────────────────────────────────────────────────

def slide_stat_hero(prs, title, value, label, *,
                     unit="", sub="", kicker="",
                     bg_dark=True, accent_color=None, pg=None):
    """대형 수치 히어로 슬라이드.

    Args:
        title: 상단 작은 제목 (Action Title)
        value: 거대한 숫자 (예: "87%")
        label: 숫자 아래 설명 (예: "타겟 도달률")
        unit: 단위 (값에 포함되지 않은 경우)
        sub: 라벨 아래 보조 설명
        kicker: 타이틀 위 eyebrow (예: "핵심 지표 01")
        bg_dark: 다크 배경 여부
    """
    s = new_slide(prs)
    accent = accent_color or tok("brand/primary")

    if bg_dark:
        bg(s, tok("surface/darker"))
        title_color = tok("text/on_dark")
        label_color = tok("text/muted")
        DOT_PATTERN(s, 0.5, 0.5, 12.0, 6.3, dot_size_pt=2,
                     gap_in=0.3, color=tok("border/dark"), alpha_pct=30)
    else:
        bg(s, tok("surface/base"))
        title_color = tok("text/on_light")
        label_color = tok("text/muted")

    # Kicker
    if kicker:
        T(s, ML, Inches(0.7), CW, Inches(0.3), kicker,
          sz=SZ["eyebrow"], c=accent, b=True, al=PP_ALIGN.CENTER,
          fn=FONT_W["bold"])

    # 타이틀
    T(s, ML, Inches(1.0 if kicker else 1.2), CW, Inches(0.6),
      title, sz=SZ["action"], c=title_color, al=PP_ALIGN.CENTER,
      fn=FONT_W["semibold"])

    # 히어로 수치
    STAT_HERO(s, l_in=1.0, t_in=2.4, w_in=11.3, h_in=3.2,
               value=value, label=label, unit=unit, sub=sub,
               color=accent,
               label_color=tok("text/on_dark") if bg_dark else tok("text/on_light"),
               align="center")

    if pg is not None:
        PN(s, pg)


def slide_manifesto(prs, statement, *, attribution="",
                      bg_dark=True, accent_color=None, pg=None):
    """대형 선언문 슬라이드 — 큰 메시지 한 문장만.

    Args:
        statement: 선언문 전체 텍스트 (여러 줄 가능)
        attribution: 하단 서명/출처
    """
    s = new_slide(prs)
    accent = accent_color or tok("brand/primary")

    if bg_dark:
        # 그라디언트 배경
        gradient_bg(s, tok("surface/darker"), tok("surface/dark"))
        text_color = tok("text/on_dark")
    else:
        gradient_bg(s, tok("surface/base"), tok("surface/raised"))
        text_color = tok("text/on_light")

    # 좌측 악센트 바 (네온 느낌)
    bar = R(s, Inches(0.8), Inches(2.2), Inches(0.08), Inches(3.1), accent)
    add_glow(bar, color=accent, blur_pt=12, alpha=60)

    # 큰 선언문
    T(s, Inches(1.4), Inches(2.2), Inches(10.6), Inches(3.1),
      statement, sz=SZ["manifesto"], c=text_color, b=True,
      fn=FONT_W["bold"], al=PP_ALIGN.LEFT)

    # 서명
    if attribution:
        T(s, Inches(1.4), Inches(5.6), Inches(10.6), Inches(0.4),
          "— " + attribution, sz=SZ["caption_sm"], c=accent,
          fn=FONT_W["medium"], al=PP_ALIGN.LEFT)

    if pg is not None:
        PN(s, pg)


def slide_comparison_3way(prs, title, columns, *, pg=None,
                            bg_dark=False, highlight_idx=None):
    """3-way 비교 슬라이드.

    columns: [{"title": str, "body": [str, ...], "badge": str?}, ...]
    highlight_idx: 강조할 컬럼 인덱스 (0~2). 해당 컬럼에 네온 프레임.
    """
    s = new_slide(prs)
    if bg_dark:
        bg(s, tok("surface/darker"))
        title_color = tok("text/on_dark")
        card_color = tok("surface/dark")
        border_color = tok("border/dark")
        body_color = tok("text/on_dark")
    else:
        bg(s, tok("surface/base"))
        title_color = tok("text/on_light")
        card_color = tok("surface/raised")
        border_color = tok("border/subtle")
        body_color = tok("text/on_light")

    # 타이틀
    TB(s, title, pg=pg)

    # 3 컬럼
    start_y = 1.6
    card_h = 5.0
    total_w = 12.0
    card_gap = 0.2
    card_w = (total_w - card_gap * 2) / 3
    start_x = 0.67

    for i, col in enumerate(columns[:3]):
        x = start_x + i * (card_w + card_gap)

        if highlight_idx == i:
            NEON_FRAME(s, Inches(x - 0.04), Inches(start_y - 0.04),
                        Inches(card_w + 0.08), Inches(card_h + 0.08),
                        color=tok("brand/primary"), glow_blur=15, line_pt=2)

        card = GLASS_CARD(s, x, start_y, card_w, card_h,
                           alpha=0, border=True, tint=card_color,
                           border_color=border_color)

        # 뱃지 (있으면)
        if col.get("badge"):
            PARALLELOGRAM_BADGE(s, x + 0.25, start_y + 0.2, 1.5, 0.35,
                                 col["badge"],
                                 color=tok("brand/primary"))

        # 컬럼 타이틀
        T(s, Inches(x + 0.3),
          Inches(start_y + (0.7 if col.get("badge") else 0.3)),
          Inches(card_w - 0.6), Inches(0.6),
          col["title"], sz=24, c=body_color, b=True,
          fn=FONT_W["bold"])

        # 본문
        MT(s, Inches(x + 0.3),
           Inches(start_y + 1.5),
           Inches(card_w - 0.6),
           Inches(card_h - 2.0),
           col.get("body", []), sz=SZ["fine"],
           c=body_color if not bg_dark else tok("text/on_dark"),
           bul=True)


def slide_dashboard(prs, title, items, *, cols=4, pg=None,
                     bg_dark=False):
    """정보 밀도 대시보드 슬라이드 — 9~12개 메트릭 카드."""
    s = new_slide(prs)
    if bg_dark:
        bg(s, tok("surface/darker"))
    else:
        bg(s, tok("surface/base"))
    TB(s, title, pg=pg)

    grid_h = 5.0
    DENSE_GRID(s, l_in=0.67, t_in=1.5, w_in=12.0, h_in=grid_h,
                items=items, cols=cols,
                card_color=tok("surface/dark") if bg_dark else tok("surface/raised"),
                text_color=tok("text/on_dark") if bg_dark else tok("text/on_light"))


def slide_timeline_ribbon(prs, title, items, *, pg=None, bg_dark=False,
                            color=None):
    """리본 스타일 타임라인 슬라이드.

    items: [("2024 Q1", "런칭 준비"), ("Q2", "MVP 출시"), ...]
    """
    s = new_slide(prs)
    if bg_dark:
        bg(s, tok("surface/darker"))
    else:
        bg(s, tok("surface/base"))
    TB(s, title, pg=pg)

    TIMELINE_RIBBON(s, items, y_in=3.3, h_in=1.1,
                     color=color or tok("brand/primary"),
                     text_color=tok("text/on_dark"))


# ───────────────────────────────────────────────────────────────
#  v4.0 / 8. 기존 컴포넌트 확장 변형 (non-breaking)
# ───────────────────────────────────────────────────────────────

def HIGHLIGHT_NEON(s, title, sub="", y_in=1.5, color=None):
    """HIGHLIGHT의 네온 버전 — 다크 배경 + 글로우."""
    color = color or tok("neon/cyan")
    h_in = 1.4 if sub else 1.0

    # 다크 배경 카드
    card = R(s, Inches(0.5), Inches(y_in), Inches(12.3), Inches(h_in),
              tok("surface/dark"))
    add_glow(card, color=color, blur_pt=14, alpha=45)

    # 왼쪽 네온 악센트
    bar = R(s, Inches(0.5), Inches(y_in), Inches(0.08), Inches(h_in), color)
    add_glow(bar, color=color, blur_pt=20, alpha=80)

    # 제목
    T(s, Inches(0.8), Inches(y_in + 0.2),
      Inches(11.8), Inches(0.6), title,
      sz=28, c=tok("text/on_dark"), b=True, fn=FONT_W["bold"])

    if sub:
        T(s, Inches(0.8), Inches(y_in + 0.85),
          Inches(11.8), Inches(0.4), sub,
          sz=SZ["fine"], c=tok("text/muted"), fn=FONT_W["regular"])


def slide_cover_neon(prs, title, subtitle, client="", *,
                      kicker="", accent_color=None):
    """네온 스타일 표지 — 다크 + 브랜드 그라디언트."""
    s = new_slide(prs)
    accent = accent_color or tok("brand/primary")

    # 그라디언트 배경
    gradient_bg(s, tok("surface/darker"), tok("surface/dark"))

    # 도트 패턴 배경
    DOT_PATTERN(s, 0.5, 0.5, 12.3, 6.5, dot_size_pt=2,
                 gap_in=0.3, color=tok("border/dark"), alpha_pct=40)

    # Kicker (있으면)
    if kicker:
        PARALLELOGRAM_BADGE(s, 1.0, 2.0, 2.2, 0.4, kicker,
                             color=accent)

    # 메인 타이틀
    T(s, Inches(1.0), Inches(2.6 if kicker else 2.2),
      Inches(11.3), Inches(1.6), title,
      sz=SZ["hero"], c=tok("text/on_dark"), b=True,
      fn=FONT_W["black"], al=PP_ALIGN.LEFT)

    # 서브타이틀
    T(s, Inches(1.0), Inches(4.4), Inches(11.3), Inches(0.8),
      subtitle, sz=SZ["action"], c=accent,
      fn=FONT_W["semibold"], al=PP_ALIGN.LEFT)

    # 클라이언트
    if client:
        bar = R(s, Inches(1.0), Inches(5.8), Inches(0.04), Inches(0.5),
                 accent)
        T(s, Inches(1.15), Inches(5.8), Inches(11.0), Inches(0.5),
          client, sz=SZ["fine"], c=tok("text/on_dark"),
          fn=FONT_W["medium"], al=PP_ALIGN.LEFT)


# ───────────────────────────────────────────────────────────────
#  v4.0 / 9. 편의 헬퍼
# ───────────────────────────────────────────────────────────────

def list_v4_components():
    """v4.0 신규 컴포넌트/템플릿 목록."""
    items = [
        ("COMPONENT", "NEON_FRAME", "네온 글로우 테두리"),
        ("COMPONENT", "GLASS_CARD", "반투명 유리 카드"),
        ("COMPONENT", "CHEVRON_FLOW", "쉐브론 프로세스"),
        ("COMPONENT", "PARALLELOGRAM_BADGE", "기울어진 뱃지"),
        ("COMPONENT", "DOT_PATTERN", "도트 패턴 배경"),
        ("COMPONENT", "STAT_HERO", "대형 수치"),
        ("COMPONENT", "NEON_KPI", "네온 KPI 카드"),
        ("COMPONENT", "SPLIT_DIAGONAL", "대각선 2분할"),
        ("COMPONENT", "DENSE_GRID", "정보 밀도 그리드"),
        ("COMPONENT", "TIMELINE_RIBBON", "리본 타임라인"),
        ("EFFECT", "add_glow", "네온 글로우"),
        ("EFFECT", "add_alpha", "투명도"),
        ("EFFECT", "gradient_text", "그라디언트 텍스트"),
        ("EFFECT", "duotone_overlay", "듀오톤 오버레이"),
        ("EFFECT", "gradient_shape_3stop", "3-stop 그라디언트"),
        ("TEMPLATE", "slide_stat_hero", "대형 수치 슬라이드"),
        ("TEMPLATE", "slide_manifesto", "선언문 슬라이드"),
        ("TEMPLATE", "slide_comparison_3way", "3-way 비교"),
        ("TEMPLATE", "slide_dashboard", "정보 대시보드"),
        ("TEMPLATE", "slide_timeline_ribbon", "리본 타임라인"),
        ("TEMPLATE", "slide_cover_neon", "네온 표지"),
        ("VARIANT", "HIGHLIGHT_NEON", "네온 하이라이트"),
    ]
    print(f"\n=== slide_kit v{__version__} — 신규 기능 ({len(items)}개) ===")
    for kind, name, desc in items:
        print(f"  [{kind:<10}] {name:<25} {desc}")
    print()


# ═══════════════════════════════════════════════════════════════════════
#  v4.1 "EDITORIAL DARK" — 레퍼런스 재분석 기반 재설계 (2026-04-17)
# ═══════════════════════════════════════════════════════════════════════
"""
레퍼런스 22장 실측 분석 결과 반영.

v4.0은 "게이밍 네온"이라는 오분석에 기반했지만, 실제 레퍼런스는
에디토리얼(MIT Tech Review / Stripe Press 계열) 다크 스타일이다.

핵심 원칙:
  1. 섹션 디바이더는 "여백 80% + 거대 영문"이 정답
  2. 콘텐츠 워크호스는 "포토 카드 3열 + 본문"
  3. 사이언 네온은 eyebrow 라벨 전용 (배경 사용 금지)
  4. 퍼플 브랜드는 쉐브론/악센트 바 등 최소 강조만
  5. 풀블리드 사진 + 하단 그라디언트 오버레이가 최고급 느낌의 원천
  6. 표는 다크 BG + 헤더 컬러 바 + 넉넉한 패딩
  7. 페이지당 색상 4개 이하

컴포넌트:
  - HEADLINE_STACK       eyebrow + pre + headline 표준 타이포 묶음
  - PHOTO_CARD_TRIO      3열 포토 카드 (워크호스)
  - STAT_ROW_HERO        거대 수치 3개 가로 배치
  - DATA_TABLE_DARK      다크 배경 표
  - PHOTO_FULL_OVERLAY   풀블리드 이미지 + 하단 그라디언트 + 캡션
  - RENDER_CAPTION       3D 렌더/평면도 + 떠있는 캡션
  - CIRCULAR_PHOTO_FLOW  원형 사진 타임라인
  - CREDENTIAL_STAGE     STAGE 배지 + 사진 + 본문 카드
  - CHEVRON_CONNECTOR    카드 사이 작은 » 연결자

슬라이드 템플릿:
  - slide_divider_hero   거의 빈 다크 + 80pt+ 영문
  - slide_hook_question  배경사진 + 거대 질문 + stat row
  - slide_summary_split  좌측 정보 + 우측 IP 영역
  - slide_cover_editorial IP풀블리드 + 타이틀
"""


# ───────────────────────────────────────────────────────────────
#  v4.1 / SZ 티어 정정 — 레퍼런스 실측 기반
# ───────────────────────────────────────────────────────────────

SZ.update({
    # 레퍼런스 실측 반영 (v4.0 예상치 정정)
    "eyebrow":      12,   # 섹션 마커 대문자 라벨 ("ACTION PLAN") — 원래 13
    "pre_headline": 18,   # 헤드라인 위 리드 문장
    "headline":     36,   # 슬라이드 메인 헤드라인
    "sub_headline": 22,   # 부 헤드라인
    "stat_big":     72,   # 3-열 stat row (160명/48H/50+)
    "stat_hero_v41": 96,  # 단일 거대 수치
    "section_hero": 100,  # 섹션 디바이더 영문 (ACTION PLAN)
    "eod":          120,  # E.O.D 최대 크기
    "body_reading": 14,   # 일반 본문 (촘촘하지 않은)
})


# ───────────────────────────────────────────────────────────────
#  v4.1 / editorial_dark 테마
# ───────────────────────────────────────────────────────────────

THEMES["editorial_dark"] = {
    "primary":   (95, 112, 252),    # 5F70FC 퍼플 — 악센트 바, 쉐브론
    "secondary": (102, 255, 255),   # 66FFFF 사이언 — eyebrow 전용
    "teal":      (105, 226, 255),   # 69E2FF
    "accent":    (255, 80, 80),     # 레드 — 크리티컬 강조만
    "dark":      (28, 31, 40),      # 1C1F28 딥 네이비 BG
    "light":     (242, 245, 250),   # F2F5FA 페이퍼 라이트
}


# ───────────────────────────────────────────────────────────────
#  v4.1 / 배경 스타일
# ───────────────────────────────────────────────────────────────

def bg_editorial_dark(s, variant="deep"):
    """에디토리얼 다크 배경 — 순수 다크 또는 살짝 그라디언트.

    레퍼런스처럼 "거의 플랫 + 미묘한 그라디언트"가 핵심.
    과한 그라디언트는 저급해보임.
    """
    if variant == "deep":
        # 순수 플랫 다크 (대부분의 콘텐츠 슬라이드)
        bg(s, tok("surface/darker"))
    elif variant == "subtle":
        # 미묘한 그라디언트 (디바이더/표지에만)
        gradient_bg(s, tok("surface/darker"), tok("surface/dark"))
    elif variant == "spotlight":
        # 상단 살짝 밝게 (히어로 슬라이드)
        gradient_bg(s, tok("surface/dark"), tok("surface/darker"))
    return s


# ───────────────────────────────────────────────────────────────
#  v4.1 / 1. HEADLINE_STACK — 표준 타이포 리듬
# ───────────────────────────────────────────────────────────────

def HEADLINE_STACK(s, *, eyebrow="", pre="", headline="", sub="",
                    x_in=None, y_in=None, w_in=None,
                    on_dark=True, align="left", eyebrow_color=None):
    """에디토리얼 다크의 표준 타이포 스택.

    레퍼런스 장표 거의 모든 콘텐츠 슬라이드의 상단 영역.

    Layout:
        [EYEBROW]  ← 사이언 대문자 12pt (섹션 마커)
        pre text  ← 헤드라인 위 가벼운 리드 18pt 회색
        HEADLINE  ← 36pt 흰색 볼드 (메인 메시지)
        subtitle  ← 22pt 회색 (부가)

    Args:
        eyebrow: 섹션 마커 (대문자 권장, 예: "ACTION PLAN")
        pre: 헤드라인 위 리드 문장
        headline: 메인 헤드라인
        sub: 서브
        on_dark: 다크 배경 여부
        w_in: None이면 CW
    """
    if x_in is None:
        x_in = ML_IN
    if y_in is None:
        y_in = round(float(SH / 914400) * 0.085, 3)  # 상단 8.5%
    if w_in is None:
        w_in = CW_IN
    eyebrow_c = eyebrow_color or tok("secondary" if on_dark else "brand/primary")
    headline_c = tok("text/on_dark") if on_dark else tok("text/on_light")
    pre_c = tok("text/muted")
    sub_c = tok("text/muted")

    al = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
          "right": PP_ALIGN.RIGHT}[align]

    y = y_in
    # 한글 너비 추정 (일반 sans 기준): 폰트 pt * 0.056 inch/char
    def _fits_one_line(text, pt, w):
        if not text:
            return True
        char_w = pt * 0.056  # 대략치
        return len(text) * char_w <= w - 0.1

    # Eyebrow
    if eyebrow:
        t = T(s, Inches(x_in), Inches(y), Inches(w_in), Inches(0.35),
              eyebrow.upper(), sz=SZ["eyebrow"], c=eyebrow_c, b=True,
              al=al, fn=FONT_W["bold"])
        try:
            set_char_spacing(t, 200)
        except Exception:
            pass
        y += 0.45

    # Pre-headline — 1줄 또는 2줄 자동 감지
    if pre:
        # 18pt → ~1.01 inch/char × len
        pre_lines = 1 if _fits_one_line(pre, SZ["pre_headline"], w_in) else 2
        pre_h = 0.45 * pre_lines
        T(s, Inches(x_in), Inches(y), Inches(w_in), Inches(pre_h),
          pre, sz=SZ["pre_headline"], c=pre_c,
          al=al, fn=FONT_W["regular"])
        y += pre_h + 0.1

    # Headline — 길이 기반 자동 크기 조절 + 공간 확보
    if headline:
        # 기본 36pt → 길면 32pt, 더 길면 28pt
        hl_sz = SZ["headline"]
        if len(headline) > 28:
            hl_sz = 32
        if len(headline) > 38:
            hl_sz = 28
        # 줄 수 예측 (컨테이너 폭 기준)
        hl_lines = 1 if _fits_one_line(headline, hl_sz, w_in) else 2
        hl_h = (hl_sz / 72) * 1.3 * hl_lines  # line height 1.3x
        T(s, Inches(x_in), Inches(y), Inches(w_in), Inches(hl_h + 0.1),
          headline, sz=hl_sz, c=headline_c, b=True,
          al=al, fn=FONT_W["bold"])
        y += hl_h + 0.2

    # Sub
    if sub:
        sub_lines = 1 if _fits_one_line(sub, SZ["sub_headline"], w_in) else 2
        sub_h = 0.4 * sub_lines + 0.1
        T(s, Inches(x_in), Inches(y), Inches(w_in), Inches(sub_h),
          sub, sz=SZ["sub_headline"], c=sub_c,
          al=al, fn=FONT_W["medium"])
        y += sub_h + 0.1

    return y  # 다음 요소의 시작 Y


# ───────────────────────────────────────────────────────────────
#  v4.1 / 2. PHOTO_CARD_TRIO — 메인 워크호스
# ───────────────────────────────────────────────────────────────

def PHOTO_CARD_TRIO(s, items, *, y_in=3.0, h_in=4.3,
                     on_dark=True, gap_in=0.25,
                     label_colors=None, img_ratio=None):
    """3열 포토 카드 — 레퍼런스 콘텐츠 슬라이드의 80%.

    각 카드 구조:
      [포토 영역 60%]          ← IMG_PH 또는 실제 이미지
      [컬러 라벨 바 2px]        ← 카테고리 구분용
      [영문 제목 (대문자)]       ← label (카드 타이틀)
      [한글 세부 설명]          ← body

    Args:
        items: [{"img": path?, "label": "GAME DEV",
                 "title": "160명의 개발자", "body": "진짜 게임 개발..."}, ...]
        label_colors: 각 라벨 바 색상 리스트. None이면 브랜드 퍼플 통일.
    """
    n = min(3, len(items))
    l_start = ML_IN
    w_total = CW_IN
    each_w = (w_total - gap_in * (n - 1)) / n

    if label_colors is None:
        label_colors = [tok("brand/primary")] * n

    # 이미지 비율 자동 결정: h_in에 따라 동적 (카드 작으면 텍스트 공간 확보)
    if img_ratio is None:
        if h_in < 3.0:
            img_ratio = 0.42   # 작은 카드 → 텍스트 우선
        elif h_in < 4.0:
            img_ratio = 0.48
        else:
            img_ratio = 0.52

    img_h = h_in * img_ratio

    for i, item in enumerate(items[:3]):
        x = l_start + i * (each_w + gap_in)

        # 포토 영역
        if item.get("img"):
            # 실제 이미지가 있으면 insert
            try:
                s.shapes.add_picture(item["img"], Inches(x), Inches(y_in),
                                      Inches(each_w), Inches(img_h))
            except Exception:
                # fallback: placeholder
                IMG_PH(s, Inches(x), Inches(y_in), Inches(each_w), Inches(img_h), label=item.get("label", "이미지"))
        else:
            IMG_PH(s, Inches(x), Inches(y_in), Inches(each_w), Inches(img_h), label=item.get("label", "이미지"))

        # 컬러 라벨 바 (포토와 텍스트 경계)
        bar_color = label_colors[i % len(label_colors)]
        R(s, Inches(x), Inches(y_in + img_h),
          Inches(each_w), Inches(0.04), bar_color)

        # 카드 내 가용 공간 (이미지 아래 전체)
        card_text_h = h_in - img_h - 0.12

        # 영문 라벨 (작게, 타이트한 간격)
        lab_h_fix = 0.28
        lab_y = y_in + img_h + 0.12
        if item.get("label"):
            T(s, Inches(x), Inches(lab_y),
              Inches(each_w), Inches(lab_h_fix),
              item["label"].upper(),
              sz=SZ["eyebrow"], c=bar_color, b=True,
              al=PP_ALIGN.LEFT, fn=FONT_W["bold"])
            lab_y += lab_h_fix + 0.02
            card_text_h -= lab_h_fix + 0.02

        # 3분할: 라벨 제외 나머지를 title:body = 40:60 비율
        #   작은 카드는 title을 더 작게, 큰 카드는 넉넉히
        title_slot = card_text_h * 0.42
        body_slot  = card_text_h - title_slot - 0.05

        # 한글 타이틀 — title_slot에 맞춰 폰트 자동 축소
        if item.get("title"):
            t_sz = SZ["sub_headline"]
            title_lines = item["title"].count("\n") + 1
            char_w = t_sz * 0.056
            per_line = max(1, int((each_w - 0.05) / char_w))
            longest = max(len(ln) for ln in item["title"].split("\n")) \
                        if item["title"] else 0
            if longest > per_line:
                title_lines += 1
            # 폰트 축소
            needed = (t_sz / 72) * 1.3 * title_lines
            while needed > title_slot and t_sz > 10:
                t_sz -= 1
                needed = (t_sz / 72) * 1.3 * title_lines
            T(s, Inches(x), Inches(lab_y),
              Inches(each_w), Inches(title_slot),
              item["title"],
              sz=t_sz, b=True,
              c=tok("text/on_dark") if on_dark else tok("text/on_light"),
              al=PP_ALIGN.LEFT, fn=FONT_W["bold"])
            lab_y += title_slot + 0.05

        # 본문 — body_slot 기반, 폰트 자동 축소
        if item.get("body"):
            body_h = min(body_slot, y_in + h_in - lab_y - 0.05)
            if body_h < 0.2:
                continue
            b_sz = SZ["body_reading"]
            body_chars = len(item["body"])
            per_line_b = max(1, int((each_w - 0.1) / (b_sz * 0.056)))
            body_lines = max(1, (body_chars // per_line_b) + 1)
            needed_h = (b_sz / 72) * 1.35 * body_lines
            while needed_h > body_h and b_sz > 8:
                b_sz -= 1
                per_line_b = max(1, int((each_w - 0.1) / (b_sz * 0.056)))
                body_lines = max(1, (body_chars // per_line_b) + 1)
                needed_h = (b_sz / 72) * 1.35 * body_lines
            T(s, Inches(x), Inches(lab_y),
              Inches(each_w), Inches(body_h),
              item["body"],
              sz=b_sz,
              c=tok("text/muted"),
              al=PP_ALIGN.LEFT, fn=FONT_W["regular"])


# ───────────────────────────────────────────────────────────────
#  v4.1 / 3. STAT_ROW_HERO — 거대 수치 3열
# ───────────────────────────────────────────────────────────────

def STAT_ROW_HERO(s, items, *, y_in=3.5, h_in=2.2,
                   on_dark=True, show_dividers=False, value_sz=None):
    """3개 거대 수치 가로 배치 (160명 | 48H | 50+).

    레퍼런스 Hook 슬라이드의 핵심 장치.

    items: [{"value": "160명", "label": "총 참가자",
              "desc": "대학(원)생 + 일반 창작자"}, ...]

    컨테이너 높이:
      - val_h = 폰트크기 × 1.25 / 72 (line-height 여유)
      - 전체 h_in 이상이면 폰트 자동 축소
    """
    n = len(items)
    l_start = ML_IN
    w_total = CW_IN
    each_w = w_total / n
    div_w = 0.01

    # 폰트 크기 자동 결정 — 실제 텍스트 폭 측정 기반
    v_sz = value_sz or SZ["stat_big"]  # 기본 72pt
    # 각 컬럼 안에 실제 들어가는지 측정 (여유 패딩 0.2")
    col_safe_w = each_w - 0.2
    # 가장 넓은 value를 기준으로 공통 폰트 축소
    max_needed = 0
    for it in items:
        val = str(it.get("value", ""))
        max_needed = max(max_needed, measure_text_width(val, v_sz, "black"))
    # 넘치면 비례 축소
    if max_needed > col_safe_w:
        scale = col_safe_w / max_needed
        v_sz = max(32, int(v_sz * scale))   # 최소 32pt (stat 느낌 유지)

    # 컨테이너 높이 = 폰트 인치 + 라벨 + 설명 + 여백
    val_h = (v_sz / 72) * 1.3    # 약 1.3" at 72pt
    label_h = 0.35
    desc_h = 0.35
    needed_h = val_h + 0.1 + label_h + 0.05 + desc_h
    # h_in이 부족하면 값 폰트 축소
    if needed_h > h_in:
        v_sz = int(v_sz * (h_in / needed_h) * 0.95)
        val_h = (v_sz / 72) * 1.3

    value_color = tok("text/on_dark") if on_dark else tok("text/on_light")
    label_color = tok("secondary")
    desc_color = tok("text/muted")

    for i, item in enumerate(items):
        x = l_start + i * each_w

        # 구분선
        if show_dividers and i > 0:
            R(s, Inches(x - div_w/2), Inches(y_in + 0.3),
              Inches(div_w), Inches(h_in - 0.6),
              tok("border/dark"))

        # 큰 수치
        val_box = s.shapes.add_textbox(Inches(x), Inches(y_in),
                                        Inches(each_w), Inches(val_h))
        tf = val_box.text_frame
        tf.margin_left = tf.margin_right = Inches(0.05)
        tf.margin_top = tf.margin_bottom = Inches(0)
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(item["value"])
        run.font.name = FONT_W["black"]
        run.font.size = Pt(v_sz)
        run.font.color.rgb = value_color
        run.font.bold = True

        # 라벨 (사이언 작게)
        lab_y = y_in + val_h + 0.1
        lab_box = s.shapes.add_textbox(Inches(x), Inches(lab_y),
                                        Inches(each_w), Inches(label_h))
        lp = lab_box.text_frame
        lp.margin_left = lp.margin_right = Inches(0.05)
        lp2 = lp.paragraphs[0]
        lp2.alignment = PP_ALIGN.CENTER
        lrun = lp2.add_run()
        lrun.text = item.get("label", "")
        lrun.font.name = FONT_W["semibold"]
        lrun.font.size = Pt(SZ["label"])
        lrun.font.color.rgb = label_color

        # 설명
        if item.get("desc"):
            d_box = s.shapes.add_textbox(Inches(x), Inches(lab_y + label_h + 0.02),
                                          Inches(each_w), Inches(desc_h))
            dtf = d_box.text_frame
            dtf.margin_left = dtf.margin_right = Inches(0.05)
            dp = dtf.paragraphs[0]
            dp.alignment = PP_ALIGN.CENTER
            dr = dp.add_run()
            dr.text = item["desc"]
            dr.font.name = FONT_W["regular"]
            dr.font.size = Pt(SZ["caption_sm"])
            dr.font.color.rgb = desc_color


# ───────────────────────────────────────────────────────────────
#  v4.1 / 4. DATA_TABLE_DARK — 다크 배경 표
# ───────────────────────────────────────────────────────────────

def DATA_TABLE_DARK(s, headers, rows, *, x_in=0.7, y_in=1.8,
                     w_in=11.95, row_h_in=0.4,
                     header_color=None, highlight_col=None):
    """다크 배경 전용 표 — 레퍼런스의 Time Table, Staff Plan 스타일.

    - 헤더: 브랜드 퍼플 배경 + 흰 글자
    - 바디: 투명 배경 + 서브틀 보더
    - 여유 있는 행 높이
    """
    header_color = header_color or tok("brand/primary")
    n_cols = len(headers)
    n_rows = len(rows)

    # python-pptx add_table
    cols_w = [w_in / n_cols] * n_cols
    tot_rows = 1 + n_rows
    total_h = row_h_in * tot_rows

    tbl_shape = s.shapes.add_table(tot_rows, n_cols,
                                     Inches(x_in), Inches(y_in),
                                     Inches(w_in), Inches(total_h))
    tbl = tbl_shape.table

    # 행 높이 설정
    for r in range(tot_rows):
        tbl.rows[r].height = Inches(row_h_in)

    # 헤더
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        tf = cell.text_frame
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.1)
        tf.margin_top = tf.margin_bottom = Inches(0.08)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(h)
        run.font.name = FONT_W["bold"]
        run.font.size = Pt(SZ["fine"])
        run.font.color.rgb = tok("text/on_dark")
        run.font.bold = True

    # 바디
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.fill.solid()
            if highlight_col is not None and c == highlight_col:
                cell.fill.fore_color.rgb = tok("surface/dark")
            else:
                cell.fill.fore_color.rgb = tok("surface/darker")

            tf = cell.text_frame
            tf.margin_left = Inches(0.15)
            tf.margin_right = Inches(0.1)
            tf.margin_top = tf.margin_bottom = Inches(0.08)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.name = FONT_W["regular"]
            run.font.size = Pt(SZ["body_reading"])
            run.font.color.rgb = tok("text/on_dark")

    return tbl_shape


# ───────────────────────────────────────────────────────────────
#  v4.1 / 5. PHOTO_FULL_OVERLAY — 풀블리드 사진 + 하단 그라디언트
# ───────────────────────────────────────────────────────────────

def PHOTO_FULL_OVERLAY(s, image_path=None, *, caption="", sub_caption="",
                        overlay_strength=65, overlay_color=None):
    """풀블리드 이미지 + 하단 그라디언트 오버레이 + 캡션.

    레퍼런스 실적 슬라이드(60p) 패턴.
    image_path 없으면 IMG_PH로 자리만 잡음.
    """
    # 이미지
    if image_path:
        try:
            s.shapes.add_picture(image_path, Inches(0), Inches(0),
                                  width=SW, height=SH)
        except Exception:
            IMG_PH(s, Inches(0), Inches(0), SW, SH, label=caption or "Photo")
    else:
        IMG_PH(s, Inches(0), Inches(0), SW, SH, label=caption or "Photo")

    # 하단 그라디언트 오버레이 (캔버스 비율 기반)
    overlay_color = overlay_color or tok("surface/darker")
    _sh = float(SH / 914400)
    overlay_h = _sh * 0.37
    overlay = R(s, Inches(0), Inches(_sh - overlay_h),
                 SW, Inches(overlay_h), overlay_color)
    gradient_shape(overlay, overlay_color, overlay_color,
                    angle=5400000)
    add_alpha(overlay, 100 - overlay_strength)

    # 캡션 (비율 기반 하단 배치)
    if caption:
        T(s, Inches(ML_IN), Inches(_sh * 0.813), Inches(CW_IN),
          Inches(_sh * 0.1),
          caption, sz=SZ["headline"], c=tok("text/on_dark"), b=True,
          fn=FONT_W["bold"], al=PP_ALIGN.LEFT)

    # 부캡션
    if sub_caption:
        T(s, Inches(ML_IN), Inches(_sh * 0.907), Inches(CW_IN),
          Inches(_sh * 0.06),
          sub_caption, sz=SZ["label"], c=tok("text/muted"),
          fn=FONT_W["regular"], al=PP_ALIGN.LEFT)


# ───────────────────────────────────────────────────────────────
#  v4.1 / 6. RENDER_CAPTION — 렌더/평면도 + 캡션
# ───────────────────────────────────────────────────────────────

def RENDER_CAPTION(s, image_path=None, *, title="", caption="",
                    accent_note="", image_area=(0.7, 1.5, 11.95, 4.5),
                    on_dark=True):
    """공간 렌더 or 평면도 + 하단 캡션 카드.

    레퍼런스: Developer Room(42p), Info Desk(40p), Space Design(38p).
    """
    l, t, w, h = image_area
    if image_path:
        try:
            s.shapes.add_picture(image_path, Inches(l), Inches(t),
                                  Inches(w), Inches(h))
        except Exception:
            IMG_PH(s, Inches(l), Inches(t), Inches(w), Inches(h), label=title or "Render")
    else:
        IMG_PH(s, Inches(l), Inches(t), Inches(w), Inches(h), label=title or "Render")

    # 캡션 영역
    cap_y = t + h + 0.3

    if title:
        T(s, Inches(l), Inches(cap_y), Inches(w), Inches(0.4),
          title, sz=SZ["sub_headline"],
          c=tok("text/on_dark") if on_dark else tok("text/on_light"),
          b=True, fn=FONT_W["bold"], al=PP_ALIGN.LEFT)

    if caption:
        T(s, Inches(l), Inches(cap_y + 0.5), Inches(w), Inches(0.4),
          caption, sz=SZ["body_reading"],
          c=tok("text/muted"),
          fn=FONT_W["regular"], al=PP_ALIGN.LEFT)

    if accent_note:
        # 우측 상단 작은 노트 (레퍼런스 "※ 넥슨 지하 1층 교실 이용")
        T(s, Inches(SW/914400 - 3.0), Inches(t + h + 0.3),
          Inches(2.8), Inches(0.3),
          "※ " + accent_note, sz=SZ["caption_sm"],
          c=tok("text/subtle"), al=PP_ALIGN.RIGHT)


# ───────────────────────────────────────────────────────────────
#  v4.1 / 7. CIRCULAR_PHOTO_FLOW — 원형 사진 타임라인
# ───────────────────────────────────────────────────────────────

def CIRCULAR_PHOTO_FLOW(s, items, *, y_in=3.2, circle_d=1.6,
                         arrow_color=None):
    """원형 사진 + 화살표 타임라인 — 레퍼런스 Event Flow(30p) 패턴.

    items: [{"img": path?, "stage": "STAGE 1", "title": "참가 접수",
              "time": "1월 19:00~20:00"}, ...]
    """
    arrow_color = arrow_color or tok("brand/primary")
    n = len(items)
    # 4개 기준 설계
    l_start = ML_IN
    w_total = CW_IN
    each_w = w_total / n

    for i, item in enumerate(items):
        cx = l_start + i * each_w + each_w / 2 - circle_d / 2

        # 원형 사진
        if item.get("img"):
            try:
                pic = s.shapes.add_picture(item["img"], Inches(cx),
                                            Inches(y_in),
                                            Inches(circle_d),
                                            Inches(circle_d))
                # 원형 크롭 (라운드 100%)
                # python-pptx는 도형 변환이 제한적. 원형 프레임 오버레이
                frame = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx),
                                            Inches(y_in), Inches(circle_d),
                                            Inches(circle_d))
                frame.fill.background()
                frame.line.color.rgb = arrow_color
                frame.line.width = Pt(2.5)
            except Exception:
                oval = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx),
                                           Inches(y_in), Inches(circle_d),
                                           Inches(circle_d))
                oval.fill.solid()
                oval.fill.fore_color.rgb = tok("surface/dark")
                oval.line.color.rgb = arrow_color
                oval.line.width = Pt(2.5)
        else:
            oval = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx),
                                       Inches(y_in), Inches(circle_d),
                                       Inches(circle_d))
            oval.fill.solid()
            oval.fill.fore_color.rgb = tok("surface/dark")
            oval.line.color.rgb = arrow_color
            oval.line.width = Pt(2.5)

        # 스테이지 라벨 (원 위)
        if item.get("stage"):
            T(s, Inches(l_start + i * each_w), Inches(y_in - 0.5),
              Inches(each_w), Inches(0.4),
              item["stage"], sz=SZ["eyebrow"],
              c=arrow_color, b=True, al=PP_ALIGN.CENTER,
              fn=FONT_W["bold"])

        # 타이틀 (원 아래)
        label_y = y_in + circle_d + 0.15
        if item.get("title"):
            T(s, Inches(l_start + i * each_w), Inches(label_y),
              Inches(each_w), Inches(0.4),
              item["title"], sz=SZ["sub_headline"],
              c=tok("text/on_dark"), b=True,
              al=PP_ALIGN.CENTER, fn=FONT_W["bold"])

        # 시간
        if item.get("time"):
            T(s, Inches(l_start + i * each_w), Inches(label_y + 0.45),
              Inches(each_w), Inches(0.3),
              item["time"], sz=SZ["caption_sm"],
              c=tok("text/muted"), al=PP_ALIGN.CENTER)

        # 화살표 (다음 요소 연결)
        if i < n - 1:
            arrow_x = l_start + (i + 1) * each_w - 0.3
            arrow_y = y_in + circle_d / 2 - 0.15
            CHEVRON_CONNECTOR(s, arrow_x, arrow_y, 0.5, 0.3,
                                color=arrow_color)


# ───────────────────────────────────────────────────────────────
#  v4.1 / 8. CHEVRON_CONNECTOR — 작은 » 연결자
# ───────────────────────────────────────────────────────────────

def CHEVRON_CONNECTOR(s, x_in, y_in, w_in=0.5, h_in=0.3, *, color=None):
    """카드 사이의 작은 쉐브론 연결자 (>> 느낌).

    레퍼런스 Check Point(5p), Communication(55p)에서 카드 그룹 사이.
    """
    color = color or tok("brand/primary")
    # 2개 겹치는 쉐브론으로 »
    ch1 = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x_in), Inches(y_in),
                              Inches(w_in * 0.5), Inches(h_in))
    ch1.fill.solid()
    ch1.fill.fore_color.rgb = color
    ch1.line.fill.background()
    ch2 = s.shapes.add_shape(MSO_SHAPE.CHEVRON,
                              Inches(x_in + w_in * 0.3), Inches(y_in),
                              Inches(w_in * 0.5), Inches(h_in))
    ch2.fill.solid()
    ch2.fill.fore_color.rgb = color
    ch2.line.fill.background()


# ───────────────────────────────────────────────────────────────
#  v4.1 / 9. CREDENTIAL_STAGE — START/GROW/EVOLVE 카드
# ───────────────────────────────────────────────────────────────

def CREDENTIAL_STAGE(s, items, *, y_in=2.6, h_in=3.8, stage_colors=None):
    """실적 단계 카드 — START 2016 / GROW / EVOLVE 2026 패턴.

    레퍼런스 NYPC 2016(7p) 스타일.

    items: [{"stage": "START, 2016", "title": "시작",
              "body": "...", "img": path?}, ...]
    """
    n = min(3, len(items))
    if stage_colors is None:
        stage_colors = [tok("brand/primary"), tok("brand/secondary"),
                         tok("secondary")]
    l_start = ML_IN
    w_total = CW_IN
    gap = 0.2
    each_w = (w_total - gap * (n - 1)) / n

    for i, item in enumerate(items[:3]):
        x = l_start + i * (each_w + gap)
        color = stage_colors[i % len(stage_colors)]

        # 사진 (상단 60%)
        img_h = h_in * 0.6
        if item.get("img"):
            try:
                s.shapes.add_picture(item["img"], Inches(x), Inches(y_in),
                                      Inches(each_w), Inches(img_h))
            except Exception:
                IMG_PH(s, Inches(x), Inches(y_in), Inches(each_w), Inches(img_h), label=item.get("stage", "이미지"))
        else:
            IMG_PH(s, Inches(x), Inches(y_in), Inches(each_w), Inches(img_h), label=item.get("stage", "이미지"))

        # 스테이지 배지 (포토 하단에 겹치게) — BADGE로 상하좌우 중앙정렬 보장
        badge_y = y_in + img_h - 0.3
        BADGE(s, x + 0.3, badge_y, 2.0, 0.5,
              item.get("stage", "").upper(),
              fill=color, sz_pt=SZ["label"])

        # 제목 (뱃지 아래, 적당한 gap)
        title_y = y_in + img_h + 0.3
        title_h_fix = 0.4
        T(s, Inches(x), Inches(title_y),
          Inches(each_w), Inches(title_h_fix),
          item.get("title", ""), sz=SZ["sub_headline"],
          c=tok("text/on_dark"), b=True,
          al=PP_ALIGN.LEFT, fn=FONT_W["bold"])

        # 본문 — 남은 공간에 맞춤
        body_y = title_y + title_h_fix + 0.05
        body_h = (y_in + h_in) - body_y - 0.05
        if body_h < 0.2:
            continue
        # 폰트 자동 축소 for overflow 방지
        b_sz = SZ["body_reading"]
        char_w = b_sz * 0.056
        per_line = max(1, int((each_w - 0.1) / char_w))
        body_lines = max(1, (len(item.get("body", "")) // per_line) + 1)
        needed = (b_sz / 72) * 1.35 * body_lines
        while needed > body_h and b_sz > 9:
            b_sz -= 1
            char_w = b_sz * 0.056
            per_line = max(1, int((each_w - 0.1) / char_w))
            body_lines = max(1, (len(item.get("body", "")) // per_line) + 1)
            needed = (b_sz / 72) * 1.35 * body_lines
        T(s, Inches(x), Inches(body_y),
          Inches(each_w), Inches(body_h),
          item.get("body", ""), sz=b_sz,
          c=tok("text/muted"),
          al=PP_ALIGN.LEFT, fn=FONT_W["regular"])


# ───────────────────────────────────────────────────────────────
#  v4.1 / 10. slide_divider_hero — 여백 디바이더
# ───────────────────────────────────────────────────────────────

def slide_divider_hero(prs, eng_title, kr_subtitle="", tagline="",
                         *, pg=None, variant="subtle"):
    """거대 영문 + 여백 80% 섹션 디바이더.

    레퍼런스 ACTION PLAN(15), MANAGEMENT(48), APPENDIX(58), E.O.D(65).

    Args:
        eng_title: 대형 영문 (대문자 권장)
        kr_subtitle: 아래 작은 한글 부제
        tagline: 더 아래 한 줄 사이언 태그
        variant: 배경 그라디언트 스타일
    """
    s = new_slide(prs)
    bg_editorial_dark(s, variant=variant)
    _sh = float(SH / 914400)

    # 영문 대형 — 수직 중앙에서 약간 위 (비율: 2.9/7.5 = 0.387)
    y_title = _sh * 0.387
    title_h = _sh * 0.20   # 1.5/7.5
    T(s, Inches(ML_IN), Inches(y_title), Inches(CW_IN), Inches(title_h),
      eng_title.upper(), sz=SZ["section_hero"],
      c=tok("text/on_dark"), b=True,
      al=PP_ALIGN.LEFT, fn=FONT_W["black"])

    # 한글 부제 (영문 바로 아래)
    if kr_subtitle:
        T(s, Inches(ML_IN + 0.05), Inches(y_title + title_h + 0.05),
          Inches(CW_IN), Inches(_sh * 0.07),
          kr_subtitle, sz=SZ["sub_headline"],
          c=tok("text/on_dark"), b=True,
          al=PP_ALIGN.LEFT, fn=FONT_W["bold"])

    # 사이언 태그라인
    if tagline:
        T(s, Inches(ML_IN + 0.05), Inches(y_title + title_h + _sh * 0.1),
          Inches(CW_IN), Inches(_sh * 0.06),
          tagline, sz=SZ["label"],
          c=tok("secondary"),
          al=PP_ALIGN.LEFT, fn=FONT_W["regular"])

    if pg is not None:
        # 우하단 작게
        T(s, Inches(SW/914400 - 0.85), Inches(SH/914400 - 0.4), Inches(0.7), Inches(0.3),
          str(pg), sz=SZ["caption_sm"],
          c=tok("text/muted"), al=PP_ALIGN.RIGHT)

    return s


# ───────────────────────────────────────────────────────────────
#  v4.1 / 11. slide_hook_question — 배경사진 + 거대 질문
# ───────────────────────────────────────────────────────────────

def slide_hook_question(prs, question, *, stats=None, bg_image=None,
                          sub="", pg=None):
    """HOOK 슬라이드 — 거대 질문 + 배경 사진 + 3개 통계.

    레퍼런스 WHY MAKERTHON?(2p) 패턴.

    Args:
        question: 거대 질문 (대문자 권장)
        stats: [{"value": "160명", "label": "총 참가자"}, ...] 3개
        bg_image: 배경 사진 경로
        sub: 질문 아래 작은 설명
    """
    s = new_slide(prs)

    # 배경 사진 (풀블리드)
    if bg_image:
        try:
            s.shapes.add_picture(bg_image, Inches(0), Inches(0),
                                  width=SW, height=SH)
        except Exception:
            bg_editorial_dark(s, "subtle")
    else:
        bg_editorial_dark(s, "subtle")

    # 다크 오버레이 (가독성)
    overlay = R(s, Inches(0), Inches(0), SW, SH, tok("surface/darker"))
    add_alpha(overlay, 35)  # 65% 불투명

    # 캔버스 크기 기반 비율 배치
    _sh = float(SH / 914400)
    _sw = float(SW / 914400)

    # 상단 eyebrow
    T(s, Inches(ML_IN), Inches(_sh * 0.07), Inches(CW_IN), Inches(0.3),
      "WHY ?", sz=SZ["eyebrow"],
      c=tok("secondary"), b=True,
      al=PP_ALIGN.LEFT, fn=FONT_W["bold"])

    # 거대 질문 — 캔버스 크기 + 텍스트 길이 기반 스케일
    base_q_sz = 72 if len(question) <= 15 else (56 if len(question) <= 25 else 44)
    q_sz = int(base_q_sz * (_sw / 13.333))   # 캔버스에 맞춤
    T(s, Inches(ML_IN), Inches(_sh * 0.18), Inches(CW_IN), Inches(_sh * 0.25),
      question.upper(), sz=q_sz,
      c=tok("text/on_dark"), b=True,
      al=PP_ALIGN.LEFT, fn=FONT_W["black"])

    # 설명
    if sub:
        T(s, Inches(ML_IN), Inches(_sh * 0.43), Inches(CW_IN), Inches(_sh * 0.15),
          sub, sz=SZ["body_reading"],
          c=tok("text/on_dark"),
          al=PP_ALIGN.LEFT, fn=FONT_W["regular"])

    # 3 stats — 남은 하단 공간에 맞춤 (페이지번호 영역 제외)
    if stats:
        stat_y = _sh * 0.60
        stat_h = _sh * 0.34   # 20~94% 사이
        STAT_ROW_HERO(s, stats, y_in=stat_y, h_in=stat_h,
                       on_dark=True, show_dividers=True)

    if pg is not None:
        T(s, Inches(SW/914400 - 0.85), Inches(SH/914400 - 0.4), Inches(0.7), Inches(0.3),
          str(pg), sz=SZ["caption_sm"],
          c=tok("text/muted"), al=PP_ALIGN.RIGHT)


# ───────────────────────────────────────────────────────────────
#  v4.1 / 12. slide_summary_split — 좌측 정보 + 우측 IP 영역
# ───────────────────────────────────────────────────────────────

def slide_summary_split(prs, *, title="", info_blocks=None,
                          ip_image=None, pg=None):
    """SUMMARY 슬라이드 — 좌측 키/밸류 + 우측 IP 일러스트.

    레퍼런스 SUMMARY(16p) 패턴.

    info_blocks: [("행사명", "2026 메이플스토리 월드 메커톤"), ...]
    """
    s = new_slide(prs)
    bg_editorial_dark(s, "deep")

    # 상단 eyebrow
    T(s, Inches(ML_IN), Inches(0.7), Inches(CW_IN), Inches(0.3),
      "SUMMARY", sz=SZ["eyebrow"],
      c=tok("secondary"), b=True,
      al=PP_ALIGN.LEFT, fn=FONT_W["bold"])

    # 좌측 정보 블록
    left_x = 0.7
    left_w = 6.5
    y = 1.4
    if info_blocks:
        for key, val in info_blocks:
            # 라벨
            T(s, Inches(left_x), Inches(y), Inches(left_w), Inches(0.35),
              key, sz=SZ["label"], c=tok("secondary"),
              b=True, fn=FONT_W["bold"])
            # 값
            T(s, Inches(left_x), Inches(y + 0.4), Inches(left_w),
              Inches(0.5),
              val, sz=SZ["sub_headline"], c=tok("text/on_dark"),
              b=True, fn=FONT_W["bold"])
            y += 1.05

    # 우측 IP 영역
    right_x = 7.5
    right_w = 5.5
    if ip_image:
        try:
            s.shapes.add_picture(ip_image, Inches(right_x), Inches(1.2),
                                  Inches(right_w), Inches(5.5))
        except Exception:
            IMG_PH(s, Inches(right_x), Inches(1.2), Inches(right_w), Inches(5.5), label="IP Visual")
    else:
        IMG_PH(s, Inches(right_x), Inches(1.2), Inches(right_w), Inches(5.5), label="IP Visual")

    if pg is not None:
        T(s, Inches(SW/914400 - 0.85), Inches(SH/914400 - 0.4), Inches(0.7), Inches(0.3),
          str(pg), sz=SZ["caption_sm"],
          c=tok("text/muted"), al=PP_ALIGN.RIGHT)


# ───────────────────────────────────────────────────────────────
#  v4.1 / 13. slide_cover_editorial — 에디토리얼 표지
# ───────────────────────────────────────────────────────────────

def slide_cover_editorial(prs, *, ip_image=None, title="",
                            subtitle="", client="", date=""):
    """표지 — IP 비주얼 풀블리드 + 하단 타이틀 블록.

    레퍼런스 Cover(1p) 패턴. IP 캐릭터/비주얼이 주인공.
    """
    s = new_slide(prs)
    has_ip_image = False

    # 배경
    if ip_image:
        try:
            s.shapes.add_picture(ip_image, Inches(0), Inches(0),
                                  width=SW, height=SH)
            has_ip_image = True
        except Exception:
            gradient_bg(s, tok("surface/dark"), tok("surface/darker"))
    else:
        gradient_bg(s, tok("surface/dark"), tok("surface/darker"))

    _sh = float(SH / 914400)
    if has_ip_image:
        # IP 이미지 있음 → 하단 오버레이 + 하단 타이틀 (비율 기반)
        overlay = R(s, Inches(0), Inches(_sh * 0.6), SW, Inches(_sh * 0.4),
                     tok("surface/darker"))
        add_alpha(overlay, 30)
        title_y = _sh * 0.680
        subtitle_y = _sh * 0.807
        y_bot = _sh * 0.907
    else:
        # IP 이미지 없음 → 중앙 센터드 레이아웃
        R(s, Inches(ML_IN), Inches(_sh * 0.27), Inches(0.8), Inches(0.04),
          tok("brand/primary"))
        title_y = _sh * 0.32     # 2.4 / 7.5 = 0.32
        subtitle_y = _sh * 0.48  # 3.6 / 7.5 = 0.48
        y_bot = _sh * 0.867      # 6.5 / 7.5

    # 타이틀 — 48pt 고정 + 2줄 예측해서 높이 확보
    title_sz = 48
    title_lines = 1
    if title:
        # 한글 char 48pt ≈ 0.48 inch. 폭 초과 시 2줄로
        char_w_est = 0.48
        if len(title) * char_w_est > CW_IN - 0.2:
            title_lines = 2
        title_h = (title_sz / 72) * 1.3 * title_lines + 0.1
        T(s, Inches(ML_IN), Inches(title_y), Inches(CW_IN),
          Inches(title_h),
          title, sz=title_sz,
          c=tok("text/on_dark"), b=True,
          al=PP_ALIGN.LEFT, fn=FONT_W["black"])

    # 서브타이틀 — 타이틀 하단 여유 간격 두고 배치
    if subtitle:
        # title이 2줄이면 subtitle을 실제 title 하단으로 밀기
        if title and title_lines == 2:
            effective_subtitle_y = title_y + (title_sz / 72) * 1.3 * 2 + 0.2
        else:
            effective_subtitle_y = subtitle_y
        T(s, Inches(ML_IN), Inches(effective_subtitle_y), Inches(CW_IN),
          Inches(0.45),
          "— " + subtitle + " —", sz=SZ["sub_headline"],
          c=tok("secondary"),
          al=PP_ALIGN.LEFT, fn=FONT_W["semibold"])

    # 날짜 + 클라이언트 (하단)
    if date:
        T(s, Inches(ML_IN), Inches(y_bot), Inches(CW_IN * 0.42), Inches(0.3),
          date, sz=SZ["label"],
          c=tok("text/on_dark"),
          al=PP_ALIGN.LEFT, fn=FONT_W["regular"])

    if client:
        T(s, Inches(ML_IN + CW_IN * 0.55), Inches(y_bot), Inches(CW_IN * 0.45), Inches(0.3),
          client, sz=SZ["label"],
          c=tok("text/on_dark"),
          al=PP_ALIGN.RIGHT, fn=FONT_W["bold"])


# ───────────────────────────────────────────────────────────────
#  v4.1 / 헬퍼
# ───────────────────────────────────────────────────────────────

# ═══════════════════════════════════════════════════════════════════════
#  VAETKI / EDITORIAL LIGHT 패턴 (2026-04-17 수주작 흡수)
# ═══════════════════════════════════════════════════════════════════════

def bg_pastel_gradient(s, c1=None, c2=None, c3=None):
    """3-stop 파스텔 그라디언트 배경 — VAETKI 시그니처.

    기본값: DEE2FB (연보라) → FFFFFF (화이트 중앙) → C2CAF8 (연블루)
    페이지 상하 극단이 컬러, 중앙이 화이트 → 부드러운 글로우 느낌

    Args:
        c1, c2, c3: 3-stop 색상 (RGBColor/hex). None이면 파스텔 기본.
    """
    if c1 is None:
        c1 = RGBColor(222, 226, 251)   # DEE2FB
    if c2 is None:
        c2 = RGBColor(255, 255, 255)   # FFFFFF
    if c3 is None:
        c3 = RGBColor(194, 202, 248)   # C2CAF8

    # 슬라이드 전체 덮는 rect + 3-stop 그라디언트
    try:
        from lxml import etree
    except ImportError:
        bg(s, c2)
        return s

    # 기존 slide <p:bg> 제거
    clear_slide_bg(s)
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    cSld = s._element.find(f'{{{ns_p}}}cSld')
    bg_el = etree.SubElement(cSld, f'{{{ns_p}}}bg')
    cSld.insert(0, bg_el)
    bgPr = etree.SubElement(bg_el, f'{{{ns_p}}}bgPr')
    gradFill = etree.SubElement(bgPr, f'{{{ns_a}}}gradFill',
                                 flip='none', rotWithShape='1')
    gsLst = etree.SubElement(gradFill, f'{{{ns_a}}}gsLst')
    for pos, c in [(0, c1), (50000, c2), (100000, c3)]:
        gs = etree.SubElement(gsLst, f'{{{ns_a}}}gs', pos=str(pos))
        hex_val = (f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"
                    if isinstance(c, (tuple, list, RGBColor))
                    else str(c).lstrip('#').upper())
        etree.SubElement(gs, f'{{{ns_a}}}srgbClr', val=hex_val)
    etree.SubElement(gradFill, f'{{{ns_a}}}lin',
                      ang="5400000", scaled='1')   # 수직 (상→하)
    return s


def gradient_headline(s, l_in, t_in, w_in, h_in, text, *,
                      c1=None, c2=None, sz_pt=36, align="center",
                      font_weight="black"):
    """2색 그라디언트 텍스트 헤드라인 — VAETKI 시그니처.

    "몰입감을 극대화하는 공간 설계" 같은 브랜드 헤드라인이
    보라블루→핑크로 흐르는 효과.

    Args:
        c1, c2: 그라디언트 시작/끝 색 (None이면 #6868F1 → #DD6495)
    """
    from pptx.enum.text import MSO_ANCHOR
    c1 = c1 or RGBColor(104, 104, 241)   # 6868F1
    c2 = c2 or RGBColor(221, 100, 149)   # DD6495

    tb = s.shapes.add_textbox(Inches(l_in), Inches(t_in),
                               Inches(w_in), Inches(h_in))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = {
        "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }[align]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_W.get(font_weight, FONT_W["black"])
    run.font.size = Pt(sz_pt)
    run.font.bold = True

    # 그라디언트 적용
    gradient_text(run, c1, c2)
    return tb


def PARALLELOGRAM_ZONE(s, l_in, t_in, w_in, h_in, text, *,
                        color=None, text_color=None, sz_pt=None,
                        alpha=100):
    """평행사변형 존 구분자 — VAETKI SPACE PLAN 스타일.

    "로비", "종합게임시연실" 같은 공간 카테고리 라벨로 쓰이는
    기울어진 큰 평행사변형.

    Args:
        color: 채우기 색 (파스텔 권장, 예: 연보라 B8B3ED / 연핑크 F2C4DA)
        alpha: 0~100 불투명도 (기본 100 = 불투명)
    """
    from pptx.enum.text import MSO_ANCHOR
    color = color or tok("brand/primary")
    text_color = text_color or tok("text/on_light")
    sz_pt = sz_pt or SZ["sub_headline"]

    shape = s.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
                                Inches(l_in), Inches(t_in),
                                Inches(w_in), Inches(h_in))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.15   # 살짝 기울어짐 (VAETKI 9p 스타일)
    except Exception:
        pass
    if alpha < 100:
        add_alpha(shape, 100 - alpha)

    tf = shape.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_W["medium"]
    run.font.size = Pt(sz_pt)
    run.font.color.rgb = text_color
    return shape


def slide_divider_light(prs, eng_title, *, pg=None,
                         gradient_bottom=True):
    """라이트 섹션 디바이더 — VAETKI "SPACE PLANNING" 패턴.

    - 좌상단 검정 영문 대형 텍스트 (여백 90%)
    - 하단 파스텔 그라디언트 (옵션)
    - 전체적으로 라이트 톤

    레퍼런스 slide 8 "SPACE PLANNING" 스타일.
    """
    s = new_slide(prs)
    if gradient_bottom:
        bg_pastel_gradient(s)
    else:
        bg(s, tok("surface/base") if on_light_mode() else RGBColor(255, 255, 255))

    # 좌상단 대형 검정 영문 (수직 중앙)
    _sh = float(SH / 914400)
    y_title = _sh * 0.42
    title_h = _sh * 0.15
    T(s, Inches(ML_IN), Inches(y_title), Inches(CW_IN), Inches(title_h),
      eng_title.upper(), sz=48, c=RGBColor(0, 0, 0), b=True,
      al=PP_ALIGN.LEFT, fn=FONT_W["black"])

    if pg is not None:
        T(s, Inches(float(SW/914400) - 0.85), Inches(_sh - 0.4),
          Inches(0.7), Inches(0.3),
          str(pg), sz=SZ["caption_sm"],
          c=RGBColor(100, 100, 100), al=PP_ALIGN.RIGHT)
    return s


def on_light_mode():
    """현재 적용된 팔레트가 라이트 모드인지 판정."""
    try:
        bg_c = tok("bg")
        luma = int(bg_c[0]) * 0.299 + int(bg_c[1]) * 0.587 + int(bg_c[2]) * 0.114
        return luma >= 180
    except Exception:
        return False


def PAGE_HEADER_LIGHT(s, *, page_title="", pre="", headline="",
                       gradient_headline_text=False,
                       y_title=None, y_center_start=None,
                       on_dark=None):
    """VAETKI 스타일 헤더 — 좌상단 라벨 + 중앙정렬 헤드라인.

    - 좌상단 페이지 라벨 ("SPACE PLAN", "TIME TABLE" 등)
    - 상단 구분선
    - 중앙 정렬 pre (뮤티드)
    - 중앙 정렬 headline (옵션: 그라디언트 텍스트)

    Args:
        gradient_headline_text: True면 headline에 브랜드 그라디언트 적용
        on_dark: None=자동감지 / True=다크 배경용 / False=라이트 배경용
    """
    _sh = float(SH / 914400)
    if y_title is None:
        y_title = _sh * 0.07
    if y_center_start is None:
        y_center_start = _sh * 0.18
    if on_dark is None:
        on_dark = _detect_dark_bg(s)

    # 텍스트 색상 팔레트 (다크/라이트 자동 스위치)
    if on_dark:
        page_title_color = tok("text/on_dark")   # 다크 배경 → 밝은 텍스트
        divider_color = RGBColor(80, 70, 110)
        pre_color = tok("text/muted")
        headline_color = tok("text/on_dark")
    else:
        page_title_color = RGBColor(0, 0, 0)
        divider_color = RGBColor(200, 200, 205)
        pre_color = RGBColor(100, 100, 115)
        headline_color = RGBColor(10, 10, 20)

    # 1. 좌상단 페이지 라벨 + 상단 구분선
    if page_title:
        T(s, Inches(ML_IN), Inches(y_title),
          Inches(CW_IN * 0.5), Inches(0.35),
          page_title.upper(), sz=SZ["sub_headline"], b=True,
          c=page_title_color, fn=FONT_W["bold"],
          al=PP_ALIGN.LEFT)
        R(s, Inches(ML_IN), Inches(y_title + 0.5),
          Inches(CW_IN), Inches(0.005), divider_color)

    y = y_center_start

    # 2. 중앙 정렬 pre
    if pre:
        pre_h = (SZ["pre_headline"] / 72) * 1.4
        T(s, Inches(ML_IN), Inches(y), Inches(CW_IN),
          Inches(pre_h + 0.1),
          pre, sz=SZ["pre_headline"],
          c=pre_color, fn=FONT_W["regular"],
          al=PP_ALIGN.CENTER)
        y += pre_h + 0.1

    # 3. 중앙 정렬 headline
    if headline:
        hl_sz = 36
        char_w = hl_sz * 0.056
        while len(headline) * char_w > CW_IN * 0.95 and hl_sz > 20:
            hl_sz -= 1
            char_w = hl_sz * 0.056
        hl_h = (hl_sz / 72) * 1.4
        if gradient_headline_text:
            gradient_headline(s, ML_IN, y, CW_IN, hl_h + 0.1,
                               headline, sz_pt=hl_sz, align="center",
                               font_weight="bold")
        else:
            T(s, Inches(ML_IN), Inches(y), Inches(CW_IN),
              Inches(hl_h + 0.1),
              headline, sz=hl_sz, b=True,
              c=headline_color, fn=FONT_W["bold"],
              al=PP_ALIGN.CENTER)
        y += hl_h + 0.15

    return y


def PAGE_HEADER(s, *, page_title="", pre="", headline="",
                  y_title=None, y_center_start=None,
                  on_dark=True, page_title_color=None):
    """레퍼런스 공통 헤더 패턴 — 좌상단 페이지제목 + 중앙정렬 pre/headline.

    Layout (NYPC 2016, MANAGEMENT, RECRUITMENT 등 레퍼런스 전반 패턴):

        ┌─────────────────────────────────────────────┐
        │  PAGE TITLE ← 좌상단 작게                     │
        │                                              │
        │         작은 리드 문장 (중앙 정렬)             │  ← pre
        │      큰 헤드라인 (중앙 정렬, 볼드)              │  ← headline
        │                                              │

    Args:
        page_title: 좌상단 페이지 제목 (예: "NYPC 2016")
        pre: 중앙 정렬 리드 문장 (설명)
        headline: 중앙 정렬 메인 헤드라인 (볼드)
        y_title: 페이지 제목 Y (기본 SH*0.08)
        y_center_start: pre 시작 Y (기본 SH*0.15)

    Returns:
        다음 요소 시작 Y (헤드라인 하단)
    """
    _sh = float(SH / 914400)
    if y_title is None:
        y_title = _sh * 0.08
    if y_center_start is None:
        y_center_start = _sh * 0.15

    text_color = tok("text/on_dark") if on_dark else tok("text/on_light")
    page_title_color = page_title_color or text_color
    pre_color = tok("text/muted") if on_dark else tok("text/subtle")

    # 1. 좌상단 페이지 제목 (작고 좌정렬)
    if page_title:
        T(s, Inches(ML_IN), Inches(y_title),
          Inches(CW_IN * 0.5), Inches(0.35),
          page_title, sz=SZ["sub_headline"], b=True,
          c=page_title_color, fn=FONT_W["bold"],
          al=PP_ALIGN.LEFT)

    # 2. 중앙 정렬 pre (작은 리드)
    y = y_center_start
    if pre:
        pre_h = (SZ["pre_headline"] / 72) * 1.4
        T(s, Inches(ML_IN), Inches(y), Inches(CW_IN),
          Inches(pre_h + 0.1),
          pre, sz=SZ["pre_headline"],
          c=pre_color, fn=FONT_W["regular"],
          al=PP_ALIGN.CENTER)
        y += pre_h + 0.1

    # 3. 중앙 정렬 headline (큰 볼드) — 길이 기반 자동 폰트
    if headline:
        hl_sz = SZ["headline"]   # 36pt (scale된 값)
        # 캔버스 폭 대비 폭 체크 (여유 있게 0.9 마진)
        char_w = hl_sz * 0.056
        while len(headline) * char_w > CW_IN * 0.95 and hl_sz > 20:
            hl_sz -= 1
            char_w = hl_sz * 0.056
        hl_h = (hl_sz / 72) * 1.4
        T(s, Inches(ML_IN), Inches(y), Inches(CW_IN),
          Inches(hl_h + 0.1),
          headline, sz=hl_sz, b=True,
          c=text_color, fn=FONT_W["bold"],
          al=PP_ALIGN.CENTER)
        y += hl_h + 0.15

    return y


def BADGE(s, l_in, t_in, w_in, h_in, text, *,
           fill=None, text_color=None, sz_pt=None, bold=True,
           font_weight="bold"):
    """컬러 배경 텍스트 뱃지 — 상하좌우 완전 중앙 정렬 보장.

    도형(R) + 텍스트박스(T) 조합으로 발생하는 중앙정렬 틀어짐 방지.
    도형 자체의 text_frame에 텍스트 삽입 + vertical_anchor=MIDDLE.
    """
    from pptx.enum.text import MSO_ANCHOR
    fill = fill or tok("brand/primary")
    text_color = text_color or tok("text/on_dark")
    sz_pt = sz_pt or SZ["label"]

    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(l_in), Inches(t_in),
                                Inches(w_in), Inches(h_in))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()

    tf = shape.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_W.get(font_weight, FONT_W["bold"])
    run.font.size = Pt(sz_pt)
    run.font.color.rgb = text_color
    run.font.bold = bold
    return shape


# ═══════════════════════════════════════════════════════════════════════
#  텍스트 폭 측정 + 덱 레이아웃 검증기
# ═══════════════════════════════════════════════════════════════════════

def measure_text_width(text, sz_pt, weight="regular"):
    """텍스트 폭을 인치로 추정 (CJK 0.9 EM, Latin 0.5 EM, 볼드 +8%).

    정확도: ±15% (폰트별 차이 고려). 오버플로우 1차 필터링 용도.
    """
    em = sz_pt / 72   # 1 EM 인치
    bold_f = 1.08 if weight in ("bold", "black", "semibold") else 1.0
    total = 0.0
    for ch in text:
        code = ord(ch)
        if code > 0x2E80:            # CJK (한글·한자·일본어)
            total += em * 0.90 * bold_f
        elif ch in " \t":            # 공백
            total += em * 0.28
        elif ch.isdigit():           # 숫자
            total += em * 0.52 * bold_f
        elif ch.isupper():           # 대문자 영문
            total += em * 0.60 * bold_f
        elif ch.isalpha():           # 소문자 영문
            total += em * 0.42 * bold_f
        else:                        # 기호·특수문자
            total += em * 0.32 * bold_f
    return total


def fit_font_to_width(text, sz_pt, max_w_in, weight="regular",
                       min_pt=8, pad_in=0.1):
    """max_w_in 너비에 맞게 폰트 크기 자동 축소.

    Returns: 조정된 pt.
    """
    sz = sz_pt
    while sz > min_pt and measure_text_width(text, sz, weight) > max_w_in - pad_in:
        sz -= 1
    return sz


def validate_deck(prs, *, verbose=True, check_overlaps=True,
                    safe_margin_in=0.05):
    """덱 전체 레이아웃 검증.

    검사 항목:
    - 슬라이드 경계 초과 (우측 / 하단 / 좌측 음수 / 상단 음수)
    - 텍스트 박스끼리 겹침 (check_overlaps=True 시)

    Args:
        safe_margin_in: 이 정도 초과는 허용 (0.05" = 약 1mm)

    Returns:
        [{"slide": N, "type": str, "shape": str, "detail": ...}, ...]
    """
    sw_in = prs.slide_width / 914400
    sh_in = prs.slide_height / 914400
    issues = []

    for idx, slide in enumerate(prs.slides, 1):
        boxes = []   # (l, t, r, b, shape_name) — 겹침 체크용
        for sp in slide.shapes:
            try:
                if sp.left is None or sp.top is None:
                    continue
                x = sp.left / 914400
                y = sp.top / 914400
                w = (sp.width / 914400) if sp.width else 0
                h = (sp.height / 914400) if sp.height else 0
                name = sp.name or sp.shape_type.name if hasattr(sp, "shape_type") else "?"

                # 경계 초과
                if x < -safe_margin_in:
                    issues.append({"slide": idx, "type": "overflow_left",
                                     "shape": name, "amount": -x})
                if y < -safe_margin_in:
                    issues.append({"slide": idx, "type": "overflow_top",
                                     "shape": name, "amount": -y})
                if x + w > sw_in + safe_margin_in:
                    issues.append({"slide": idx, "type": "overflow_right",
                                     "shape": name, "amount": x + w - sw_in,
                                     "at": (x, y, w, h)})
                if y + h > sh_in + safe_margin_in:
                    issues.append({"slide": idx, "type": "overflow_bottom",
                                     "shape": name, "amount": y + h - sh_in,
                                     "at": (x, y, w, h)})

                # 겹침 검사 대상 (텍스트가 있는 도형만)
                if check_overlaps and hasattr(sp, "text_frame") and sp.text_frame.text:
                    boxes.append((x, y, x + w, y + h, name, sp.text_frame.text[:20]))
            except Exception:
                continue

        # 텍스트 박스 겹침 (양쪽 모두 텍스트 있을 때만)
        if check_overlaps:
            for i in range(len(boxes)):
                for j in range(i + 1, len(boxes)):
                    a, b = boxes[i], boxes[j]
                    # 교차 영역 계산
                    ix1 = max(a[0], b[0])
                    iy1 = max(a[1], b[1])
                    ix2 = min(a[2], b[2])
                    iy2 = min(a[3], b[3])
                    if ix2 - ix1 > 0.1 and iy2 - iy1 > 0.1:
                        # 의미있는 겹침 (0.1" 이상)
                        # 단, 포함 관계(작은게 큰거 안에)는 무시 (카드 내 텍스트 정상 케이스)
                        a_area = (a[2]-a[0]) * (a[3]-a[1])
                        b_area = (b[2]-b[0]) * (b[3]-b[1])
                        smaller = min(a_area, b_area)
                        overlap_area = (ix2-ix1) * (iy2-iy1)
                        if overlap_area / smaller < 0.90:   # 90% 이상 포함 아니면 이슈
                            issues.append({
                                "slide": idx, "type": "text_overlap",
                                "shape": f"{a[4]} x {b[4]}",
                                "texts": (a[5], b[5]),
                                "overlap": (ix2 - ix1, iy2 - iy1),
                            })

    if verbose:
        print_validation_report(issues, sw_in, sh_in)
    return issues


def print_validation_report(issues, sw_in=None, sh_in=None):
    """검증 결과 콘솔 출력."""
    if not issues:
        print(f"\n[OK] 레이아웃 이슈 없음\n")
        return
    print(f"\n{'='*60}")
    print(f"레이아웃 검증 리포트 — {len(issues)}건")
    if sw_in and sh_in:
        print(f"캔버스: {sw_in:.2f}\" × {sh_in:.2f}\"")
    print(f"{'='*60}")
    by_slide = {}
    for i in issues:
        by_slide.setdefault(i["slide"], []).append(i)
    for slide_num in sorted(by_slide.keys()):
        items = by_slide[slide_num]
        print(f"\nSlide {slide_num} — {len(items)}건")
        for it in items[:5]:   # 슬라이드당 최대 5개
            t = it["type"]
            if t == "overflow_right":
                print(f"  [우측 초과 {it['amount']:.2f}\"] {it['shape'][:30]}")
            elif t == "overflow_bottom":
                print(f"  [하단 초과 {it['amount']:.2f}\"] {it['shape'][:30]}")
            elif t == "overflow_left":
                print(f"  [좌측 초과 {it['amount']:.2f}\"] {it['shape'][:30]}")
            elif t == "overflow_top":
                print(f"  [상단 초과 {it['amount']:.2f}\"] {it['shape'][:30]}")
            elif t == "text_overlap":
                ow, oh = it["overlap"]
                ta, tb = it["texts"]
                print(f"  [텍스트 겹침 {ow:.2f}x{oh:.2f}\"]")
                print(f"    '{ta}' ↔ '{tb}'")
        if len(items) > 5:
            print(f"  ... +{len(items) - 5}건")
    print(f"\n{'='*60}\n")


def auto_fix_overflow(prs, *, shrink_min_pt=10):
    """검증 후 overflow 발견 시 자동 수정 (텍스트 폰트 축소).

    현재 구현: text box의 폰트를 비례 축소하여 컨테이너 안에 맞춤.
    Returns: (fixed_count, remaining_issues)
    """
    issues = validate_deck(prs, verbose=False, check_overlaps=False)
    overflow_issues = [i for i in issues
                        if i["type"] in ("overflow_right", "overflow_bottom")]
    fixed = 0

    for issue in overflow_issues:
        slide = prs.slides[issue["slide"] - 1]
        target_name = issue["shape"]
        for sp in slide.shapes:
            if (sp.name or "") != target_name:
                continue
            if not hasattr(sp, "text_frame") or not sp.text_frame.text:
                continue
            # 현재 폰트 사이즈 찾기 → 축소
            for para in sp.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        cur = run.font.size.pt
                        new = max(shrink_min_pt, int(cur * 0.85))
                        if new < cur:
                            run.font.size = Pt(new)
                            fixed += 1
            break

    remaining = validate_deck(prs, verbose=False)
    return fixed, remaining


def list_v41_components():
    """v4.1 에디토리얼 다크 컴포넌트 목록."""
    items = [
        ("COMPONENT",  "HEADLINE_STACK",      "eyebrow+pre+headline+sub 표준 타이포"),
        ("COMPONENT",  "PHOTO_CARD_TRIO",     "3열 포토 카드 (메인 워크호스)"),
        ("COMPONENT",  "STAT_ROW_HERO",       "3 거대 수치 가로 배치"),
        ("COMPONENT",  "DATA_TABLE_DARK",     "다크 배경 표"),
        ("COMPONENT",  "PHOTO_FULL_OVERLAY",  "풀블리드 사진 + 하단 오버레이"),
        ("COMPONENT",  "RENDER_CAPTION",      "3D 렌더 + 캡션"),
        ("COMPONENT",  "CIRCULAR_PHOTO_FLOW", "원형 사진 타임라인"),
        ("COMPONENT",  "CHEVRON_CONNECTOR",   "작은 » 연결자"),
        ("COMPONENT",  "CREDENTIAL_STAGE",    "STAGE 배지 카드"),
        ("BG",         "bg_editorial_dark",   "에디토리얼 다크 배경"),
        ("TEMPLATE",   "slide_divider_hero",  "거대 영문 디바이더"),
        ("TEMPLATE",   "slide_hook_question", "HOOK 질문 슬라이드"),
        ("TEMPLATE",   "slide_summary_split", "SUMMARY 분할 슬라이드"),
        ("TEMPLATE",   "slide_cover_editorial","IP 표지"),
        ("THEME",      "editorial_dark",      "에디토리얼 다크 (공식 권장)"),
    ]
    print(f"\n=== slide_kit v{__version__} / EDITORIAL DARK ({len(items)}) ===")
    for kind, name, desc in items:
        print(f"  [{kind:<10}] {name:<26} {desc}")
    print()


# 별칭
_p = new_slide
_set_char_spacing = set_char_spacing
