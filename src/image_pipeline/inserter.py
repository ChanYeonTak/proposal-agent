"""
PPTX 이미지 삽입기

IMG_PH 도형을 실제 이미지로 교체합니다.
"""

from __future__ import annotations

from pathlib import Path
from typing import Dict

from ..utils.logger import get_logger

logger = get_logger("image_inserter")


class ImageInserter:
    """
    PPTX에 이미지 삽입

    IMG_PH 도형의 위치/크기를 유지하면서 실제 이미지로 교체합니다.
    """

    @staticmethod
    def insert_images(
        pptx_path: Path,
        image_map: Dict[str, Path],
        output_path: Optional[Path] = None,
    ) -> Path:
        """
        PPTX 파일의 IMG_PH를 실제 이미지로 교체

        Args:
            pptx_path: 원본 PPTX 경로
            image_map: {placeholder_text: image_path} 매핑
            output_path: 출력 경로 (None이면 원본에 _with_images 접미사)

        Returns:
            출력 파일 경로
        """
        if not image_map:
            logger.info("교체할 이미지가 없음")
            return pptx_path

        try:
            from pptx import Presentation
            from pptx.util import Inches, Emu

            prs = Presentation(str(pptx_path))
            replaced = 0

            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue

                    text = shape.text_frame.text.strip()

                    # IMG_PH 매칭
                    matched_path = None
                    for placeholder, img_path in image_map.items():
                        if placeholder in text or text in placeholder:
                            matched_path = img_path
                            break

                    if matched_path and matched_path.exists():
                        # 도형 위치/크기 저장
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height

                        # 기존 도형 삭제
                        sp = shape._element
                        sp.getparent().remove(sp)

                        # 이미지 삽입
                        slide.shapes.add_picture(
                            str(matched_path),
                            left, top, width, height,
                        )
                        replaced += 1

            out_path = output_path or pptx_path.with_stem(pptx_path.stem + "_with_images")
            prs.save(str(out_path))

            logger.info(f"이미지 삽입 완료: {replaced}개 교체 -> {out_path}")
            return out_path

        except ImportError:
            logger.error("python-pptx 미설치")
            return pptx_path
        except Exception as e:
            logger.error(f"이미지 삽입 실패: {e}")
            return pptx_path
