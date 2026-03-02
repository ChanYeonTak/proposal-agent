"""
디자인 브릿지 (v7.1) — Gamma MCP ↔ DesignAgent 연동 + API 폴링 자동 다운로드

Gamma MCP 도구(get_themes, generate)와 Python 코드 사이의 브릿지.
Claude Code가 Gamma MCP를 호출한 결과를 이 모듈이 가공하여
DesignAgent에 전달합니다.

워크플로우 (v7.1 — API 폴링 자동 다운로드):
    1. bridge.run_gamma_pipeline(pptx, brief) → GammaPipelineResult
       (텍스트 추출 + 파라미터 빌드 + 중간 파일 저장)
    2. Claude Code: Gamma MCP generate(**result.params) 호출
    3. bridge.handle_gamma_response(response, result) → 다운로드 전략 반환
    4. bridge.poll_and_download_gamma(generation_id, result) → PPTX 자동 다운로드
       (Gamma Public API 폴링 → exportUrl 획득 → HTTP 다운로드 → 프로젝트 저장)

    ※ GAMMA_API_KEY 미설정 시 fallback:
    4-alt. Claude Code: 브라우저로 gammaUrl 이동 + 수동 PPTX 내보내기
    5-alt. bridge.save_gamma_pptx(downloaded_path) → 프로젝트 디렉토리에 저장

    단순 호출 (기존 호환):
    1. Claude Code: Gamma MCP get_themes() 호출
    2. design_bridge.process_gamma_themes(결과) → ThemeRecommendation[]
    3. DesignAgent.merge(tt_brief, gamma_recs) → MergedDesignBrief
    4. DesignAgent.register_to_slide_kit(brief) → slide_kit 테마 적용

v7.1 변경:
    - poll_and_download_gamma() — Gamma Public API 폴링 + exportUrl PPTX 자동 다운로드
    - _download_export_url() — exportUrl → HTTP 다운로드 → 프로젝트 저장
    - handle_gamma_response() — api_poll 다운로드 전략 추가 (GAMMA_API_KEY 보유 시)
    - 모듈 워크플로우 설명 v7.1로 업데이트

v7.0 변경:
    - GammaPipelineResult 데이터 클래스 추가
    - run_gamma_pipeline() — 텍스트 추출 ~ 파라미터 빌드 원스텝 실행
    - handle_gamma_response() — generate() 응답 처리 + 다운로드 지침
    - save_gamma_pptx() — 다운로드된 PPTX를 프로젝트에 저장
    - get_gamma_export_url() — gammaUrl → 다운로드 URL 변환

v6.0:
    - extract_user_edits() 감지 고도화 (14개 카테고리)
    - 도형 매칭 알고리즘 (이름 + 위치 기반)
    - 세부 속성 비교: 지오메트리, 채우기, 선/테두리, 타이포그래피,
      단락 서식, 그림자/효과, 이미지 교체, 도형 유형 변경
    - 글로벌 디자인 통계 (최다 사용 색상, 폰트 등) 자동 추출
    - build_gamma_params() additionalInstructions 지원
"""

from __future__ import annotations

import hashlib
import json
import os
import shutil
import time
import urllib.request
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Set, Tuple
from xml.etree import ElementTree as ET

from ..utils.logger import get_logger

logger = get_logger("design_bridge")

# ═══════════════════════════════════════════════════════════════
#  EMU 변환 상수
# ═══════════════════════════════════════════════════════════════
_EMU_PER_INCH = 914400
_EMU_PER_PT = 12700


class DesignBridgeResult:
    """디자인 브릿지 결과"""

    def __init__(
        self,
        success: bool = False,
        provider: str = "",
        output_url: str = "",
        output_path: Optional[Path] = None,
        theme_id: Optional[str] = None,
        error: str = "",
    ):
        self.success = success
        self.provider = provider
        self.output_url = output_url
        self.output_path = output_path
        self.theme_id = theme_id
        self.error = error


class GammaPipelineResult:
    """
    Gamma 파이프라인 전체 상태 추적 객체.

    run_gamma_pipeline()에서 생성 → handle_gamma_response()에서 업데이트
    → save_gamma_pptx()에서 최종 완료.

    Claude Code가 각 단계별 상태를 확인하고
    다음 액션을 결정하는 데 사용합니다.
    """

    def __init__(self) -> None:
        # ── 입력 ──
        self.source_pptx: str = ""              # 원본 slide_kit PPTX 경로
        self.input_text: str = ""               # Gamma에 전송할 텍스트
        self.input_text_path: str = ""          # 저장된 input_text 파일 경로
        self.params: Dict[str, Any] = {}        # Gamma generate() 파라미터
        self.params_path: str = ""              # 저장된 params JSON 경로

        # ── Gamma 응답 ──
        self.gamma_url: str = ""                # Gamma 프레젠테이션 URL
        self.generation_id: str = ""            # Gamma generationId
        self.export_format: str = "pptx"        # 요청한 내보내기 형식
        self.status: str = "not_started"        # not_started / prepared / sent / completed / downloaded / error

        # ── 다운로드 ──
        self.export_url: str = ""               # PPTX 직접 다운로드 URL (있으면)
        self.downloaded_path: str = ""          # 다운로드된 파일의 로컬 경로
        self.output_path: str = ""              # 최종 저장 경로 (project_dir 내)

        # ── 메타 ──
        self.theme_id: str = ""                 # 사용된 Gamma 테마 ID
        self.theme_name: str = ""               # 사용된 Gamma 테마명
        self.project_name: str = ""             # 프로젝트명
        self.created_at: str = ""               # 생성 시각
        self.error: str = ""                    # 에러 메시지

    def to_dict(self) -> Dict[str, Any]:
        """JSON 직렬화 가능한 딕셔너리 반환."""
        return {
            "source_pptx": self.source_pptx,
            "input_text_chars": len(self.input_text),
            "input_text_path": self.input_text_path,
            "params_path": self.params_path,
            "gamma_url": self.gamma_url,
            "generation_id": self.generation_id,
            "export_format": self.export_format,
            "status": self.status,
            "export_url": self.export_url,
            "downloaded_path": self.downloaded_path,
            "output_path": self.output_path,
            "theme_id": self.theme_id,
            "theme_name": self.theme_name,
            "project_name": self.project_name,
            "created_at": self.created_at,
            "error": self.error,
        }

    def __repr__(self) -> str:
        return (
            f"GammaPipelineResult(status={self.status!r}, "
            f"generation_id={self.generation_id!r}, "
            f"output_path={self.output_path!r})"
        )


class GammaMCPBridge:
    """
    Gamma MCP ↔ Python 브릿지

    Claude Code가 Gamma MCP 도구를 호출한 결과를
    Python 코드(DesignAgent, ImagePipeline)에서 활용할 수 있도록 변환합니다.
    """

    def __init__(self, project_dir: Optional[Path] = None):
        self.project_dir = project_dir or Path("output")
        self._gamma_meta_path = self.project_dir / ".gamma_meta.json"
        self._gamma_meta: Dict[str, Any] = {}
        self._load_meta()

    # ── 1. 테마 컨설팅 ───────────────────────────────────────

    @staticmethod
    def process_gamma_themes(
        gamma_themes_raw: List[Dict[str, Any]],
        project_keywords: List[str],
    ) -> List["ThemeRecommendation"]:
        """
        Gamma get_themes() 결과 → ThemeRecommendation 변환.

        Claude Code가 Gamma MCP get_themes()를 호출한 결과를
        DesignAgent.interpret_gamma_themes()에 전달 가능한 형태로 정규화.

        Args:
            gamma_themes_raw: Gamma get_themes() 원본 결과
            project_keywords: 프로젝트 키워드 ["게임", "부스", "AGF"]

        Returns:
            List[ThemeRecommendation]
        """
        from src.agents.design_agent import DesignAgent

        agent = DesignAgent()
        return agent.interpret_gamma_themes(gamma_themes_raw, project_keywords)

    # ── 2. Gamma generate() 입력 준비 ─────────────────────────

    @staticmethod
    def prepare_content_for_gamma(
        pptx_path: Path,
        max_slides: int = 0,
    ) -> str:
        """
        PPTX에서 텍스트 추출 → Gamma generate() inputText 변환.

        slide_kit으로 생성한 초안 PPTX의 텍스트 콘텐츠를 추출하여
        Gamma MCP generate()의 inputText로 사용합니다.

        Args:
            pptx_path: slide_kit 초안 PPTX 경로
            max_slides: 추출 슬라이드 수 제한 (0=전체)

        Returns:
            Gamma generate() inputText 문자열
        """
        try:
            from pptx import Presentation

            prs = Presentation(str(pptx_path))
            lines = []

            for i, slide in enumerate(prs.slides):
                if max_slides and i >= max_slides:
                    break

                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text and text not in ("", " "):
                            # IMG_PH 플레이스홀더는 이미지 설명으로 변환
                            if "IMG_PH" in text:
                                desc = text.replace("IMG_PH", "").strip(" ()（）:：")
                                if desc:
                                    slide_texts.append(f"[Image: {desc}]")
                            else:
                                slide_texts.append(text)

                if slide_texts:
                    lines.append(f"--- Slide {i + 1} ---")
                    lines.extend(slide_texts)
                    lines.append("")

            result = "\n".join(lines)
            logger.info(
                f"PPTX 텍스트 추출: {len(prs.slides)}슬라이드, "
                f"{len(result)}자"
            )
            return result

        except ImportError:
            logger.error("python-pptx 미설치")
            return ""
        except Exception as e:
            logger.error(f"PPTX 텍스트 추출 실패: {e}")
            return ""

    @staticmethod
    def build_gamma_params(
        input_text: str,
        brief: Optional[Any] = None,
        num_cards: Optional[int] = None,
        export_as: Optional[str] = None,
        additional_instructions: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Gamma MCP generate() 호출 파라미터 생성.

        Claude Code가 이 결과를 Gamma MCP generate() 도구에 전달합니다.

        파라미터 구성 전략:
            themeId        — brief.gamma_recommendations 기반 (기구현)
            imageOptions   — 프로젝트 맥락 반영 (source + style + model)
            textOptions    — brief.content_tone 기반 tone 동적 설정
            cardOptions    — 16x9 + headerFooter로 슬라이드 간 일관성 확보
            additionalInstructions — 디자인 규칙/시각 위계/레이아웃 힌트

        Args:
            input_text: prepare_content_for_gamma() 결과
            brief: MergedDesignBrief (테마/이미지 스타일 반영)
            num_cards: 슬라이드 수
            export_as: "pptx" 또는 "pdf"
            additional_instructions: 사용자 직접 지정 디자인 지시사항
                (지정 시 자동 생성 대신 이 값 사용)

        Returns:
            Gamma generate() 파라미터 딕셔너리
        """
        params: Dict[str, Any] = {
            "inputText": input_text,
            "textMode": "preserve",
        }

        if num_cards:
            params["numCards"] = num_cards

        if export_as:
            params["exportAs"] = export_as

        if not brief:
            return params

        # ── themeId (기구현 — brief 기반) ──
        gamma_theme_id = getattr(brief, "gamma_theme_id", None)
        if not gamma_theme_id and getattr(brief, "gamma_recommendations", None):
            gamma_theme_id = brief.gamma_recommendations[0].theme_id
        if gamma_theme_id:
            params["themeId"] = gamma_theme_id

        # ── imageOptions (v6.0 고도화) ──
        params["imageOptions"] = _build_gamma_image_options(brief)

        # ── textOptions (기구현 — tone 유지) ──
        tone_text = "professional"
        content_tone = getattr(brief, "content_tone", None)
        if content_tone:
            level = getattr(content_tone, "emotional_tone_level", 3)
            if level >= 4:
                tone_text = "narrative and emotionally engaging"
            elif level <= 2:
                tone_text = "formal and data-driven"
            else:
                tone_text = "professional with balanced insight"

        params["textOptions"] = {
            "language": "ko",
            "tone": tone_text,
            "amount": "medium",
        }

        # ── cardOptions (v6.0 — headerFooter 추가) ──
        params["cardOptions"] = _build_gamma_card_options(brief)

        # ── additionalInstructions (v6.0 고도화) ──
        if additional_instructions:
            params["additionalInstructions"] = additional_instructions[:2000]
        else:
            hints = _build_additional_instructions(brief)
            if hints:
                params["additionalInstructions"] = hints[:2000]

        return params

    # ── 3. Gamma 결과 저장/조회 ───────────────────────────────

    def store_gamma_result(
        self,
        gamma_url: str,
        theme_id: str = "",
        theme_name: str = "",
        project_name: str = "",
    ) -> None:
        """
        Gamma generate() 결과 메타데이터 저장.

        다음 버전 생성 시 참조용으로 저장합니다.
        """
        self._gamma_meta = {
            "gamma_url": gamma_url,
            "theme_id": theme_id,
            "theme_name": theme_name,
            "project_name": project_name,
        }
        self._save_meta()
        logger.info(f"Gamma 결과 저장: {gamma_url}")

    def get_gamma_meta(self) -> Dict[str, Any]:
        """저장된 Gamma 메타데이터 조회."""
        return dict(self._gamma_meta)

    def _load_meta(self):
        if self._gamma_meta_path.exists():
            try:
                self._gamma_meta = json.loads(
                    self._gamma_meta_path.read_text(encoding="utf-8")
                )
            except Exception:
                self._gamma_meta = {}

    def _save_meta(self):
        self.project_dir.mkdir(parents=True, exist_ok=True)
        self._gamma_meta_path.write_text(
            json.dumps(self._gamma_meta, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )

    # ── 4. Gamma 자동 파이프라인 (v7.0) ──────────────────────

    def run_gamma_pipeline(
        self,
        pptx_path: Path,
        brief: Optional[Any] = None,
        num_cards: Optional[int] = None,
        export_as: str = "pptx",
        project_name: str = "",
    ) -> GammaPipelineResult:
        """
        Gamma 전송을 위한 전체 준비 단계를 원스텝으로 실행.

        텍스트 추출 → 파라미터 빌드 → 중간 파일 저장까지
        모든 준비를 완료하고, Claude Code가 Gamma MCP generate()를
        호출할 수 있는 상태의 GammaPipelineResult를 반환합니다.

        사용법:
            bridge = GammaMCPBridge(project_dir=Path("output/PROJECT"))
            result = bridge.run_gamma_pipeline(pptx_path, brief, num_cards=52)
            # → Claude Code: Gamma MCP generate(**result.params)
            # → bridge.handle_gamma_response(response, result)
            # → bridge.poll_and_download_gamma(result.generation_id, result)

        Args:
            pptx_path: slide_kit 초안 PPTX 경로
            brief: MergedDesignBrief (테마/이미지 스타일 반영)
            num_cards: 슬라이드 수
            export_as: 내보내기 형식 ("pptx" 또는 "pdf")
            project_name: 프로젝트명 (파일 이름 + 메타데이터용)

        Returns:
            GammaPipelineResult — params 필드에 generate() 파라미터 포함
        """
        result = GammaPipelineResult()
        result.created_at = datetime.now().isoformat()
        result.source_pptx = str(pptx_path)
        result.export_format = export_as
        result.project_name = (
            project_name
            or getattr(brief, "project_name", "")
            or pptx_path.stem
        )

        try:
            # 1. 텍스트 추출
            result.input_text = self.prepare_content_for_gamma(pptx_path)
            if not result.input_text:
                result.status = "error"
                result.error = "PPTX 텍스트 추출 실패 (빈 텍스트)"
                return result

            # 2. 파라미터 빌드
            result.params = self.build_gamma_params(
                input_text=result.input_text,
                brief=brief,
                num_cards=num_cards,
                export_as=export_as,
            )

            # themeId 추적
            result.theme_id = result.params.get("themeId", "")

            # 3. 중간 파일 저장 (디버깅/재실행용)
            self.project_dir.mkdir(parents=True, exist_ok=True)

            text_path = self.project_dir / "gamma_input_text.txt"
            text_path.write_text(result.input_text, encoding="utf-8")
            result.input_text_path = str(text_path)

            params_path = self.project_dir / "gamma_params.json"
            params_path.write_text(
                json.dumps(
                    result.params, ensure_ascii=False, indent=2, default=str
                ),
                encoding="utf-8",
            )
            result.params_path = str(params_path)

            # 4. 파이프라인 상태 저장
            result.status = "prepared"
            self._save_pipeline_state(result)

            slide_count = result.input_text.count("--- Slide ")
            logger.info(
                f"Gamma 파이프라인 준비 완료: "
                f"{slide_count}슬라이드, "
                f"{len(result.input_text)}자, "
                f"themeId={result.theme_id or '(없음)'}"
            )
            return result

        except Exception as e:
            result.status = "error"
            result.error = str(e)
            logger.error(f"Gamma 파이프라인 준비 실패: {e}")
            return result

    def handle_gamma_response(
        self,
        gamma_response: Dict[str, Any],
        pipeline_result: Optional[GammaPipelineResult] = None,
    ) -> Dict[str, Any]:
        """
        Gamma MCP generate() 응답을 처리하고 다운로드 지침을 반환.

        generate() 호출 후 Claude Code가 이 메서드를 호출하면:
        1. 응답에서 gammaUrl, generationId 추출
        2. 메타데이터 저장
        3. 다운로드 방법 (브라우저 지침) 반환

        Args:
            gamma_response: Gamma MCP generate() 반환값
                예: {"gammaUrl": "https://...", "generationId": "abc",
                     "status": "pending", "exportUrl": "https://..."}
            pipeline_result: run_gamma_pipeline()에서 반환된 결과 (상태 추적용)

        Returns:
            {
                "status": "ready_to_download" | "poll_available" | "pending",
                "gamma_url": "https://gamma.app/...",
                "generation_id": "...",
                "download_method": "direct" | "api_poll" | "browser",
                "export_url": "..." (direct일 때만),
                "instruction": "..." (api_poll일 때만),
                "browser_steps": [...] (browser일 때만),
                "suggested_filename": "PROJECT_gamma.pptx",
                "output_dir": "output/PROJECT/gamma_pptx/"
            }
        """
        gamma_url = gamma_response.get("gammaUrl", "")
        generation_id = gamma_response.get("generationId", "")
        status = gamma_response.get("status", "unknown")
        export_url = gamma_response.get("exportUrl", "")

        # 프로젝트명 결정
        project_name = ""
        if pipeline_result:
            project_name = pipeline_result.project_name
        if not project_name:
            project_name = self._gamma_meta.get("project_name", "project")

        # 메타데이터 저장
        theme_id = ""
        theme_name = ""
        if pipeline_result:
            theme_id = pipeline_result.theme_id
            theme_name = pipeline_result.theme_name
        self.store_gamma_result(
            gamma_url=gamma_url,
            theme_id=theme_id,
            theme_name=theme_name,
            project_name=project_name,
        )

        # pipeline_result 업데이트
        if pipeline_result:
            pipeline_result.gamma_url = gamma_url
            pipeline_result.generation_id = generation_id
            pipeline_result.export_url = export_url
            pipeline_result.status = "sent"

        # 출력 경로 준비
        gamma_pptx_dir = self.project_dir / "gamma_pptx"
        gamma_pptx_dir.mkdir(parents=True, exist_ok=True)

        safe_name = project_name.replace(" ", "_").replace("/", "_")
        timestamp = datetime.now().strftime("%m%d_%H%M")
        export_ext = (
            pipeline_result.export_format if pipeline_result else "pptx"
        )
        suggested_filename = f"{safe_name}_gamma_{timestamp}.{export_ext}"

        if pipeline_result:
            pipeline_result.output_path = str(
                gamma_pptx_dir / suggested_filename
            )

        # 다운로드 전략 결정
        #
        # Gamma MCP generate() 응답에 exportUrl이 있으면 직접 다운로드.
        # 없으면 Gamma Public API 폴링으로 exportUrl 획득 후 다운로드.
        # API 키가 없으면 브라우저 수동 내보내기 안내.
        download_info: Dict[str, Any] = {
            "gamma_url": gamma_url,
            "generation_id": generation_id,
            "suggested_filename": suggested_filename,
            "output_dir": str(gamma_pptx_dir),
            "output_path": str(gamma_pptx_dir / suggested_filename),
        }

        if export_url:
            # 직접 다운로드 URL이 응답에 포함된 경우
            download_info["status"] = "ready_to_download"
            download_info["download_method"] = "direct"
            download_info["export_url"] = export_url
        elif generation_id and os.environ.get("GAMMA_API_KEY"):
            # API 키가 있으면 폴링 + 다운로드 가능
            download_info["status"] = "poll_available"
            download_info["download_method"] = "api_poll"
            download_info["instruction"] = (
                "bridge.poll_and_download_gamma("
                f"'{generation_id}', pipeline_result) 호출하여 "
                "폴링 + 자동 다운로드"
            )
        else:
            # API 키 없음 → 브라우저 수동 내보내기
            download_info["status"] = (
                "ready_to_download"
                if status in ("completed", "done", "ready")
                else "pending"
            )
            download_info["download_method"] = "browser"
            download_info["browser_steps"] = [
                f"navigate to {gamma_url}",
                "wait for generation to complete",
                "click '...' menu → '내보내기...' → 'PowerPoint로 내보내기'",
                f"downloaded file → bridge.save_gamma_pptx(path, result)",
            ]

        # 파이프라인 상태 저장
        if pipeline_result:
            self._save_pipeline_state(pipeline_result)

        logger.info(
            f"Gamma 응답 처리: generation_id={generation_id}, "
            f"status={status}, method={download_info.get('download_method', 'unknown')}"
        )
        return download_info

    def poll_and_download_gamma(
        self,
        generation_id: str,
        pipeline_result: Optional[GammaPipelineResult] = None,
        poll_interval: float = 5.0,
        timeout: float = 300.0,
        api_key: str = "",
    ) -> Path:
        """
        Gamma Public API를 폴링하여 생성 완료 시 PPTX를 다운로드.

        Gamma MCP generate() 호출 후 반환된 generationId로
        API를 폴링하고, 완료되면 exportUrl에서 PPTX를 다운로드합니다.

        API 엔드포인트:
            GET https://public-api.gamma.app/v1.0/generations/{generationId}
            Header: X-API-KEY: {api_key}

        완료 응답에서 downloadLink를 추출하여 PPTX를 다운로드합니다.
        (POST 시 exportAs: "pptx" 필수 — 없으면 downloadLink 미반환)

        다운로드 URL 필드 우선순위:
            downloadLink → exportUrl → pptxUrl → pdfUrl → files[0].url

        Args:
            generation_id: Gamma generationId (generate() 응답에서 추출)
            pipeline_result: run_gamma_pipeline()에서 반환된 결과 (상태 추적용)
            poll_interval: 폴링 간격 (초, 기본 5.0 — Gamma 공식 권장)
            timeout: 최대 대기 시간 (초, 기본 300 = 5분)
            api_key: Gamma API 키 (빈 문자열이면 환경변수에서 읽음)

        Returns:
            다운로드된 PPTX 파일의 최종 저장 경로

        Raises:
            ValueError: API 키가 없거나 generation_id가 비어있을 때
            TimeoutError: 폴링 타임아웃
            RuntimeError: 다운로드 실패
        """
        if not generation_id:
            raise ValueError("generation_id가 비어있습니다.")

        if not api_key:
            api_key = os.environ.get("GAMMA_API_KEY", "")
        if not api_key:
            raise ValueError(
                "GAMMA_API_KEY가 설정되지 않았습니다. "
                ".env 또는 환경변수에 GAMMA_API_KEY를 설정하세요."
            )

        api_base = "https://public-api.gamma.app/v1.0/generations"
        poll_url = f"{api_base}/{generation_id}"
        headers = {
            "X-API-KEY": api_key,
            "User-Agent": "GammaBridge/1.0",
            "Accept": "application/json",
        }

        logger.info(
            f"Gamma 폴링 시작: generation_id={generation_id}, "
            f"interval={poll_interval}s, timeout={timeout}s"
        )

        start_time = time.time()
        poll_count = 0
        consecutive_errors = 0
        max_consecutive_errors = 10

        while True:
            elapsed = time.time() - start_time
            if elapsed > timeout:
                msg = (
                    f"Gamma 폴링 타임아웃 ({timeout}s): "
                    f"generation_id={generation_id}, "
                    f"polls={poll_count}"
                )
                logger.error(msg)
                if pipeline_result:
                    pipeline_result.status = "error"
                    pipeline_result.error = msg
                    self._save_pipeline_state(pipeline_result)
                raise TimeoutError(msg)

            poll_count += 1
            try:
                req = urllib.request.Request(poll_url, headers=headers)
                with urllib.request.urlopen(req, timeout=30) as resp:
                    data = json.loads(resp.read().decode("utf-8"))
                consecutive_errors = 0  # 성공 시 리셋
            except Exception as e:
                consecutive_errors += 1
                logger.warning(
                    f"Gamma 폴링 {poll_count}회 HTTP 오류 "
                    f"({consecutive_errors}/{max_consecutive_errors}): {e}"
                )
                if consecutive_errors >= max_consecutive_errors:
                    msg = (
                        f"Gamma 폴링 연속 {max_consecutive_errors}회 HTTP 오류로 중단: "
                        f"마지막 오류={e}"
                    )
                    logger.error(msg)
                    if pipeline_result:
                        pipeline_result.status = "error"
                        pipeline_result.error = msg
                        self._save_pipeline_state(pipeline_result)
                    raise RuntimeError(msg) from e
                time.sleep(poll_interval)
                continue

            status = data.get("status", "unknown")
            logger.info(
                f"Gamma 폴링 {poll_count}회: status={status} "
                f"({elapsed:.0f}s/{timeout:.0f}s)"
            )

            if status in ("failed", "error"):
                raw_error = data.get("error", data.get("message", "알 수 없는 오류"))
                # error가 dict인 경우 (예: {"message": "...", "statusCode": 422})
                if isinstance(raw_error, dict):
                    error_msg = raw_error.get("message", str(raw_error))
                else:
                    error_msg = str(raw_error)
                msg = f"Gamma 생성 실패: {error_msg}"
                logger.error(msg)
                if pipeline_result:
                    pipeline_result.status = "error"
                    pipeline_result.error = msg
                    self._save_pipeline_state(pipeline_result)
                raise RuntimeError(msg)

            if status in ("completed", "done", "ready"):
                # 다운로드 URL 추출 (여러 필드 후보 — 우선순위 순)
                export_url = (
                    data.get("downloadLink")
                    or data.get("exportUrl")
                    or data.get("pptxUrl")
                    or data.get("pdfUrl")
                )
                if not export_url:
                    # nested 구조 탐색
                    files = data.get("files") or []
                    if not files:
                        result_obj = data.get("result", {})
                        files = result_obj.get("files") if isinstance(result_obj, dict) else []
                    if files and isinstance(files, list) and len(files) > 0:
                        export_url = files[0].get("url", "")

                if not export_url:
                    msg = (
                        f"Gamma 생성 완료되었지만 exportUrl 없음: "
                        f"keys={list(data.keys())}"
                    )
                    logger.error(msg)
                    if pipeline_result:
                        pipeline_result.status = "error"
                        pipeline_result.error = msg
                        pipeline_result.gamma_url = data.get("gammaUrl", "")
                        self._save_pipeline_state(pipeline_result)
                    raise RuntimeError(msg)

                # pipeline_result 업데이트
                if pipeline_result:
                    pipeline_result.export_url = export_url
                    pipeline_result.gamma_url = data.get("gammaUrl", "")
                    pipeline_result.status = "completed"
                    self._save_pipeline_state(pipeline_result)

                logger.info(f"Gamma 생성 완료: exportUrl 획득")

                # PPTX 다운로드
                return self._download_export_url(
                    export_url, pipeline_result
                )

            # 아직 진행 중 → 대기 후 재폴링
            time.sleep(poll_interval)

    def _download_export_url(
        self,
        export_url: str,
        pipeline_result: Optional[GammaPipelineResult] = None,
    ) -> Path:
        """
        exportUrl에서 PPTX 파일을 다운로드하여 프로젝트 디렉토리에 저장.

        Args:
            export_url: Gamma exportUrl (직접 다운로드 URL)
            pipeline_result: 상태 추적용

        Returns:
            저장된 PPTX 파일 경로
        """
        gamma_pptx_dir = self.project_dir / "gamma_pptx"
        gamma_pptx_dir.mkdir(parents=True, exist_ok=True)

        # 파일명 결정
        if pipeline_result and pipeline_result.output_path:
            output_path = Path(pipeline_result.output_path)
        else:
            project_name = self._gamma_meta.get("project_name", "project")
            safe_name = project_name.replace(" ", "_").replace("/", "_")
            timestamp = datetime.now().strftime("%m%d_%H%M")
            output_path = gamma_pptx_dir / f"{safe_name}_gamma_{timestamp}.pptx"

        # 임시 파일에 다운로드
        tmp_path = output_path.with_suffix(".pptx.tmp")

        try:
            logger.info(f"PPTX 다운로드 시작: {export_url}")
            req = urllib.request.Request(export_url)
            req.add_header("User-Agent", "GammaBridge/1.0")
            # API 키가 필요할 수 있으므로 환경변수에서 읽기
            api_key = os.environ.get("GAMMA_API_KEY", "")
            if api_key:
                req.add_header("X-API-KEY", api_key)

            with urllib.request.urlopen(req, timeout=600) as resp:
                with open(str(tmp_path), "wb") as f:
                    shutil.copyfileobj(resp, f)

            # 파일 크기 검증
            file_size = tmp_path.stat().st_size
            if file_size < 1024:  # 1KB 미만이면 오류 응답일 가능성
                content = tmp_path.read_bytes()
                try:
                    # JSON 에러 응답인지 확인
                    error_data = json.loads(content)
                    msg = (
                        f"다운로드된 파일이 에러 응답: "
                        f"{error_data.get('error', content[:200])}"
                    )
                    tmp_path.unlink(missing_ok=True)
                    raise RuntimeError(msg)
                except (json.JSONDecodeError, UnicodeDecodeError):
                    pass  # 바이너리 파일이면 정상

            # 임시 → 최종 파일로 이동
            tmp_path.rename(output_path)

            file_size_mb = output_path.stat().st_size / (1024 * 1024)
            logger.info(
                f"Gamma PPTX 다운로드 완료: {output_path} "
                f"({file_size_mb:.1f}MB)"
            )

            # 메타데이터 업데이트
            self._gamma_meta["output_pptx"] = str(output_path)
            self._gamma_meta["downloaded_at"] = datetime.now().isoformat()
            self._gamma_meta["file_size_bytes"] = output_path.stat().st_size
            self._gamma_meta["export_url"] = export_url
            self._save_meta()

            # pipeline_result 업데이트
            if pipeline_result:
                pipeline_result.downloaded_path = str(output_path)
                pipeline_result.output_path = str(output_path)
                pipeline_result.status = "downloaded"
                self._save_pipeline_state(pipeline_result)

            return output_path

        except Exception as e:
            tmp_path.unlink(missing_ok=True)
            msg = f"PPTX 다운로드 실패: {e}"
            logger.error(msg)
            if pipeline_result:
                pipeline_result.status = "error"
                pipeline_result.error = msg
                self._save_pipeline_state(pipeline_result)
            raise RuntimeError(msg) from e

    def save_gamma_pptx(
        self,
        source_path: Path,
        pipeline_result: Optional[GammaPipelineResult] = None,
        version_tag: str = "",
    ) -> Path:
        """
        다운로드된 Gamma PPTX를 프로젝트 디렉토리에 저장.

        브라우저에서 다운로드된 파일을 프로젝트 구조 내 적절한 위치로
        복사하고 메타데이터를 업데이트합니다.

        Args:
            source_path: 다운로드된 PPTX 파일 경로
                (예: ~/Downloads/Something.pptx)
            pipeline_result: run_gamma_pipeline()에서 반환된 결과
            version_tag: 버전 태그 (예: "v3", "final")

        Returns:
            최종 저장 경로
        """
        source_path = Path(source_path)
        if not source_path.exists():
            raise FileNotFoundError(f"소스 파일 없음: {source_path}")

        # 출력 경로 결정
        gamma_pptx_dir = self.project_dir / "gamma_pptx"
        gamma_pptx_dir.mkdir(parents=True, exist_ok=True)

        if pipeline_result and pipeline_result.output_path:
            # handle_gamma_response()에서 결정한 경로 사용
            output_path = Path(pipeline_result.output_path)
        else:
            # 이름 생성
            project_name = self._gamma_meta.get("project_name", "project")
            safe_name = project_name.replace(" ", "_").replace("/", "_")
            timestamp = datetime.now().strftime("%m%d_%H%M")
            tag = f"_{version_tag}" if version_tag else ""
            ext = source_path.suffix or ".pptx"
            output_path = (
                gamma_pptx_dir / f"{safe_name}_gamma{tag}_{timestamp}{ext}"
            )

        # 파일 복사
        shutil.copy2(str(source_path), str(output_path))

        # 메타데이터 업데이트
        self._gamma_meta["output_pptx"] = str(output_path)
        self._gamma_meta["downloaded_at"] = datetime.now().isoformat()
        self._gamma_meta["file_size_bytes"] = output_path.stat().st_size
        self._save_meta()

        # pipeline_result 업데이트
        if pipeline_result:
            pipeline_result.downloaded_path = str(source_path)
            pipeline_result.output_path = str(output_path)
            pipeline_result.status = "downloaded"
            self._save_pipeline_state(pipeline_result)

        file_size_mb = output_path.stat().st_size / (1024 * 1024)
        logger.info(
            f"Gamma PPTX 저장 완료: {output_path} ({file_size_mb:.1f}MB)"
        )
        return output_path

    @staticmethod
    def find_downloaded_gamma_pptx(
        downloads_dir: Optional[Path] = None,
        after_time: Optional[str] = None,
        keyword: str = "",
    ) -> Optional[Path]:
        """
        브라우저 Downloads 폴더에서 Gamma 내보내기 PPTX를 탐색.

        gammaUrl 페이지에서 '내보내기 → PowerPoint' 클릭 후
        다운로드된 파일을 찾습니다.

        탐색 전략:
            1. downloads_dir 내 *.pptx 파일을 수정 시간 역순 정렬
            2. after_time 이후 생성된 파일만 필터
            3. keyword가 있으면 파일명에 포함된 것 우선
            4. 가장 최근 PPTX 반환

        Args:
            downloads_dir: 다운로드 폴더 (기본: ~/Downloads)
            after_time: 이 시각 이후 파일만 (ISO format, 예: "2026-03-02T12:00:00")
            keyword: 파일명에 포함되어야 할 키워드 (예: "NIKKE")

        Returns:
            발견된 PPTX 경로, 없으면 None
        """
        import os

        if downloads_dir is None:
            # 플랫폼별 기본 다운로드 경로
            home = Path.home()
            downloads_dir = home / "Downloads"
            if not downloads_dir.exists():
                # Windows 한글 경로 등 대안
                for alt in ["다운로드", "Download"]:
                    alt_dir = home / alt
                    if alt_dir.exists():
                        downloads_dir = alt_dir
                        break

        if not downloads_dir.exists():
            logger.warning(f"다운로드 폴더 없음: {downloads_dir}")
            return None

        # PPTX 파일 수집 + 수정시간 기준 정렬
        pptx_files: List[Tuple[Path, float]] = []
        for f in downloads_dir.iterdir():
            if f.suffix.lower() == ".pptx" and f.is_file():
                pptx_files.append((f, f.stat().st_mtime))

        if not pptx_files:
            logger.info("다운로드 폴더에 PPTX 파일 없음")
            return None

        # 수정 시간 역순 정렬 (최신 먼저)
        pptx_files.sort(key=lambda x: x[1], reverse=True)

        # after_time 필터
        if after_time:
            try:
                cutoff = datetime.fromisoformat(after_time).timestamp()
                pptx_files = [
                    (f, t) for f, t in pptx_files if t >= cutoff
                ]
            except ValueError:
                pass

        if not pptx_files:
            logger.info(f"after_time({after_time}) 이후 PPTX 파일 없음")
            return None

        # keyword 매칭 우선
        if keyword:
            kw_lower = keyword.lower()
            keyword_matches = [
                (f, t) for f, t in pptx_files
                if kw_lower in f.name.lower()
            ]
            if keyword_matches:
                found = keyword_matches[0][0]
                logger.info(
                    f"Gamma PPTX 발견 (keyword '{keyword}'): {found.name} "
                    f"({found.stat().st_size / (1024*1024):.1f}MB)"
                )
                return found

        # keyword 매칭 없으면 가장 최근 PPTX
        found = pptx_files[0][0]
        logger.info(
            f"Gamma PPTX 발견 (최신): {found.name} "
            f"({found.stat().st_size / (1024*1024):.1f}MB)"
        )
        return found

    @staticmethod
    def get_gamma_export_url(
        gamma_url: str,
        export_format: str = "pptx",
    ) -> str:
        """
        gammaUrl에서 직접 다운로드 URL을 추론 (DEPRECATED).

        ⚠️ 실제 Gamma API의 exportUrl 패턴
        (https://assets.api.gamma.app/export/pptx/{id}/...)과
        이 함수의 추론 패턴이 일치하지 않습니다.
        poll_and_download_gamma()를 사용하세요.

        Gamma 응답에 exportUrl이 없을 때의 레거시 fallback입니다.
        Gamma의 URL 패턴 기반 추론 — 실제 다운로드 가능 여부는 보장하지 않음.

        Args:
            gamma_url: Gamma 프레젠테이션 URL
            export_format: "pptx" 또는 "pdf"

        Returns:
            추론된 다운로드 URL (빈 문자열이면 브라우저 사용 필요)
        """
        if not gamma_url:
            return ""

        # Gamma URL 패턴:
        #   https://gamma.app/generations/{generationId}
        #   https://gamma.app/docs/{slug}-{id}
        # 내보내기 URL 패턴 (추론):
        #   https://gamma.app/api/export/{id}?format=pptx

        # generationId 추출 시도
        generation_id = ""
        if "/generations/" in gamma_url:
            generation_id = gamma_url.split("/generations/")[-1].split("?")[0]
        elif "/docs/" in gamma_url:
            # slug-id 형태에서 id 부분 추출
            slug_id = gamma_url.split("/docs/")[-1].split("?")[0]
            parts = slug_id.rsplit("-", 1)
            if len(parts) == 2:
                generation_id = parts[1]

        if not generation_id:
            return ""

        # 추론된 내보내기 URL (Gamma API 패턴 기반)
        return (
            f"https://gamma.app/api/export/{generation_id}"
            f"?format={export_format}"
        )

    def get_pipeline_status(self) -> Dict[str, Any]:
        """
        현재 파이프라인 상태 조회.

        저장된 파이프라인 상태 파일에서 읽어옵니다.
        이전 실행 결과를 확인하거나 중단된 작업을 이어갈 때 사용.

        Returns:
            파이프라인 상태 딕셔너리 (없으면 빈 dict)
        """
        state_path = self.project_dir / ".gamma_pipeline_state.json"
        if not state_path.exists():
            return {}
        try:
            return json.loads(state_path.read_text(encoding="utf-8"))
        except Exception:
            return {}

    def _save_pipeline_state(self, result: GammaPipelineResult) -> None:
        """파이프라인 상태를 파일에 저장."""
        state_path = self.project_dir / ".gamma_pipeline_state.json"
        self.project_dir.mkdir(parents=True, exist_ok=True)
        state_path.write_text(
            json.dumps(result.to_dict(), ensure_ascii=False, indent=2),
            encoding="utf-8",
        )

    # ── 5. 사용자 편집 역파싱 (v6.0 고도화) ─────────────────────

    @staticmethod
    def extract_user_edits(
        original_pptx: Path,
        edited_pptx: Path,
    ) -> Dict[str, Any]:
        """
        사용자가 편집한 PPTX와 원본을 비교하여 변경사항 추출.

        v6.0: 14개 카테고리 전수 비교
            1. 배경 (solid/gradient/image)
            2. 도형 수 (추가/삭제)
            3. 텍스트 내용
            4. 도형 지오메트리 (위치/크기/회전)
            5. 도형 채우기 (solid/gradient/pattern)
            6. 선/테두리 (색상/두께/스타일)
            7. 타이포그래피 (폰트명/크기/굵기/색상/기울임)
            8. 단락 서식 (정렬/줄간격/들여쓰기/불릿)
            9. 도형 유형 변경 (사각형→라운드 등)
            10. 그림자/효과 (shadow XML)
            11. 이미지 교체 (blob hash)
            12. 텍스트 프레임 속성 (여백/자동크기/앵커)
            13. 그룹 도형 변경
            14. 슬라이드 레이아웃 참조 변경

        Args:
            original_pptx: slide_kit으로 생성한 원본
            edited_pptx: 사용자가 수정한 버전

        Returns:
            {
                "modified_slides": [
                    {
                        "index": 3,
                        "changes": ["background", "typography", ...],
                        "design_overrides": { ... },
                        "shape_changes": [ ... ]
                    }
                ],
                "global_changes": {
                    "slide_count_change": 0,
                    "dominant_colors": {...},
                    "dominant_fonts": {...},
                    "color_palette_shift": {...}
                },
                "version": "6.0"
            }
        """
        try:
            from pptx import Presentation

            orig = Presentation(str(original_pptx))
            edit = Presentation(str(edited_pptx))

            result: Dict[str, Any] = {
                "modified_slides": [],
                "global_changes": {},
                "version": "6.0",
            }

            orig_count = len(orig.slides)
            edit_count = len(edit.slides)

            if edit_count != orig_count:
                result["global_changes"]["slide_count_change"] = (
                    edit_count - orig_count
                )

            # 전체 프레젠테이션 색상/폰트 수집 (글로벌 통계용)
            all_orig_colors: List[str] = []
            all_edit_colors: List[str] = []
            all_orig_fonts: List[str] = []
            all_edit_fonts: List[str] = []

            # ── 슬라이드별 비교 (공통 범위) ──
            for i in range(min(orig_count, edit_count)):
                orig_slide = orig.slides[i]
                edit_slide = edit.slides[i]

                slide_diff = _compare_slides(
                    orig_slide, edit_slide, slide_index=i
                )

                # 글로벌 통계 수집
                all_orig_colors.extend(
                    slide_diff.pop("_orig_colors", [])
                )
                all_edit_colors.extend(
                    slide_diff.pop("_edit_colors", [])
                )
                all_orig_fonts.extend(
                    slide_diff.pop("_orig_fonts", [])
                )
                all_edit_fonts.extend(
                    slide_diff.pop("_edit_fonts", [])
                )

                if slide_diff["changes"]:
                    result["modified_slides"].append(slide_diff)

            # ── 추가된 슬라이드 ──
            if edit_count > orig_count:
                for i in range(orig_count, edit_count):
                    new_slide = edit.slides[i]
                    new_slide_info = _extract_full_slide_info(
                        new_slide, slide_index=i
                    )
                    result["modified_slides"].append({
                        "index": i,
                        "changes": ["new_slide"],
                        "design_overrides": {
                            "background": _extract_bg_info(new_slide),
                            "preserve": True,
                        },
                        "shape_changes": [],
                        "new_slide_snapshot": new_slide_info,
                    })

            # ── 삭제된 슬라이드 ──
            if edit_count < orig_count:
                result["global_changes"]["deleted_slide_indices"] = list(
                    range(edit_count, orig_count)
                )

            # ── 글로벌 디자인 통계 ──
            result["global_changes"]["dominant_colors"] = {
                "original": _top_items(all_orig_colors, 5),
                "edited": _top_items(all_edit_colors, 5),
            }
            result["global_changes"]["dominant_fonts"] = {
                "original": _top_items(all_orig_fonts, 3),
                "edited": _top_items(all_edit_fonts, 3),
            }

            # 색상 팔레트 변화 감지
            orig_color_set = set(all_orig_colors)
            edit_color_set = set(all_edit_colors)
            new_colors = edit_color_set - orig_color_set
            removed_colors = orig_color_set - edit_color_set
            if new_colors or removed_colors:
                result["global_changes"]["color_palette_shift"] = {
                    "added": sorted(new_colors),
                    "removed": sorted(removed_colors),
                }

            # 폰트 변화 감지
            orig_font_set = set(all_orig_fonts)
            edit_font_set = set(all_edit_fonts)
            new_fonts = edit_font_set - orig_font_set
            if new_fonts:
                result["global_changes"]["new_fonts_introduced"] = sorted(
                    new_fonts
                )

            mod_count = len(result["modified_slides"])
            logger.info(
                f"사용자 편집 분석 (v6.0): "
                f"{mod_count}개 슬라이드 변경, "
                f"{len(new_colors)}개 새 색상, "
                f"{len(new_fonts)}개 새 폰트"
            )
            return result

        except ImportError:
            logger.error("python-pptx 미설치")
            return {
                "modified_slides": [],
                "global_changes": {},
                "version": "6.0",
            }
        except Exception as e:
            logger.error(f"편집 비교 실패: {e}")
            return {
                "modified_slides": [],
                "global_changes": {},
                "version": "6.0",
            }

    @staticmethod
    def save_design_overrides(
        overrides: Dict[str, Any],
        output_dir: Path,
        project_name: str = "",
    ) -> Path:
        """디자인 오버라이드를 JSON으로 저장 (다음 버전 참조용)."""
        out_path = (
            output_dir
            / f"{project_name or 'project'}_design_overrides.json"
        )
        out_path.write_text(
            json.dumps(
                overrides, ensure_ascii=False, indent=2, default=str
            ),
            encoding="utf-8",
        )
        logger.info(f"디자인 오버라이드 저장: {out_path}")
        return out_path

    @staticmethod
    def load_design_overrides(override_path: Path) -> Dict[str, Any]:
        """저장된 디자인 오버라이드 로드."""
        if not override_path.exists():
            return {}
        try:
            return json.loads(
                override_path.read_text(encoding="utf-8")
            )
        except Exception:
            return {}

    @staticmethod
    def summarize_overrides(overrides: Dict[str, Any]) -> Dict[str, Any]:
        """
        오버라이드 결과를 DesignAgent.merge()에 전달 가능한 요약으로 변환.

        extract_user_edits()의 상세 결과를 custom_colors, custom_fonts 등
        DesignAgent가 이해하는 형태로 압축합니다.

        Returns:
            {
                "custom_colors": {"primary": (R,G,B), ...},
                "custom_fonts": {"primary": "FontName", ...},
                "preserve_slides": [3, 7, 12],
                "design_hints": "..."
            }
        """
        summary: Dict[str, Any] = {}

        global_changes = overrides.get("global_changes", {})

        # 색상 팔레트 변화 → custom_colors 추출
        dominant = global_changes.get("dominant_colors", {})
        edit_colors = dominant.get("edited", {})
        if edit_colors:
            # 가장 많이 사용된 색상을 primary/secondary로 매핑
            top_colors = sorted(
                edit_colors.items(), key=lambda x: x[1], reverse=True
            )
            custom_colors: Dict[str, Tuple[int, int, int]] = {}
            for idx, (hex_color, _count) in enumerate(top_colors[:3]):
                try:
                    r = int(hex_color[0:2], 16)
                    g = int(hex_color[2:4], 16)
                    b = int(hex_color[4:6], 16)
                    if idx == 0:
                        custom_colors["primary"] = (r, g, b)
                    elif idx == 1:
                        custom_colors["secondary"] = (r, g, b)
                    elif idx == 2:
                        custom_colors["accent"] = (r, g, b)
                except (ValueError, IndexError):
                    pass
            if custom_colors:
                summary["custom_colors"] = custom_colors

        # 폰트 변화 → custom_fonts
        new_fonts = global_changes.get("new_fonts_introduced", [])
        if new_fonts:
            summary["custom_fonts"] = {"primary": new_fonts[0]}

        # 사용자가 직접 수정한 슬라이드 → 보존 대상
        preserve = []
        for slide_diff in overrides.get("modified_slides", []):
            if slide_diff.get("design_overrides", {}).get("preserve"):
                preserve.append(slide_diff["index"])
        if preserve:
            summary["preserve_slides"] = preserve

        return summary


# ═══════════════════════════════════════════════════════════════
#  내부 헬퍼 — 슬라이드 레벨 비교
# ═══════════════════════════════════════════════════════════════


def _compare_slides(
    orig_slide, edit_slide, slide_index: int
) -> Dict[str, Any]:
    """두 슬라이드를 전수 비교하여 변경 카테고리와 상세 정보 반환."""

    changes: List[str] = []
    overrides: Dict[str, Any] = {}
    shape_changes: List[Dict[str, Any]] = []

    # 통계 수집용
    orig_colors: List[str] = []
    edit_colors: List[str] = []
    orig_fonts: List[str] = []
    edit_fonts: List[str] = []

    # ── (1) 배경 ──
    orig_bg = _extract_bg_info(orig_slide)
    edit_bg = _extract_bg_info(edit_slide)
    if orig_bg != edit_bg:
        changes.append("background")
        overrides["background"] = edit_bg

    # ── (2) 도형 수 ──
    orig_shape_count = len(orig_slide.shapes)
    edit_shape_count = len(edit_slide.shapes)
    if edit_shape_count != orig_shape_count:
        changes.append("shape_count")
        overrides["shape_delta"] = edit_shape_count - orig_shape_count

    # ── (3~14) 도형별 상세 비교 ──
    orig_shapes = list(orig_slide.shapes)
    edit_shapes = list(edit_slide.shapes)

    # 도형 매칭 (이름 기반 → 위치/유형 폴백)
    matched_pairs, unmatched_orig, unmatched_edit = _match_shapes(
        orig_shapes, edit_shapes
    )

    for orig_sp, edit_sp in matched_pairs:
        sp_diff = _compare_shapes(orig_sp, edit_sp)

        # 색상/폰트 통계 수집
        orig_colors.extend(sp_diff.pop("_orig_colors", []))
        edit_colors.extend(sp_diff.pop("_edit_colors", []))
        orig_fonts.extend(sp_diff.pop("_orig_fonts", []))
        edit_fonts.extend(sp_diff.pop("_edit_fonts", []))

        if sp_diff["changes"]:
            shape_changes.append(sp_diff)
            for ch in sp_diff["changes"]:
                if ch not in changes:
                    changes.append(ch)

    # 새로 추가된 도형
    if unmatched_edit:
        if "new_shapes" not in changes:
            changes.append("new_shapes")
        for sp in unmatched_edit:
            sp_info = _extract_shape_snapshot(sp)
            shape_changes.append({
                "action": "added",
                "shape_name": sp.name,
                "changes": ["new_shape"],
                "snapshot": sp_info,
            })
            edit_colors.extend(sp_info.get("_colors", []))
            edit_fonts.extend(sp_info.get("_fonts", []))

    # 삭제된 도형
    if unmatched_orig:
        if "deleted_shapes" not in changes:
            changes.append("deleted_shapes")
        for sp in unmatched_orig:
            shape_changes.append({
                "action": "deleted",
                "shape_name": sp.name,
                "changes": ["deleted_shape"],
            })

    # ── (3) 텍스트 내용 (기존 호환) ──
    orig_text = _extract_all_text(orig_slide)
    edit_text = _extract_all_text(edit_slide)
    if orig_text != edit_text and "text_content" not in changes:
        changes.append("text_content")

    # ── (14) 슬라이드 레이아웃 참조 ──
    try:
        orig_layout = orig_slide.slide_layout.name
        edit_layout = edit_slide.slide_layout.name
        if orig_layout != edit_layout:
            changes.append("slide_layout")
            overrides["slide_layout"] = {
                "original": orig_layout,
                "edited": edit_layout,
            }
    except Exception:
        pass

    # 디자인 오버라이드에 도형별 변경 통계 추가
    if shape_changes:
        # 타이포 오버라이드 집계
        typo_overrides = []
        color_overrides = []
        for sc in shape_changes:
            detail = sc.get("detail", {})
            if "typography" in detail:
                typo_overrides.append(detail["typography"])
            if "fill" in detail:
                color_overrides.append(detail["fill"])
        if typo_overrides:
            overrides["typography_changes"] = typo_overrides
        if color_overrides:
            overrides["fill_changes"] = color_overrides

    return {
        "index": slide_index,
        "changes": changes,
        "design_overrides": overrides,
        "shape_changes": shape_changes,
        # 내부 통계 (caller가 pop)
        "_orig_colors": orig_colors,
        "_edit_colors": edit_colors,
        "_orig_fonts": orig_fonts,
        "_edit_fonts": edit_fonts,
    }


# ═══════════════════════════════════════════════════════════════
#  내부 헬퍼 — 도형 매칭
# ═══════════════════════════════════════════════════════════════


def _match_shapes(
    orig_shapes: List, edit_shapes: List
) -> Tuple[List[Tuple], List, List]:
    """
    원본/편집 도형을 매칭.

    전략:
    1. shape.name이 동일하면 매칭 (python-pptx 자동 생성명 포함)
    2. 남은 도형은 (shape_type, left, top) 근접도로 매칭
    """
    matched = []
    used_edit_indices: Set[int] = set()

    # Pass 1: 이름 기반 정확 매칭
    edit_by_name: Dict[str, List[int]] = {}
    for idx, sp in enumerate(edit_shapes):
        edit_by_name.setdefault(sp.name, []).append(idx)

    unmatched_orig_indices = []
    for oi, osp in enumerate(orig_shapes):
        candidates = edit_by_name.get(osp.name, [])
        found = False
        for ei in candidates:
            if ei not in used_edit_indices:
                matched.append((osp, edit_shapes[ei]))
                used_edit_indices.add(ei)
                found = True
                break
        if not found:
            unmatched_orig_indices.append(oi)

    # Pass 2: 위치/유형 근접 매칭 (남은 도형)
    remaining_edit = [
        (i, sp)
        for i, sp in enumerate(edit_shapes)
        if i not in used_edit_indices
    ]

    for oi in unmatched_orig_indices:
        osp = orig_shapes[oi]
        best_match_idx = None
        best_score = float("inf")

        for ei, esp in remaining_edit:
            if ei in used_edit_indices:
                continue
            # 동일 유형 우선
            type_penalty = 0 if osp.shape_type == esp.shape_type else 500000
            # 위치 거리 (EMU)
            try:
                dist = abs(osp.left - esp.left) + abs(osp.top - esp.top)
            except Exception:
                dist = 999999999
            score = type_penalty + dist

            if score < best_score:
                best_score = score
                best_match_idx = ei

        # 임계값: 2인치(1828800 EMU) 이내 + 동일 유형이면 매칭
        if best_match_idx is not None and best_score < 2000000:
            matched.append((osp, edit_shapes[best_match_idx]))
            used_edit_indices.add(best_match_idx)
        # else: 이 원본 도형은 삭제된 것으로 간주

    unmatched_orig = [
        orig_shapes[oi]
        for oi in unmatched_orig_indices
        if not any(m[0] is orig_shapes[oi] for m in matched)
    ]
    unmatched_edit = [
        sp
        for i, sp in enumerate(edit_shapes)
        if i not in used_edit_indices
    ]

    return matched, unmatched_orig, unmatched_edit


# ═══════════════════════════════════════════════════════════════
#  내부 헬퍼 — 도형 속성 비교
# ═══════════════════════════════════════════════════════════════


def _compare_shapes(
    orig_sp, edit_sp
) -> Dict[str, Any]:
    """
    매칭된 두 도형의 모든 시각 속성을 비교.

    Returns:
        {
            "shape_name": "...",
            "action": "modified",
            "changes": ["geometry", "fill", "typography", ...],
            "detail": {
                "geometry": {"left": {"orig": ..., "edit": ...}, ...},
                "fill": {...},
                ...
            },
            "_orig_colors": [...],
            "_edit_colors": [...],
            "_orig_fonts": [...],
            "_edit_fonts": [...]
        }
    """
    changes: List[str] = []
    detail: Dict[str, Any] = {}
    orig_colors: List[str] = []
    edit_colors: List[str] = []
    orig_fonts: List[str] = []
    edit_fonts: List[str] = []

    # ── (4) 지오메트리 ──
    geo_diff = _compare_geometry(orig_sp, edit_sp)
    if geo_diff:
        changes.append("geometry")
        detail["geometry"] = geo_diff

    # ── (9) 도형 유형 ──
    try:
        if orig_sp.shape_type != edit_sp.shape_type:
            changes.append("shape_type")
            detail["shape_type"] = {
                "original": str(orig_sp.shape_type),
                "edited": str(edit_sp.shape_type),
            }
    except Exception:
        pass

    # ── (5) 채우기 ──
    orig_fill = _extract_fill_info(orig_sp)
    edit_fill = _extract_fill_info(edit_sp)
    if orig_fill != edit_fill:
        changes.append("fill")
        detail["fill"] = {"original": orig_fill, "edited": edit_fill}
    # 색상 수집
    for c in _colors_from_fill(orig_fill):
        orig_colors.append(c)
    for c in _colors_from_fill(edit_fill):
        edit_colors.append(c)

    # ── (6) 선/테두리 ──
    orig_line = _extract_line_info(orig_sp)
    edit_line = _extract_line_info(edit_sp)
    if orig_line != edit_line:
        changes.append("line")
        detail["line"] = {"original": orig_line, "edited": edit_line}

    # ── (7) 타이포그래피 + (8) 단락 서식 + (12) 텍스트 프레임 ──
    if orig_sp.has_text_frame and edit_sp.has_text_frame:
        typo_diff, para_diff, frame_diff = _compare_text_properties(
            orig_sp.text_frame, edit_sp.text_frame,
            orig_colors, edit_colors,
            orig_fonts, edit_fonts,
        )
        if typo_diff:
            changes.append("typography")
            detail["typography"] = typo_diff
        if para_diff:
            changes.append("paragraph_format")
            detail["paragraph_format"] = para_diff
        if frame_diff:
            changes.append("text_frame")
            detail["text_frame"] = frame_diff

    # ── (10) 그림자/효과 (XML 레벨) ──
    shadow_diff = _compare_shadow_effects(orig_sp, edit_sp)
    if shadow_diff:
        changes.append("shadow_effects")
        detail["shadow_effects"] = shadow_diff

    # ── (11) 이미지 교체 ──
    if orig_sp.shape_type == 13 and edit_sp.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
        img_diff = _compare_images(orig_sp, edit_sp)
        if img_diff:
            changes.append("image_replaced")
            detail["image"] = img_diff

    return {
        "shape_name": orig_sp.name,
        "action": "modified" if changes else "unchanged",
        "changes": changes,
        "detail": detail,
        "_orig_colors": orig_colors,
        "_edit_colors": edit_colors,
        "_orig_fonts": orig_fonts,
        "_edit_fonts": edit_fonts,
    }


# ═══════════════════════════════════════════════════════════════
#  내부 헬퍼 — 개별 속성 추출/비교
# ═══════════════════════════════════════════════════════════════


def _compare_geometry(orig_sp, edit_sp) -> Optional[Dict[str, Any]]:
    """도형 위치/크기/회전 비교. 변경 있으면 diff dict, 없으면 None."""
    diff: Dict[str, Any] = {}
    threshold = 9144  # 0.01인치 (무시 가능한 미세 변경 필터)

    for attr in ("left", "top", "width", "height"):
        try:
            ov = getattr(orig_sp, attr, None)
            ev = getattr(edit_sp, attr, None)
            if ov is not None and ev is not None:
                if abs(ov - ev) > threshold:
                    diff[attr] = {
                        "original": ov,
                        "edited": ev,
                        "delta_inches": round((ev - ov) / _EMU_PER_INCH, 3),
                    }
        except Exception:
            pass

    # 회전
    try:
        orig_rot = getattr(orig_sp, "rotation", 0) or 0
        edit_rot = getattr(edit_sp, "rotation", 0) or 0
        if abs(orig_rot - edit_rot) > 0.5:  # 0.5도 이상
            diff["rotation"] = {
                "original": orig_rot,
                "edited": edit_rot,
            }
    except Exception:
        pass

    return diff if diff else None


def _extract_fill_info(shape) -> Dict[str, Any]:
    """도형의 채우기 정보 추출 (solid/gradient/pattern/picture/none)."""
    info: Dict[str, Any] = {"type": "none"}
    try:
        fill = shape.fill
        if fill.type is None:
            return info

        fill_type = str(fill.type)
        info["type"] = fill_type

        # SOLID (1)
        if "SOLID" in fill_type.upper() or fill_type == "1":
            try:
                if fill.fore_color and fill.fore_color.rgb:
                    info["color"] = str(fill.fore_color.rgb)
            except Exception:
                pass

        # GRADIENT (4) — 그래디언트 스톱 추출
        elif "GRADIENT" in fill_type.upper() or fill_type == "4":
            try:
                stops = []
                if hasattr(fill, "gradient_stops"):
                    for stop in fill.gradient_stops:
                        stop_info: Dict[str, Any] = {}
                        try:
                            stop_info["position"] = stop.position
                        except Exception:
                            pass
                        try:
                            if stop.color and stop.color.rgb:
                                stop_info["color"] = str(stop.color.rgb)
                        except Exception:
                            pass
                        if stop_info:
                            stops.append(stop_info)
                if stops:
                    info["gradient_stops"] = stops
                # 그래디언트 유형/방향 (XML에서)
                _xml_gradient = _extract_gradient_xml(shape)
                if _xml_gradient:
                    info.update(_xml_gradient)
            except Exception:
                pass

        # PATTERNED (2) / PICTURE (6)
        elif "PATTERN" in fill_type.upper() or fill_type == "2":
            try:
                if fill.fore_color and fill.fore_color.rgb:
                    info["fg_color"] = str(fill.fore_color.rgb)
                if fill.back_color and fill.back_color.rgb:
                    info["bg_color"] = str(fill.back_color.rgb)
            except Exception:
                pass

    except Exception:
        pass

    return info


def _extract_gradient_xml(shape) -> Optional[Dict[str, Any]]:
    """도형 XML에서 그래디언트 세부 정보 추출."""
    try:
        sp_xml = shape._element
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        grad = sp_xml.find(".//a:gradFill", ns)
        if grad is not None:
            info: Dict[str, Any] = {}
            lin = grad.find("a:lin", ns)
            if lin is not None:
                ang = lin.get("ang")
                if ang:
                    info["angle"] = int(ang) / 60000  # 각도 변환
                info["gradient_type"] = "linear"
            else:
                path = grad.find("a:path", ns)
                if path is not None:
                    info["gradient_type"] = path.get("path", "circle")
            return info if info else None
    except Exception:
        pass
    return None


def _extract_line_info(shape) -> Dict[str, Any]:
    """도형의 선/테두리 속성 추출."""
    info: Dict[str, Any] = {"has_line": False}
    try:
        line = shape.line
        if line is None:
            return info

        # 테두리 유무
        if line.fill and line.fill.type is not None:
            info["has_line"] = True

        # 색상
        try:
            if line.color and line.color.rgb:
                info["color"] = str(line.color.rgb)
                info["has_line"] = True
        except Exception:
            pass

        # 두께
        try:
            if line.width is not None:
                info["width_pt"] = round(line.width / _EMU_PER_PT, 2)
                info["has_line"] = True
        except Exception:
            pass

        # 대시 스타일
        try:
            if line.dash_style is not None:
                info["dash_style"] = str(line.dash_style)
        except Exception:
            pass

    except Exception:
        pass
    return info


def _compare_text_properties(
    orig_tf, edit_tf,
    orig_colors: List[str],
    edit_colors: List[str],
    orig_fonts: List[str],
    edit_fonts: List[str],
) -> Tuple[
    Optional[Dict[str, Any]],
    Optional[Dict[str, Any]],
    Optional[Dict[str, Any]],
]:
    """
    텍스트 프레임의 타이포/단락/프레임 속성 비교.

    Returns: (typo_diff, para_diff, frame_diff) — 변경 없으면 None
    """
    typo_diff: Dict[str, Any] = {}
    para_diff: Dict[str, Any] = {}
    frame_diff: Dict[str, Any] = {}

    # ── 텍스트 프레임 속성 (12) ──
    for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        try:
            ov = getattr(orig_tf, attr, None)
            ev = getattr(edit_tf, attr, None)
            if ov != ev and ov is not None and ev is not None:
                frame_diff[attr] = {
                    "original": ov,
                    "edited": ev,
                }
        except Exception:
            pass

    try:
        if orig_tf.word_wrap != edit_tf.word_wrap:
            frame_diff["word_wrap"] = {
                "original": orig_tf.word_wrap,
                "edited": edit_tf.word_wrap,
            }
    except Exception:
        pass

    try:
        orig_auto = str(getattr(orig_tf, "auto_size", None))
        edit_auto = str(getattr(edit_tf, "auto_size", None))
        if orig_auto != edit_auto:
            frame_diff["auto_size"] = {
                "original": orig_auto,
                "edited": edit_auto,
            }
    except Exception:
        pass

    # ── 단락별 비교 ──
    orig_paras = list(orig_tf.paragraphs)
    edit_paras = list(edit_tf.paragraphs)

    para_diffs: List[Dict[str, Any]] = []
    run_diffs: List[Dict[str, Any]] = []

    max_paras = max(len(orig_paras), len(edit_paras))
    for pi in range(max_paras):
        if pi < len(orig_paras) and pi < len(edit_paras):
            op = orig_paras[pi]
            ep = edit_paras[pi]

            # 단락 서식 비교 (8)
            pd = _compare_paragraph_format(op, ep)
            if pd:
                para_diffs.append({"para_index": pi, **pd})

            # 런 비교 (7) — 타이포그래피
            orig_runs = list(op.runs)
            edit_runs = list(ep.runs)
            max_runs = max(len(orig_runs), len(edit_runs))

            for ri in range(max_runs):
                if ri < len(orig_runs) and ri < len(edit_runs):
                    rd = _compare_run_typography(
                        orig_runs[ri], edit_runs[ri]
                    )
                    # 색상/폰트 수집
                    _collect_run_stats(
                        orig_runs[ri], orig_colors, orig_fonts
                    )
                    _collect_run_stats(
                        edit_runs[ri], edit_colors, edit_fonts
                    )
                    if rd:
                        run_diffs.append({
                            "para_index": pi,
                            "run_index": ri,
                            **rd,
                        })
                elif ri < len(edit_runs):
                    # 새로 추가된 런
                    _collect_run_stats(
                        edit_runs[ri], edit_colors, edit_fonts
                    )
                elif ri < len(orig_runs):
                    _collect_run_stats(
                        orig_runs[ri], orig_colors, orig_fonts
                    )
        elif pi < len(edit_paras):
            # 새로 추가된 단락
            for r in edit_paras[pi].runs:
                _collect_run_stats(r, edit_colors, edit_fonts)

    if run_diffs:
        typo_diff["run_changes"] = run_diffs
    if para_diffs:
        para_diff["paragraph_changes"] = para_diffs

    return (
        typo_diff if typo_diff else None,
        para_diff if para_diff else None,
        frame_diff if frame_diff else None,
    )


def _compare_paragraph_format(
    orig_para, edit_para
) -> Optional[Dict[str, Any]]:
    """단락 서식 비교 (정렬, 줄간격, 들여쓰기, 불릿)."""
    diff: Dict[str, Any] = {}

    # 정렬
    try:
        oa = str(getattr(orig_para, "alignment", None))
        ea = str(getattr(edit_para, "alignment", None))
        if oa != ea:
            diff["alignment"] = {"original": oa, "edited": ea}
    except Exception:
        pass

    # 줄간격
    pf_attrs = [
        "line_spacing", "space_before", "space_after",
        "level",
    ]
    for attr in pf_attrs:
        try:
            opf = orig_para.paragraph_format
            epf = edit_para.paragraph_format
            ov = getattr(opf, attr, None)
            ev = getattr(epf, attr, None)
            if ov != ev and (ov is not None or ev is not None):
                diff[attr] = {"original": str(ov), "edited": str(ev)}
        except Exception:
            pass

    # 들여쓰기
    for attr in ("first_line_indent", "left_indent"):
        try:
            opf = orig_para.paragraph_format
            epf = edit_para.paragraph_format
            ov = getattr(opf, attr, None)
            ev = getattr(epf, attr, None)
            if ov != ev and (ov is not None or ev is not None):
                diff[attr] = {"original": ov, "edited": ev}
        except Exception:
            pass

    # 불릿 (XML 레벨)
    try:
        orig_bullet = _extract_bullet_info(orig_para)
        edit_bullet = _extract_bullet_info(edit_para)
        if orig_bullet != edit_bullet:
            diff["bullet"] = {
                "original": orig_bullet,
                "edited": edit_bullet,
            }
    except Exception:
        pass

    return diff if diff else None


def _extract_bullet_info(para) -> Dict[str, Any]:
    """단락의 불릿 정보를 XML에서 추출."""
    info: Dict[str, Any] = {"has_bullet": False}
    try:
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        pPr = para._p.find("a:pPr", ns)
        if pPr is not None:
            buNone = pPr.find("a:buNone", ns)
            if buNone is not None:
                return info

            buChar = pPr.find("a:buChar", ns)
            if buChar is not None:
                info["has_bullet"] = True
                info["type"] = "char"
                info["char"] = buChar.get("char", "")

            buAutoNum = pPr.find("a:buAutoNum", ns)
            if buAutoNum is not None:
                info["has_bullet"] = True
                info["type"] = "auto_num"
                info["num_type"] = buAutoNum.get("type", "")

            buFont = pPr.find("a:buFont", ns)
            if buFont is not None:
                info["bullet_font"] = buFont.get("typeface", "")

            buColor = pPr.find("a:buClr", ns)
            if buColor is not None:
                srgb = buColor.find("a:srgbClr", ns)
                if srgb is not None:
                    info["bullet_color"] = srgb.get("val", "")
    except Exception:
        pass
    return info


def _compare_run_typography(
    orig_run, edit_run
) -> Optional[Dict[str, Any]]:
    """런(텍스트 조각) 수준의 타이포그래피 비교."""
    diff: Dict[str, Any] = {}

    # 폰트 이름
    try:
        of = getattr(orig_run.font, "name", None)
        ef = getattr(edit_run.font, "name", None)
        if of != ef and (of is not None or ef is not None):
            diff["font_name"] = {"original": of, "edited": ef}
    except Exception:
        pass

    # 크기
    try:
        os = getattr(orig_run.font, "size", None)
        es = getattr(edit_run.font, "size", None)
        if os != es and (os is not None or es is not None):
            os_pt = round(os / _EMU_PER_PT, 1) if os else None
            es_pt = round(es / _EMU_PER_PT, 1) if es else None
            diff["font_size_pt"] = {"original": os_pt, "edited": es_pt}
    except Exception:
        pass

    # 굵기
    try:
        ob = getattr(orig_run.font, "bold", None)
        eb = getattr(edit_run.font, "bold", None)
        if ob != eb:
            diff["bold"] = {"original": ob, "edited": eb}
    except Exception:
        pass

    # 기울임
    try:
        oi = getattr(orig_run.font, "italic", None)
        ei = getattr(edit_run.font, "italic", None)
        if oi != ei:
            diff["italic"] = {"original": oi, "edited": ei}
    except Exception:
        pass

    # 밑줄
    try:
        ou = getattr(orig_run.font, "underline", None)
        eu = getattr(edit_run.font, "underline", None)
        if ou != eu:
            diff["underline"] = {"original": str(ou), "edited": str(eu)}
    except Exception:
        pass

    # 글꼴 색상
    try:
        oc = None
        ec = None
        if orig_run.font.color and orig_run.font.color.rgb:
            oc = str(orig_run.font.color.rgb)
        if edit_run.font.color and edit_run.font.color.rgb:
            ec = str(edit_run.font.color.rgb)
        if oc != ec and (oc is not None or ec is not None):
            diff["font_color"] = {"original": oc, "edited": ec}
    except Exception:
        pass

    # 자간 (character spacing) — XML 레벨
    try:
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        orig_rPr = orig_run._r.find("a:rPr", ns)
        edit_rPr = edit_run._r.find("a:rPr", ns)
        os_val = orig_rPr.get("spc") if orig_rPr is not None else None
        es_val = edit_rPr.get("spc") if edit_rPr is not None else None
        if os_val != es_val:
            diff["char_spacing"] = {
                "original": os_val,
                "edited": es_val,
            }
    except Exception:
        pass

    return diff if diff else None


def _collect_run_stats(
    run, colors: List[str], fonts: List[str]
) -> None:
    """런에서 색상/폰트 통계 수집."""
    try:
        if run.font.color and run.font.color.rgb:
            colors.append(str(run.font.color.rgb))
    except Exception:
        pass
    try:
        if run.font.name:
            fonts.append(run.font.name)
    except Exception:
        pass


def _compare_shadow_effects(orig_sp, edit_sp) -> Optional[Dict[str, Any]]:
    """그림자/효과를 XML 레벨에서 비교."""
    try:
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

        orig_shadow = _extract_effect_xml(orig_sp._element, ns)
        edit_shadow = _extract_effect_xml(edit_sp._element, ns)

        if orig_shadow != edit_shadow:
            return {"original": orig_shadow, "edited": edit_shadow}
    except Exception:
        pass
    return None


def _extract_effect_xml(
    sp_elem, ns: Dict[str, str]
) -> Dict[str, Any]:
    """도형 XML에서 효과(그림자/반사/글로우) 정보 추출."""
    effects: Dict[str, Any] = {}

    # 외부 그림자
    outer = sp_elem.find(".//a:outerShdw", ns)
    if outer is not None:
        shadow: Dict[str, Any] = {"type": "outer"}
        for attr in ("blurRad", "dist", "dir", "algn"):
            val = outer.get(attr)
            if val:
                shadow[attr] = val
        # 그림자 색상
        srgb = outer.find("a:srgbClr", ns)
        if srgb is not None:
            shadow["color"] = srgb.get("val", "")
            alpha = srgb.find("a:alpha", ns)
            if alpha is not None:
                shadow["alpha"] = alpha.get("val", "")
        effects["shadow"] = shadow

    # 내부 그림자
    inner = sp_elem.find(".//a:innerShdw", ns)
    if inner is not None:
        effects["inner_shadow"] = {"type": "inner"}
        for attr in ("blurRad", "dist", "dir"):
            val = inner.get(attr)
            if val:
                effects["inner_shadow"][attr] = val

    # 글로우
    glow = sp_elem.find(".//a:glow", ns)
    if glow is not None:
        effects["glow"] = {"rad": glow.get("rad", "")}
        srgb = glow.find("a:srgbClr", ns)
        if srgb is not None:
            effects["glow"]["color"] = srgb.get("val", "")

    # 반사
    refl = sp_elem.find(".//a:reflection", ns)
    if refl is not None:
        effects["reflection"] = {
            attr: refl.get(attr, "")
            for attr in ("blurRad", "stA", "endA", "dist", "dir")
            if refl.get(attr)
        }

    # 소프트 에지
    soft = sp_elem.find(".//a:softEdge", ns)
    if soft is not None:
        effects["soft_edge"] = {"rad": soft.get("rad", "")}

    return effects


def _compare_images(orig_sp, edit_sp) -> Optional[Dict[str, Any]]:
    """이미지 도형의 실제 이미지 변경 감지 (blob hash 비교)."""
    try:
        orig_hash = _image_blob_hash(orig_sp)
        edit_hash = _image_blob_hash(edit_sp)

        if orig_hash and edit_hash and orig_hash != edit_hash:
            return {
                "original_hash": orig_hash[:16],
                "edited_hash": edit_hash[:16],
                "replaced": True,
            }
    except Exception:
        pass
    return None


def _image_blob_hash(shape) -> Optional[str]:
    """이미지 도형의 blob SHA256 해시."""
    try:
        if hasattr(shape, "image") and shape.image:
            blob = shape.image.blob
            return hashlib.sha256(blob).hexdigest()
    except Exception:
        pass
    return None


# ═══════════════════════════════════════════════════════════════
#  내부 헬퍼 — 배경/텍스트 추출 (기존 호환 + 고도화)
# ═══════════════════════════════════════════════════════════════


def _extract_bg_info(slide) -> Dict[str, Any]:
    """
    슬라이드 배경 정보 추출.
    v6.0: gradient 지원 추가.
    """
    info: Dict[str, Any] = {"type": "none"}
    try:
        bg = slide.background
        fill = bg.fill
        if fill.type is None:
            return info

        fill_type = str(fill.type)
        info["type"] = fill_type

        # SOLID
        if "SOLID" in fill_type.upper() or fill_type == "1":
            try:
                if fill.fore_color and fill.fore_color.rgb:
                    info["color"] = str(fill.fore_color.rgb)
            except Exception:
                pass

        # GRADIENT
        elif "GRADIENT" in fill_type.upper() or fill_type == "4":
            try:
                stops = []
                if hasattr(fill, "gradient_stops"):
                    for stop in fill.gradient_stops:
                        stop_info: Dict[str, Any] = {}
                        try:
                            stop_info["position"] = stop.position
                        except Exception:
                            pass
                        try:
                            if stop.color and stop.color.rgb:
                                stop_info["color"] = str(stop.color.rgb)
                        except Exception:
                            pass
                        if stop_info:
                            stops.append(stop_info)
                if stops:
                    info["gradient_stops"] = stops
            except Exception:
                pass

            # XML에서 추가 정보
            try:
                ns = {
                    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
                }
                bg_elem = slide.background._element
                grad = bg_elem.find(".//a:gradFill", ns)
                if grad is not None:
                    lin = grad.find("a:lin", ns)
                    if lin is not None:
                        ang = lin.get("ang")
                        if ang:
                            info["angle"] = int(ang) / 60000
                        info["gradient_type"] = "linear"
                    else:
                        path = grad.find("a:path", ns)
                        if path is not None:
                            info["gradient_type"] = path.get(
                                "path", "circle"
                            )
            except Exception:
                pass

        # PICTURE (배경 이미지)
        elif "PICTURE" in fill_type.upper() or fill_type == "6":
            info["type"] = "picture"
            # 이미지 해시는 비용이 높으므로 "있음"만 표시

    except Exception:
        pass
    return info


def _extract_all_text(slide) -> str:
    """슬라이드의 모든 텍스트를 결합."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    return "\n".join(texts)


def _extract_full_slide_info(
    slide, slide_index: int
) -> Dict[str, Any]:
    """새로 추가된 슬라이드의 전체 디자인 스냅샷."""
    info: Dict[str, Any] = {
        "index": slide_index,
        "background": _extract_bg_info(slide),
        "shape_count": len(slide.shapes),
        "shapes": [],
    }

    for sp in slide.shapes:
        sp_snap = _extract_shape_snapshot(sp)
        info["shapes"].append(sp_snap)

    return info


def _extract_shape_snapshot(shape) -> Dict[str, Any]:
    """단일 도형의 전체 속성 스냅샷."""
    snap: Dict[str, Any] = {
        "name": shape.name,
        "shape_type": str(shape.shape_type),
    }
    colors: List[str] = []
    fonts: List[str] = []

    # 지오메트리
    for attr in ("left", "top", "width", "height"):
        try:
            snap[attr] = getattr(shape, attr, None)
        except Exception:
            pass
    try:
        snap["rotation"] = getattr(shape, "rotation", 0)
    except Exception:
        pass

    # 채우기
    snap["fill"] = _extract_fill_info(shape)
    for c in _colors_from_fill(snap["fill"]):
        colors.append(c)

    # 선
    snap["line"] = _extract_line_info(shape)

    # 텍스트
    if shape.has_text_frame:
        snap["text"] = shape.text_frame.text[:200]  # 200자 제한
        # 첫 런의 타이포
        try:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    _collect_run_stats(run, colors, fonts)
                    try:
                        snap["first_font"] = run.font.name
                        snap["first_size_pt"] = round(
                            run.font.size / _EMU_PER_PT, 1
                        ) if run.font.size else None
                        if run.font.color and run.font.color.rgb:
                            snap["first_color"] = str(
                                run.font.color.rgb
                            )
                    except Exception:
                        pass
                    break  # 첫 런만
                break  # 첫 단락만
        except Exception:
            pass

    # 이미지
    if shape.shape_type == 13:
        snap["image_hash"] = _image_blob_hash(shape)

    snap["_colors"] = colors
    snap["_fonts"] = fonts

    return snap


# ═══════════════════════════════════════════════════════════════
#  내부 헬퍼 — 유틸리티
# ═══════════════════════════════════════════════════════════════


def _colors_from_fill(fill_info: Dict[str, Any]) -> List[str]:
    """fill_info에서 색상 문자열 목록 추출."""
    colors = []
    if "color" in fill_info:
        colors.append(fill_info["color"])
    if "fg_color" in fill_info:
        colors.append(fill_info["fg_color"])
    if "bg_color" in fill_info:
        colors.append(fill_info["bg_color"])
    for stop in fill_info.get("gradient_stops", []):
        if "color" in stop:
            colors.append(stop["color"])
    return colors


def _top_items(items: List[str], n: int = 5) -> Dict[str, int]:
    """리스트에서 빈도 상위 n개 항목 반환."""
    counts: Dict[str, int] = {}
    for item in items:
        if item:
            counts[item] = counts.get(item, 0) + 1
    sorted_items = sorted(
        counts.items(), key=lambda x: x[1], reverse=True
    )
    return dict(sorted_items[:n])


def _build_gamma_image_options(brief) -> Dict[str, Any]:
    """
    MergedDesignBrief → Gamma imageOptions 파라미터.

    프로젝트 맥락(산업, 유형, 키워드, 톤)을 반영하여
    source + style + model을 결정합니다.

    source 결정:
        photorealistic  → pexels (실사진)
        illustration    → aiGenerated (AI 생성)
        minimal/mixed   → aiGenerated (더 다양한 결과)

    style 결정 (자유형식 — 프로젝트 맥락 반영):
        산업/유형별 시각 분위기 + 키워드 + 톤 레벨 종합

    model 결정:
        기본 imagen-4-pro (Gamma API 유효 모델 기준)
    """
    img_style = getattr(brief, "image_style", None)
    project_type = getattr(brief, "project_type", "event")
    industry = getattr(brief, "industry", "")
    content_tone = getattr(brief, "content_tone", None)
    keywords = getattr(img_style, "keywords", []) if img_style else []
    primary_style = getattr(img_style, "primary_style", "") if img_style else ""

    opts: Dict[str, Any] = {}

    # ── source ──
    if primary_style == "photorealistic":
        opts["source"] = "pexels"
    elif primary_style == "illustration":
        opts["source"] = "aiGenerated"
    elif primary_style == "minimal":
        opts["source"] = "aiGenerated"
    else:
        opts["source"] = "aiGenerated"

    # ── style (프로젝트 맥락 반영 자유형식) ──
    style_parts: List[str] = []

    # 산업별 시각 분위기 베이스
    industry_styles = {
        "game_event": "dark cinematic game art, dramatic lighting, sci-fi futuristic atmosphere",
        "event": "vibrant event photography, dynamic stage lighting, immersive venue atmosphere",
        "marketing_pr": "clean modern lifestyle, vibrant brand-forward visuals, social media aesthetic",
        "it_system": "minimal flat design, clean UI mockups, tech-forward diagrams",
        "public": "warm community-focused imagery, diverse people, civic and public space",
        "consulting": "professional corporate photography, clean data visualization, executive boardroom",
    }
    base_style = industry_styles.get(
        industry,
        industry_styles.get(project_type, "professional and clean visuals")
    )
    style_parts.append(base_style)

    # 키워드 반영 (최대 4개)
    if keywords:
        kw_str = ", ".join(keywords[:4])
        style_parts.append(f"themed around: {kw_str}")

    # 톤 레벨 반영
    if content_tone:
        level = getattr(content_tone, "emotional_tone_level", 3)
        if level >= 4:
            style_parts.append("emotionally evocative, storytelling-driven imagery")
        elif level <= 2:
            style_parts.append("factual, diagram-oriented, minimal decoration")

    # 색상 무드 반영
    colors = getattr(brief, "colors", {})
    if colors:
        dark = colors.get("dark", (33, 33, 33))
        primary = colors.get("primary", (0, 44, 95))
        # 다크 테마 감지
        if dark[0] < 40 and dark[1] < 40 and dark[2] < 40:
            if primary[2] > primary[0]:  # 블루 계열
                style_parts.append("dark background with blue accent lighting")
            elif primary[0] > primary[2]:  # 레드/웜 계열
                style_parts.append("dark background with warm accent tones")

    opts["style"] = ", ".join(style_parts)

    # ── model (AI 이미지 생성 시만 사용) ──
    # Gamma API 유효 모델: imagen-3-flash, imagen-3-pro, imagen-4-pro,
    #   imagen-4-ultra, ideogram-v3-turbo, flux-1-pro, flux-1-quick 등
    if opts["source"] == "aiGenerated":
        opts["model"] = "imagen-4-pro"

    return opts


def _build_gamma_card_options(brief) -> Dict[str, Any]:
    """
    MergedDesignBrief → Gamma cardOptions 파라미터.

    headerFooter로 전 슬라이드 일관된 헤더/푸터를 설정하여
    Gamma 결과물의 시각적 일관성을 확보합니다.
    """
    card_opts: Dict[str, Any] = {
        "dimensions": "16x9",
    }

    project_name = getattr(brief, "project_name", "")

    # headerFooter — 전 슬라이드 일관된 상하단 요소
    header_footer: Dict[str, Any] = {
        "hideFromFirstCard": True,   # 표지 제외
        "hideFromLastCard": True,    # 클로징 제외
    }

    # 좌상단: 프로젝트명 (일관된 위치)
    if project_name:
        header_footer["topLeft"] = {
            "type": "text",
            "value": project_name,
        }

    # 우하단: 카드 번호 (페이지 번호)
    header_footer["bottomRight"] = {
        "type": "cardNumber",
    }

    card_opts["headerFooter"] = header_footer

    return card_opts


def _build_additional_instructions(brief) -> str:
    """
    MergedDesignBrief → Gamma additionalInstructions.

    Gamma가 슬라이드를 생성할 때 참조하는 디자인 규칙.
    themeId/imageOptions/textOptions로 전달할 수 없는
    레이아웃 일관성, 시각 위계, 콘텐츠 구조 힌트를 전달합니다.

    최대 2000자 제한 내에서 우선순위별 구성:
        1. 시각 위계 규칙 (필수)
        2. 레이아웃 일관성 (필수)
        3. 프로젝트 맥락 (선택)
        4. 배경 패턴 (선택)
    """
    lines: List[str] = []

    # ═══ 1. 시각 위계 (Visual Hierarchy) ═══
    lines.append(
        "VISUAL HIERARCHY RULES: "
        "Maintain consistent visual hierarchy across ALL slides. "
        "Section divider titles: largest size, centered. "
        "Content slide titles: same size and position on every slide. "
        "Body text: uniform size throughout. "
        "Never vary title font size between content slides."
    )

    # ═══ 2. 레이아웃 일관성 ═══
    lines.append(
        "LAYOUT CONSISTENCY: "
        "Keep title position fixed at the same coordinates on every content slide. "
        "Use consistent margins (left, top, right). "
        "Align visual elements to a grid. "
        "Section dividers should use a distinct layout from content slides."
    )

    # ═══ 3. 프로젝트 맥락 ═══
    project_name = getattr(brief, "project_name", "")
    project_type = getattr(brief, "project_type", "")
    industry = getattr(brief, "industry", "")

    if project_name or industry:
        ctx_parts = []
        if project_name:
            ctx_parts.append(f"Project: {project_name}")
        if industry:
            industry_labels = {
                "game_event": "Gaming/Entertainment event",
                "event": "Event/Exhibition",
                "marketing_pr": "Marketing & PR campaign",
                "it_system": "IT/System project",
                "public": "Public sector project",
                "consulting": "Consulting engagement",
            }
            ctx_parts.append(
                f"Industry: {industry_labels.get(industry, industry)}"
            )
        lines.append(f"CONTEXT: {'. '.join(ctx_parts)}.")

    # ═══ 4. 색상 가이드 (themeId 보완) ═══
    colors = getattr(brief, "colors", {})
    if colors:
        color_parts = []
        for key in ("primary", "secondary", "accent"):
            c = colors.get(key)
            if c:
                hex_val = f"#{c[0]:02X}{c[1]:02X}{c[2]:02X}"
                color_parts.append(f"{key}: {hex_val}")
        if color_parts:
            lines.append(
                f"COLOR PALETTE: {', '.join(color_parts)}. "
                "Use these colors consistently for headings, accents, and key visuals."
            )

    # ═══ 5. 배경 패턴 ═══
    bg_sched = getattr(brief, "background_schedule", [])
    if bg_sched:
        # 배경 유형별 빈도 요약
        bg_counts: Dict[str, int] = {}
        for bg in bg_sched:
            bg_counts[bg] = bg_counts.get(bg, 0) + 1
        sorted_bgs = sorted(bg_counts.items(), key=lambda x: x[1], reverse=True)
        bg_desc = ", ".join(f"{k}({v})" for k, v in sorted_bgs[:4])
        lines.append(
            f"BACKGROUND PATTERN: Alternate backgrounds — {bg_desc}. "
            "Use dark backgrounds for section dividers, "
            "white/light for content slides."
        )

    # ═══ 6. 콘텐츠 구조 힌트 ═══
    content_tone = getattr(brief, "content_tone", None)
    if content_tone:
        level = getattr(content_tone, "emotional_tone_level", 3)
        framing = getattr(content_tone, "narrative_framing_style", "")
        if level >= 4:
            lines.append(
                "CONTENT STYLE: Narrative and emotionally engaging. "
                "Use full-bleed images for impact. "
                "Allow breathing room for key messages."
            )
        elif level <= 2:
            lines.append(
                "CONTENT STYLE: Data-driven and concise. "
                "Prioritize charts, tables, and metrics. "
                "Minimize decorative elements."
            )
        if framing == "worldview_based":
            metaphor = getattr(content_tone, "core_metaphor", "")
            if metaphor:
                lines.append(
                    f"NARRATIVE: Maintain '{metaphor}' as recurring visual motif."
                )

    # ═══ 7. 컴포넌트 밀도 힌트 ═══
    visual_density = getattr(brief, "visual_density_targets", {})
    if visual_density:
        img_pct = visual_density.get("image_slides_pct", 0)
        diag_pct = visual_density.get("diagram_slides_pct", 0)
        if img_pct > 0.25:
            lines.append(
                "VISUAL DENSITY: Image-heavy presentation. "
                f"~{int(img_pct*100)}% of slides should feature prominent images."
            )
        if diag_pct > 0.30:
            lines.append(
                f"~{int(diag_pct*100)}% of slides should contain diagrams or charts."
            )

    result = "\n".join(lines)

    # 2000자 제한 — 우선순위 높은 앞부분 보존
    if len(result) > 2000:
        result = result[:1997] + "..."

    return result


# ── 하위호환: 기존 DesignBridge 클래스 유지 ──

class DesignBridge(GammaMCPBridge):
    """하위호환 별칭 (기존 코드에서 DesignBridge 임포트 시)."""

    def __init__(
        self, provider: str = "gamma", project_dir: Optional[Path] = None
    ):
        super().__init__(project_dir=project_dir)
        self.provider = provider
