"""
PPTX 레퍼런스 입수

PPTX 문서에서 텍스트, 레이아웃, 디자인 프로파일을 추출합니다.
"""

from __future__ import annotations

from pathlib import Path
from typing import Dict, List, Optional

from ..models import (
    ReferenceDocument,
    SectionStructure,
    DocType,
    Industry,
)
from .dedup_checker import compute_file_hash
from .content_extractor import ContentExtractor
from .design_extractor import DesignExtractor
from src.utils.logger import get_logger

logger = get_logger("pptx_ingester")


class PPTXIngester:
    """PPTX 레퍼런스 입수"""

    def __init__(self):
        self.design_extractor = DesignExtractor()
        self.content_extractor = ContentExtractor()

    def ingest(
        self,
        file_path: Path,
        doc_type: DocType = DocType.PROPOSAL,
        industry: Industry = Industry.OTHER,
        project_type: str = "",
        won_bid: bool = False,
        tags: Optional[List[str]] = None,
        notes: str = "",
    ) -> ReferenceDocument:
        """
        PPTX 파일을 분석하여 ReferenceDocument 생성

        Args:
            file_path: PPTX 파일 경로
            doc_type: 문서 유형
            industry: 산업 분류
            project_type: 프로젝트 유형
            won_bid: 수주 성공 여부
            tags: 태그
            notes: 메모

        Returns:
            ReferenceDocument: 추출된 레퍼런스 데이터
        """
        from pptx import Presentation

        file_hash = compute_file_hash(file_path)
        prs = Presentation(str(file_path))

        # 텍스트 추출
        texts, slide_texts = self._extract_all_text(prs)
        total_pages = len(prs.slides)

        # 섹션 구조 분석
        sections = self._analyze_sections(slide_texts)

        # 디자인 프로파일 추출
        design_profile = self.design_extractor.extract(prs)

        # 콘텐츠 패턴 + 프로그램 템플릿 추출
        content_patterns = self.content_extractor.extract_content_patterns(
            prs, slide_texts, sections
        )
        program_templates = self.content_extractor.extract_program_templates(
            prs, slide_texts, sections
        )

        doc = ReferenceDocument(
            file_path=str(file_path.absolute()),
            file_hash=file_hash,
            file_name=file_path.name,
            file_size=file_path.stat().st_size,
            doc_type=doc_type,
            industry=industry,
            project_type=project_type,
            won_bid=won_bid,
            total_pages=total_pages,
            sections=sections,
            design_profile=design_profile,
            content_patterns=content_patterns,
            program_templates=program_templates,
            full_text=texts[:100000],
            tags=tags or [],
            notes=notes,
        )

        logger.info(
            f"PPTX 입수 완료: {file_path.name} "
            f"({total_pages}슬라이드, {len(design_profile.colors)}색상, "
            f"{len(design_profile.fonts)}폰트, "
            f"{len(content_patterns)}패턴, {len(program_templates)}템플릿)"
        )
        return doc

    def _extract_all_text(self, prs) -> tuple[str, List[Dict]]:
        """전체 텍스트 + 슬라이드별 텍스트 추출"""
        all_texts = []
        slide_texts = []

        for i, slide in enumerate(prs.slides):
            slide_text_parts = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_text_parts.append(text)

                # 테이블 텍스트
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text = cell.text.strip()
                            if text:
                                slide_text_parts.append(text)

            slide_dict = {
                "slide_number": i + 1,
                "texts": slide_text_parts,
                "full_text": "\n".join(slide_text_parts),
            }
            slide_texts.append(slide_dict)
            all_texts.extend(slide_text_parts)

        return "\n\n".join(all_texts), slide_texts

    def _analyze_sections(self, slide_texts: List[Dict]) -> List[SectionStructure]:
        """슬라이드별 텍스트에서 섹션 구조 분석"""
        sections = []
        current_section = None
        current_count = 0

        for st in slide_texts:
            texts = st["texts"]
            if not texts:
                current_count += 1
                continue

            first_text = texts[0]

            # 섹션 구분자 감지 (짧은 텍스트, 대문자/번호)
            is_divider = (
                len(texts) <= 3
                and len(first_text) < 40
                and (
                    first_text.isupper()
                    or first_text.startswith(("0", "1", "2", "3", "4", "5", "6", "7", "8", "9"))
                    or any(kw in first_text.upper() for kw in [
                        "PHASE", "PART", "SECTION", "CHAPTER",
                        "INSIGHT", "CONCEPT", "STRATEGY", "ACTION",
                        "MANAGEMENT", "WHY US", "INVESTMENT",
                        "HOOK", "SUMMARY", "APPENDIX",
                    ])
                )
            )

            if is_divider:
                if current_section:
                    sections.append(SectionStructure(
                        name=current_section,
                        slide_count=current_count,
                    ))
                current_section = first_text
                current_count = 1
            else:
                current_count += 1

        # 마지막 섹션
        if current_section:
            sections.append(SectionStructure(
                name=current_section,
                slide_count=current_count,
            ))

        # 섹션이 없으면 전체를 하나로
        if not sections and slide_texts:
            sections.append(SectionStructure(
                name="전체",
                slide_count=len(slide_texts),
            ))

        # 가중치 계산
        total_slides = sum(s.slide_count for s in sections)
        if total_slides > 0:
            for s in sections:
                s.weight_pct = round(s.slide_count / total_slides * 100, 1)

        return sections
