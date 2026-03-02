"""
디자인 프로파일 추출기

PPTX에서 컬러 팔레트, 폰트 계층, 레이아웃 패턴을 자동 추출합니다.
"""

from __future__ import annotations

from collections import Counter
from typing import Dict, List

from ..models import (
    ColorInfo,
    DesignProfile,
    FontInfo,
    LayoutPattern,
)
from src.utils.logger import get_logger

logger = get_logger("design_extractor")


class DesignExtractor:
    """PPTX 디자인 프로파일 추출"""

    def extract(self, prs) -> DesignProfile:
        """
        프레젠테이션에서 디자인 프로파일 추출

        Args:
            prs: python-pptx Presentation 객체

        Returns:
            DesignProfile: 추출된 디자인 정보
        """
        colors = self._extract_colors(prs)
        fonts = self._extract_fonts(prs)
        font_hierarchy = self._build_font_hierarchy(fonts)
        layout_patterns = self._analyze_layouts(prs)

        # 슬라이드 크기
        slide_width = prs.slide_width.inches if hasattr(prs.slide_width, 'inches') else 13.33
        slide_height = prs.slide_height.inches if hasattr(prs.slide_height, 'inches') else 7.5

        # 배경 스타일 판별
        bg_style = self._detect_bg_style(prs)

        return DesignProfile(
            colors=colors,
            fonts=fonts,
            font_hierarchy=font_hierarchy,
            layout_patterns=layout_patterns,
            bg_style=bg_style,
            aspect_ratio=f"{round(slide_width / slide_height, 2)}:1",
            slide_dimensions={"width": slide_width, "height": slide_height},
        )

    def _extract_colors(self, prs) -> List[ColorInfo]:
        """모든 슬라이드에서 사용된 색상 추출"""
        from pptx.util import Pt
        from pptx.dml.color import RGBColor

        color_counter: Counter = Counter()

        for slide in prs.slides:
            for shape in slide.shapes:
                # 텍스트 색상
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                if run.font.color and run.font.color.rgb:
                                    hex_color = f"#{run.font.color.rgb}"
                                    color_counter[hex_color] += 1
                            except Exception:
                                pass

                # 도형 채우기 색상
                try:
                    fill = shape.fill
                    if fill and fill.type is not None:
                        if hasattr(fill, 'fore_color') and fill.fore_color and fill.fore_color.rgb:
                            hex_color = f"#{fill.fore_color.rgb}"
                            color_counter[hex_color] += 1
                except Exception:
                    pass

        # 빈도순 정렬
        total = sum(color_counter.values()) or 1
        colors = []
        for hex_color, count in color_counter.most_common(20):
            freq = count / total
            # 사용 용도 추정
            usage = self._guess_color_usage(hex_color, freq)
            colors.append(ColorInfo(
                hex=hex_color,
                usage=usage,
                frequency=round(freq, 3),
            ))

        return colors

    def _guess_color_usage(self, hex_color: str, frequency: float) -> str:
        """색상 용도 추정"""
        hex_lower = hex_color.lower().lstrip("#")

        # 흰색 계열
        if hex_lower in ("ffffff", "fefefe", "f5f5f5", "fafafa"):
            return "background"
        # 검정 계열
        if hex_lower in ("000000", "1a1a1a", "333333", "2d2d2d"):
            return "text"
        # 회색 계열
        if all(abs(int(hex_lower[i:i+2], 16) - int(hex_lower[0:2], 16)) < 20 for i in (0, 2, 4)):
            mid = int(hex_lower[0:2], 16)
            if mid > 200:
                return "background"
            elif mid > 100:
                return "secondary_text"
            else:
                return "text"

        # 가장 빈도 높은 유채색 = primary
        if frequency > 0.1:
            return "primary"
        elif frequency > 0.05:
            return "secondary"
        else:
            return "accent"

    def _extract_fonts(self, prs) -> List[FontInfo]:
        """폰트 정보 추출"""
        from pptx.util import Pt

        font_counter: Counter = Counter()  # (name, size_pt, bold) → count

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                name = run.font.name or "Unknown"
                                size = run.font.size
                                size_pt = round(size.pt, 1) if size else 0
                                bold = run.font.bold or False
                                font_counter[(name, size_pt, bold)] += 1
                            except Exception:
                                pass

        # 빈도순 정렬
        fonts = []
        for (name, size_pt, bold), count in font_counter.most_common(15):
            usage = self._guess_font_usage(size_pt, bold)
            fonts.append(FontInfo(
                name=name,
                size_pt=size_pt,
                bold=bold,
                usage=usage,
            ))

        return fonts

    def _guess_font_usage(self, size_pt: float, bold: bool) -> str:
        """폰트 용도 추정"""
        if size_pt >= 36:
            return "title"
        elif size_pt >= 24:
            return "subtitle"
        elif size_pt >= 18:
            return "heading"
        elif size_pt >= 14:
            return "body"
        elif size_pt >= 10:
            return "caption"
        else:
            return "footnote"

    def _build_font_hierarchy(self, fonts: List[FontInfo]) -> Dict[str, str]:
        """폰트 계층 구조 구축"""
        hierarchy = {}
        for font in fonts:
            key = font.usage
            if key not in hierarchy:
                bold_str = " Bold" if font.bold else ""
                hierarchy[key] = f"{font.name}{bold_str} {font.size_pt}pt"
        return hierarchy

    def _analyze_layouts(self, prs) -> List[LayoutPattern]:
        """레이아웃 패턴 분석"""
        pattern_counter: Counter = Counter()
        total_slides = len(prs.slides)

        for slide in prs.slides:
            pattern = self._classify_slide_layout(slide)
            pattern_counter[pattern] += 1

        patterns = []
        for pattern_type, count in pattern_counter.most_common():
            patterns.append(LayoutPattern(
                pattern_type=pattern_type,
                frequency=round(count / total_slides, 3) if total_slides > 0 else 0,
                description=self._describe_pattern(pattern_type),
            ))

        return patterns

    def _classify_slide_layout(self, slide) -> str:
        """슬라이드 레이아웃 유형 분류 (12 유형)"""
        shapes = list(slide.shapes)
        text_shapes = [s for s in shapes if s.has_text_frame]
        image_shapes = [s for s in shapes if s.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE
        table_shapes = [s for s in shapes if s.has_table]
        auto_shapes = [s for s in shapes
                       if hasattr(s, 'shape_type') and s.shape_type not in (13, 19)]

        n_shapes = len(shapes)
        n_text = len(text_shapes)
        n_image = len(image_shapes)
        n_table = len(table_shapes)
        n_auto = len(auto_shapes)

        # 이미지 커버리지 계산
        slide_area = 1.0
        try:
            sw = slide.part.package.presentation.slide_width
            sh = slide.part.package.presentation.slide_height
            if sw and sh:
                slide_area = (sw / 914400) * (sh / 914400)
        except Exception:
            slide_area = 13.333 * 7.5

        image_area = 0.0
        for img in image_shapes:
            try:
                if img.width and img.height:
                    image_area += (img.width / 914400) * (img.height / 914400)
            except Exception:
                pass
        image_coverage = image_area / slide_area if slide_area > 0 else 0

        # 1. full_bleed_image: 이미지가 슬라이드 50%+ 차지
        if image_coverage > 0.5:
            return "full_bleed_image"
        # 2. split_image_text: 이미지 + 텍스트 좌우/상하 분할
        if n_image > 0 and n_text > 2 and image_coverage > 0.15:
            return "split_image_text"
        # 3. multi_image_grid: 다중 이미지 그리드
        if n_image >= 3:
            return "multi_image_grid"
        # 4. image_focused: 이미지 있고 텍스트 적음
        if n_image > 0 and n_text <= 2:
            return "image_focused"
        # 5. table_based: 테이블 포함
        if n_table > 0:
            return "table_based"
        # 6. hero_typography: 대형 텍스트 1~2개 (표지/구분자/인용)
        if n_text <= 2 and n_shapes <= 4:
            return "hero_typography"
        # 7. data_card_array: 카드 배열 (6+개 소형 도형)
        if n_auto >= 6 and n_table == 0 and n_image == 0:
            return "data_card_array"
        # 8. process_flow: 흐름도/프로세스 (중간 도형 수)
        if 5 <= n_shapes <= 12 and n_auto >= 3:
            return "process_flow"
        # 9. layered_composition: 고밀도 겹침 레이아웃
        if n_shapes > 15:
            return "layered_composition"
        # 10. complex_diagram: 복합 도식
        if n_shapes > 10:
            return "complex_diagram"
        # 11. multi_column: 다단 텍스트
        if n_text > 4:
            return "multi_column"
        # 12. content_standard: 일반 콘텐츠
        if n_text > 2:
            return "content_standard"

        return "other"

    def _describe_pattern(self, pattern_type: str) -> str:
        """패턴 설명"""
        descriptions = {
            "full_bleed_image": "전폭 이미지 레이아웃 (이미지 50%+ 커버리지)",
            "split_image_text": "이미지+텍스트 좌우/상하 분할 레이아웃",
            "multi_image_grid": "다중 이미지 그리드 (3+장)",
            "image_focused": "이미지 중심 레이아웃",
            "table_based": "테이블 기반 데이터 표시",
            "hero_typography": "히어로 타이포그래피 (표지, 구분자, 인용)",
            "data_card_array": "카드 배열 데이터 (6+개 소형 도형)",
            "process_flow": "흐름도/프로세스 다이어그램",
            "layered_composition": "고밀도 레이어드 구성 (15+개 도형)",
            "complex_diagram": "복합 도식/다이어그램",
            "multi_column": "다단 레이아웃",
            "content_standard": "일반 콘텐츠 (제목 + 본문)",
            "other": "기타",
        }
        return descriptions.get(pattern_type, "")

    def _detect_bg_style(self, prs) -> str:
        """배경 스타일 감지"""
        dark_count = 0
        light_count = 0

        for slide in prs.slides:
            try:
                bg = slide.background
                fill = bg.fill
                if fill and fill.type is not None:
                    if hasattr(fill, 'fore_color') and fill.fore_color and fill.fore_color.rgb:
                        rgb = str(fill.fore_color.rgb)
                        r, g, b = int(rgb[0:2], 16), int(rgb[2:4], 16), int(rgb[4:6], 16)
                        brightness = (r + g + b) / 3
                        if brightness < 128:
                            dark_count += 1
                        else:
                            light_count += 1
            except Exception:
                light_count += 1  # 기본은 밝은 배경

        if dark_count > light_count:
            return "dark"
        elif dark_count > 0:
            return "mixed"
        return "light"
